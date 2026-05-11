#!/usr/bin/env python3
"""inject_notes.py — Inject speaker notes into a PPTX file using python-pptx.

Usage:
    python inject_notes.py <input.pptx> <notes.json> [output.pptx]

Notes JSON format:
    {
        "slide_notes": {
            "1": "Speaker notes for slide 1...",
            "2": "Speaker notes for slide 2..."
        }
    }

Dependencies:
    pip install python-pptx
"""

import json
import sys
from pathlib import Path
from pptx import Presentation


def inject_notes(pptx_path, notes_json_path, output_path=None):
    """Inject speaker notes from a JSON file into a PPTX presentation.

    Performance: single-pass inject + in-memory verify (no re-read from disk).
    """
    pptx_path = Path(pptx_path)
    if output_path is None:
        output_path = pptx_path.parent / f'{pptx_path.stem}_with_notes.pptx'
    else:
        output_path = Path(output_path)

    # Load notes JSON
    with open(notes_json_path, 'r', encoding='utf-8') as f:
        slide_notes = json.load(f).get('slide_notes', {})

    if not slide_notes:
        print("Error: No slide_notes found in JSON.")
        sys.exit(1)

    # Open presentation (single load)
    prs = Presentation(str(pptx_path))
    total_slides = len(prs.slides)
    results = []  # collect output lines for batch print

    results.append(f"Input:  {pptx_path}")
    results.append(f"Notes:  {notes_json_path} ({len(slide_notes)} entries)")
    results.append(f"Output: {output_path}")
    results.append(f"Slides: {total_slides}")
    results.append("")

    injected = 0
    skipped = 0

    for idx, slide in enumerate(prs.slides, start=1):
        key = str(idx)
        if key not in slide_notes:
            results.append(f"  Slide {idx:2d}: ⏭  no notes in JSON, skipped")
            skipped += 1
            continue

        note_text = slide_notes[key]

        # Access or create notes slide (auto-creates if needed)
        tf = slide.notes_slide.notes_text_frame
        tf.clear()

        # Split by newlines and add paragraphs
        lines = note_text.split('\n')
        tf.paragraphs[0].text = lines[0]
        for line in lines[1:]:
            tf.add_paragraph().text = line

        injected += 1
        results.append(f"  Slide {idx:2d}: ✅ injected ({len(note_text)} chars)")

    # Save
    prs.save(str(output_path))
    results.append(f"\n{'='*50}")
    results.append(f"Injection complete: {injected}/{total_slides} slides injected, {skipped} skipped")

    # ===== In-memory Verification (no disk re-read) =====
    results.append(f"\n{'='*50}")
    results.append("Verification (in-memory, same Presentation object):")
    verify_ok = 0
    verify_fail = 0

    for idx, slide in enumerate(prs.slides, start=1):
        key = str(idx)
        if key not in slide_notes:
            continue
        if slide.has_notes_slide:
            text = slide.notes_slide.notes_text_frame.text
            if len(text) > 0:
                verify_ok += 1
                results.append(f"  Slide {idx:2d}: ✅ verified ({len(text)} chars)")
            else:
                verify_fail += 1
                results.append(f"  Slide {idx:2d}: ❌ EMPTY after inject!")
        else:
            verify_fail += 1
            results.append(f"  Slide {idx:2d}: ❌ NO NOTES SLIDE!")

    if verify_fail > 0:
        results.append(f"\n⚠️  VERIFICATION FAILED: {verify_fail} slide(s) have issues!")
        print('\n'.join(results))
        sys.exit(2)
    else:
        results.append(f"\n✅ All {verify_ok} slides verified successfully!")
        results.append(f"Output: {output_path}")

    # Batch print all output at once
    print('\n'.join(results))
    return str(output_path)


if __name__ == '__main__':
    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(1)
    inject_notes(
        sys.argv[1],
        sys.argv[2],
        sys.argv[3] if len(sys.argv) > 3 else None
    )