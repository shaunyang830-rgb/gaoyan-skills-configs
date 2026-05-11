# Copyright 2024-2026 Kaku Li (https://github.com/likaku)
# Licensed under the Apache License, Version 2.0 — see LICENSE and NOTICE.
# Part of "Mck-ppt-design-skill" (McKinsey PPT Design Framework).
# NOTICE: This file must be retained in all copies or substantial portions.
#
"""Low-level drawing primitives for McKinsey PPT framework.

All coordinate calculations, shape creation, XML cleanup happen here.
Higher-level code should NOT import pptx directly — use these helpers.
"""
import os, zipfile
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

from .constants import *


# ═══════════════════════════════════════════
# XML CLEANUP
# ═══════════════════════════════════════════

def _clean_shape(shape):
    """Remove p:style from a shape to prevent theme effect leaks."""
    sp = shape._element
    style = sp.find(qn('p:style'))
    if style is not None:
        sp.remove(style)


def full_cleanup(outpath):
    """Post-save nuclear sanitization: remove ALL p:style + theme shadows/3D."""
    tmppath = outpath + '.tmp'
    with zipfile.ZipFile(outpath, 'r') as zin:
        with zipfile.ZipFile(tmppath, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename.endswith('.xml'):
                    root = etree.fromstring(data)
                    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                    for style in root.findall(f'.//{{{ns_p}}}style'):
                        style.getparent().remove(style)
                    if 'theme' in item.filename.lower():
                        for tag in ['outerShdw', 'innerShdw', 'scene3d', 'sp3d']:
                            for el in root.findall(f'.//{{{ns_a}}}{tag}'):
                                el.getparent().remove(el)
                    data = etree.tostring(root, xml_declaration=True,
                                          encoding='UTF-8', standalone=True)
                zout.writestr(item, data)
    os.replace(tmppath, outpath)


# ═══════════════════════════════════════════
# FONT HELPERS
# ═══════════════════════════════════════════

def set_ea_font(run, typeface='KaiTi'):
    """Set East Asian font on a text run for Chinese rendering."""
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', typeface)


# ═══════════════════════════════════════════
# DRAWING PRIMITIVES
# ═══════════════════════════════════════════

def add_text(slide, left, top, width, height, text,
             font_size=BODY_SIZE, font_name=FONT_BODY,
             font_color=DARK_GRAY, bold=False,
             alignment=PP_ALIGN.LEFT, ea_font=FONT_EA,
             anchor=MSO_ANCHOR.TOP, line_spacing=Pt(6)):
    """Add a text box. Pass str for single line, list[str] for multi-line."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    anchor_map = {MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b', MSO_ANCHOR.TOP: 't'}
    bodyPr.set('anchor', anchor_map.get(anchor, 't'))
    for attr in ['lIns', 'tIns', 'rIns', 'bIns']:
        bodyPr.set(attr, '45720')
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.name = font_name
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = alignment
        p.space_before = line_spacing if i > 0 else Pt(0)
        p.space_after = Pt(0)
        # v1.10.3: titles (≥18pt) use 0.93× multiple spacing for tighter look;
        # body text uses fixed Pt spacing to prevent CJK overlap
        p.line_spacing = 0.93 if font_size.pt >= 18 else Pt(font_size.pt * 1.35)
        for run in p.runs:
            set_ea_font(run, ea_font)
    return txBox


def add_rect(slide, left, top, width, height, fill_color):
    """Add a flat rectangle with no border and no p:style."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    _clean_shape(shape)
    return shape


def add_hline(slide, x, y, length, color=BLACK, thickness=Pt(0.5)):
    """Draw a horizontal line as a thin rectangle (never use connectors)."""
    h = max(int(thickness), Emu(6350))
    return add_rect(slide, x, y, length, h, color)


def add_oval(slide, x, y, letter, size=Inches(0.45), bg=NAVY, fg=WHITE):
    """Add a circle label with a letter/number (e.g. 'A', '1')."""
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    c.fill.solid()
    c.fill.fore_color.rgb = bg
    c.line.fill.background()
    tf = c.text_frame
    tf.paragraphs[0].text = letter
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.name = 'Arial'
    tf.paragraphs[0].font.color.rgb = fg
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for run in tf.paragraphs[0].runs:
        set_ea_font(run, 'Arial')
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 'ctr')
    # 圆形内文字边距全部设为0，确保文字居中紧贴
    for attr in ['lIns', 'tIns', 'rIns', 'bIns']:
        bodyPr.set(attr, '0')
    _clean_shape(c)
    return c


def add_image_placeholder(slide, left, top, width, height, label='Image'):
    """Gray box with crosshair lines + centered label (user replaces with real image)."""
    PG = RGBColor(0xD9, 0xD9, 0xD9)
    rect = add_rect(slide, left, top, width, height, PG)
    add_hline(slide, left, top + height // 2, width, RGBColor(0xBB, 0xBB, 0xBB), Pt(0.5))
    vw = Pt(0.5)
    add_rect(slide, left + width // 2 - vw // 2, top, vw, height, RGBColor(0xBB, 0xBB, 0xBB))
    add_text(slide, left, top + height // 2 - Inches(0.2), width, Inches(0.4),
             f'[ {label} ]', font_size=SMALL_SIZE, font_color=RGBColor(0x99, 0x99, 0x99),
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return rect


# ═══════════════════════════════════════════
# SLIDE-LEVEL HELPERS (used by all layouts)
# ═══════════════════════════════════════════

def add_action_title(slide, text):
    """White background action title bar + separator line.
    Text is bottom-aligned so single-line titles sit flush against the separator.
    """
    add_text(slide, LM, TITLE_TOP, Inches(11.7), TITLE_H, text,
             font_size=ACTION_TITLE_SIZE, font_color=BLACK, bold=True,
             font_name=FONT_HEADER, anchor=MSO_ANCHOR.BOTTOM)
    add_hline(slide, LM, TITLE_LINE_Y, Inches(11.7), BLACK, Pt(0.5))


def add_source(slide, text, y=SOURCE_Y):
    """Footnote-sized source attribution at bottom."""
    add_text(slide, LM, y, Inches(11), Inches(0.3), text,
             font_size=FOOTNOTE_SIZE, font_color=MED_GRAY)


def add_page_number(slide, num, total):
    """Page number in bottom-right corner."""
    add_text(slide, PAGE_NUM_X, Inches(7.1), Inches(1), Inches(0.3),
             f"{num}/{total}", font_size=FOOTNOTE_SIZE, font_color=MED_GRAY,
             alignment=PP_ALIGN.RIGHT)


def add_bottom_bar(slide, label, text, y=BOTTOM_BAR_Y):
    """Gray summary bar at bottom with bold label + description."""
    add_rect(slide, LM, y, CW, BOTTOM_BAR_H, BG_GRAY)
    add_text(slide, LM + Inches(0.3), y, Inches(1.8), BOTTOM_BAR_H,
             label, font_size=BODY_SIZE, font_color=NAVY, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, LM + Inches(2.2), y, CW - Inches(2.5), BOTTOM_BAR_H,
             text, font_size=BODY_SIZE, font_color=DARK_GRAY,
             anchor=MSO_ANCHOR.MIDDLE)


# ═══════════════════════════════════════════
# BLOCK_ARC HELPER (v2.0 — charts engine)
# ═══════════════════════════════════════════

def add_block_arc(slide, cx, cy, outer_r, start_deg, sweep_deg, fill_color,
                  inner_ratio=None, ring_width=Pt(10)):
    """Draw one BLOCK_ARC segment (used for donut, pie, gauge).

    Parameters
    ----------
    cx, cy : Emu – center of the arc.
    outer_r : Emu – outer radius.
    start_deg : float – start angle in *math convention* (CCW from 3-o'clock).
    sweep_deg : float – sweep amount in degrees (positive = CCW).
    fill_color : RGBColor.
    inner_ratio : int|None – explicit adj3 (0-50000).  If None, auto-calculated
                  from ring_width.
    ring_width : Emu – thickness of the ring (only used when inner_ratio is None).
                 Set ring_width=outer_r for a solid pie segment.

    Returns the shape object.
    """
    d = outer_r * 2
    sh = slide.shapes.add_shape(
        MSO_SHAPE.BLOCK_ARC,
        cx - outer_r, cy - outer_r, d, d,
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    sh.line.fill.background()
    _clean_shape(sh)

    # Calculate PPT angles (CW from 12 o'clock, in 60000ths of degree)
    ppt_start = (90 - start_deg - sweep_deg) % 360
    ppt_end = (90 - start_deg) % 360

    if inner_ratio is None:
        inner_r = max(outer_r - ring_width, 0)
        inner_ratio = int((inner_r / outer_r) * 50000) if outer_r else 0

    sp = sh._element.find(qn('p:spPr'))
    prstGeom = sp.find(qn('a:prstGeom'))
    if prstGeom is not None:
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = prstGeom.makeelement(qn('a:avLst'), {})
            prstGeom.append(avLst)
        for gd in avLst.findall(qn('a:gd')):
            avLst.remove(gd)
        gd1 = avLst.makeelement(qn('a:gd'),
                                {'name': 'adj1', 'fmla': f'val {int(ppt_start * 60000)}'})
        gd2 = avLst.makeelement(qn('a:gd'),
                                {'name': 'adj2', 'fmla': f'val {int(ppt_end * 60000)}'})
        gd3 = avLst.makeelement(qn('a:gd'),
                                {'name': 'adj3', 'fmla': f'val {inner_ratio}'})
        avLst.append(gd1)
        avLst.append(gd2)
        avLst.append(gd3)
    return sh


def add_color_legend(slide, items, x, y, spacing=Inches(1.8),
                     swatch=Inches(0.2), font_size=SMALL_SIZE):
    """Add a horizontal color-legend row.
    items: list of (label, color).
    """
    for i, (label, color) in enumerate(items):
        lx = x + spacing * i
        add_rect(slide, lx, y + Pt(2), swatch, swatch, color)
        add_text(slide, lx + swatch + Pt(6), y, spacing - swatch - Pt(10),
                 Inches(0.3), label, font_size=font_size, font_color=DARK_GRAY)


def draw_harvey_ball(slide, x, y, score, size=Inches(0.35)):
    """Draw a Harvey Ball indicator (0=empty … 4=full).
    Uses white rectangle masking for clean quarter/half/three-quarter fills.
    """
    score = max(0, min(4, round(score)))
    # Background circle
    o = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    o.fill.solid()
    o.fill.fore_color.rgb = BG_GRAY
    o.line.color.rgb = NAVY
    o.line.width = Pt(1.0)
    _clean_shape(o)
    if score == 0:
        return
    # Full filled circle
    f = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    f.fill.solid()
    f.fill.fore_color.rgb = NAVY
    f.line.fill.background()
    _clean_shape(f)
    if score == 4:
        return
    ext = Pt(2)
    half = size // 2
    if score == 1:
        add_rect(slide, x + half, y - ext, half + ext, size + ext * 2, WHITE)
        add_rect(slide, x - ext, y + half, half + ext, half + ext, WHITE)
    elif score == 2:
        add_rect(slide, x + half, y - ext, half + ext, size + ext * 2, WHITE)
    elif score == 3:
        add_rect(slide, x + half, y - ext, half + ext, half + ext, WHITE)