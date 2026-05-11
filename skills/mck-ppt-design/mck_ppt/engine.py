# Copyright 2024-2026 Kaku Li (https://github.com/likaku)
# Licensed under the Apache License, Version 2.0 — see LICENSE and NOTICE.
# Part of "Mck-ppt-design-skill" (McKinsey PPT Design Framework).
# NOTICE: This file must be retained in all copies or substantial portions.
#
"""MckEngine — The presentation engine that wraps python-pptx with high-level layout methods.

Usage:
    eng = MckEngine(total_slides=30)
    eng.cover(title='Title', subtitle='Sub')
    eng.toc(items=[('1','Topic','Desc'), ...])
    eng.save('output/deck.pptx')

Every layout method creates one slide and auto-increments page numbers.
"""
import math
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from .constants import *
from .core import (
    _clean_shape, set_ea_font,
    add_text, add_rect, add_hline, add_oval, add_image_placeholder,
    add_action_title, add_source, add_page_number, add_bottom_bar,
    add_block_arc, add_color_legend, draw_harvey_ball,
    full_cleanup,
)


class MckEngine:
    """Presentation engine with high-level layout methods."""

    def __init__(self, total_slides=30):
        self.prs = Presentation()
        self.prs.slide_width = SW
        self.prs.slide_height = SH
        self._blank_layout = self.prs.slide_layouts[6]
        self._page = 0
        self.total = total_slides

    # ─── internal ──────────────────────────────
    def _ns(self):
        """Create new blank slide, increment page counter."""
        self._page += 1
        return self.prs.slides.add_slide(self._blank_layout)

    def _footer(self, s, source=None):
        """Add source + page number (common to most content slides)."""
        if source:
            add_source(s, source)
        add_page_number(s, self._page, self.total)

    # ═══════════════════════════════════════════
    # STRUCTURE LAYOUTS (#1, #5, #6, #7, #36)
    # ═══════════════════════════════════════════

    def cover(self, title, subtitle='', author='', date='', cover_image=None):
        """#1 Cover Slide — title, subtitle, author, date, accent line.

        Parameters
        ----------
        cover_image : str or None
            - None  : 不插入图片（默认，保持原有布局）
            - 'auto': 调用腾讯混元 API 自动生成封面图
            - 路径   : 直接使用指定的图片文件
        """
        s = self._ns()

        # ── 确定是否有封面图片 ────────────────────────────────
        img_path = None
        if cover_image == 'auto':
            from .cover_image import generate_cover_image
            img_path = generate_cover_image(title)
        elif cover_image and os.path.isfile(cover_image):
            img_path = cover_image

        # ── 图片全幅垫底（先添加，后续所有元素在其上方） ────
        if img_path:
            s.shapes.add_picture(img_path, 0, 0, SW, SH)

        # ── 顶部 navy 细线 ───────────────────────────────────
        add_rect(s, 0, 0, SW, Inches(0.05), NAVY)

        # ── 布局参数：有图片时文字收左 ───────────────────────
        if img_path:
            text_left = Inches(0.9)
            text_width = Inches(7.2)
        else:
            text_left = Inches(1)
            text_width = Inches(11)

        # ── 标题 ──────────────────────────────────────────────
        lines = title.split('\n') if isinstance(title, str) else title
        n_lines = len(lines) if isinstance(lines, list) else title.count('\n') + 1
        title_h = Inches(0.8 + 0.62 * max(n_lines - 1, 0))
        add_text(s, text_left, Inches(1.15), text_width, title_h,
                 title, font_size=COVER_TITLE_SIZE, font_name=FONT_HEADER,
                 font_color=NAVY, bold=True)
        sub_y = Inches(1.15) + title_h + Inches(0.24)
        if subtitle:
            add_text(s, text_left, sub_y, text_width, Inches(0.62),
                     subtitle, font_size=Pt(22), font_color=DARK_GRAY)
            sub_y += Inches(0.95)
        else:
            sub_y += Inches(0.2)
        y = sub_y + Inches(0.28)
        if author:
            add_text(s, text_left, y, text_width, Inches(0.34),
                     author, font_size=BODY_SIZE, font_color=MED_GRAY)
            y += Inches(0.67)
        if date:
            add_text(s, text_left, y, text_width, Inches(0.34),
                     date, font_size=BODY_SIZE, font_color=MED_GRAY)
        add_hline(s, text_left, Inches(6.8), Inches(4.6), NAVY, Pt(2))

        return s

    def section_divider(self, section_label, title, subtitle=''):
        """#5 Section Divider — navy left bar, large title."""
        s = self._ns()
        add_rect(s, 0, 0, Inches(0.6), SH, NAVY)
        add_text(s, Inches(1.2), Inches(2.0), Inches(10), Inches(0.8),
                 section_label, font_size=SUB_HEADER_SIZE,
                 font_color=MED_GRAY, font_name=FONT_HEADER)
        add_text(s, Inches(1.2), Inches(2.8), Inches(10), Inches(1.2),
                 title, font_size=SECTION_TITLE_SIZE, font_color=NAVY,
                 bold=True, font_name=FONT_HEADER)
        if subtitle:
            add_text(s, Inches(1.2), Inches(4.2), Inches(10), Inches(0.6),
                     subtitle, font_size=BODY_SIZE, font_color=DARK_GRAY)
        add_page_number(s, self._page, self.total)
        return s

    def toc(self, title='目录', items=None, source=''):
        """#6 Table of Contents — numbered items with descriptions.
        items: list of (num, title, description)
        """
        s = self._ns()
        add_action_title(s, title)
        iy = Inches(1.5)
        for num, item_title, desc in (items or []):
            add_oval(s, LM, iy, str(num))
            add_text(s, LM + Inches(0.7), iy, Inches(4.0), Inches(0.4),
                     item_title, font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True)
            add_text(s, Inches(5.5), iy + Inches(0.05), Inches(6.5), Inches(0.4),
                     desc, font_size=BODY_SIZE, font_color=MED_GRAY)
            iy += Inches(0.7)
            add_hline(s, LM, iy, CW, LINE_GRAY)
            iy += Inches(0.3)
        self._footer(s, source)
        return s

    def closing(self, title, message='', source_text=''):
        """#36 Closing / Thank You slide."""
        s = self._ns()
        add_rect(s, 0, 0, SW, Inches(0.05), NAVY)
        add_text(s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.0),
                 title, font_size=SECTION_TITLE_SIZE, font_color=NAVY,
                 bold=True, font_name=FONT_HEADER, alignment=PP_ALIGN.CENTER)
        add_hline(s, Inches(5.5), Inches(3.3), Inches(2.3), NAVY, Pt(1.5))
        if message:
            add_text(s, Inches(1.5), Inches(3.8), Inches(10.3), Inches(2.0),
                     message, font_size=SUB_HEADER_SIZE, font_color=DARK_GRAY,
                     alignment=PP_ALIGN.CENTER)
        add_hline(s, LM, Inches(6.8), CW, NAVY, Pt(2))
        if source_text:
            add_text(s, Inches(1), Inches(6.2), Inches(11), Inches(0.4),
                     source_text, font_size=SMALL_SIZE, font_color=MED_GRAY,
                     alignment=PP_ALIGN.CENTER)
        return s

    # ═══════════════════════════════════════════
    # DATA LAYOUTS (#8, #9, #10, #11, #12, #23)
    # ═══════════════════════════════════════════

    def big_number(self, title, number, unit='', description='',
                   detail_items=None, source='', bottom_bar=None):
        """#8 Big Number — large stat with context.
        detail_items: list[str] bullet points shown below.
        bottom_bar: (label, text) or None.
        """
        s = self._ns()
        add_action_title(s, title)
        # Navy big number box
        box_w = Inches(3.5)
        add_rect(s, LM, CONTENT_TOP + Inches(0.1), box_w, Inches(1.8), NAVY)
        add_text(s, LM + Inches(0.2), CONTENT_TOP + Inches(0.2), box_w - Inches(0.4), Inches(0.8),
                 str(number), font_size=COVER_TITLE_SIZE, font_color=WHITE, bold=True,
                 font_name=FONT_HEADER, alignment=PP_ALIGN.CENTER)
        if unit:
            add_text(s, LM + Inches(0.2), CONTENT_TOP + Inches(1.0), box_w - Inches(0.4), Inches(0.7),
                     unit, font_size=SMALL_SIZE, font_color=WHITE, alignment=PP_ALIGN.CENTER)
        # Right description
        if description:
            right_x = Inches(5.0)
            right_w = Inches(7.5)
            add_text(s, right_x, CONTENT_TOP + Inches(0.2), right_w, Inches(2.5),
                     description if isinstance(description, list) else [description],
                     font_size=BODY_SIZE, line_spacing=Pt(10))
        # Detail area
        if detail_items:
            add_rect(s, LM, Inches(4.5), CW, Inches(2.2), BG_GRAY)
            add_text(s, LM + Inches(0.3), Inches(4.6), Inches(1.8), Inches(0.4),
                     '解决路径' if not unit else '详细说明',
                     font_size=BODY_SIZE, font_color=NAVY, bold=True)
            add_text(s, LM + Inches(0.3), Inches(5.1), CW - Inches(0.6), Inches(1.4),
                     detail_items, font_size=BODY_SIZE, line_spacing=Pt(8))
        if bottom_bar:
            add_bottom_bar(s, bottom_bar[0], bottom_bar[1])
        self._footer(s, source)
        return s

    def two_stat(self, title, stats, detail_items=None, source=''):
        """#9 Two-Stat Comparison — two big numbers side by side.
        stats: list of (number, label, is_navy:bool)
        """
        s = self._ns()
        add_action_title(s, title)
        sw_stat = Inches(5.5)
        sg = Inches(0.733)
        for i, (big, label, is_navy) in enumerate(stats):
            sx = LM + (sw_stat + sg) * i
            fill = NAVY if is_navy else BG_GRAY
            bc = WHITE if is_navy else NAVY
            sc = WHITE if is_navy else DARK_GRAY
            add_rect(s, sx, Inches(1.5), sw_stat, Inches(2.0), fill)
            add_text(s, sx + Inches(0.3), Inches(1.6), sw_stat - Inches(0.6), Inches(0.9),
                     str(big), font_size=COVER_TITLE_SIZE, font_color=bc, bold=True,
                     font_name=FONT_HEADER, alignment=PP_ALIGN.CENTER)
            add_text(s, sx + Inches(0.3), Inches(2.6), sw_stat - Inches(0.6), Inches(0.5),
                     label, font_size=BODY_SIZE, font_color=sc, alignment=PP_ALIGN.CENTER)
        if detail_items:
            add_text(s, LM, Inches(4.0), CW, Inches(2.5),
                     detail_items, font_size=BODY_SIZE, line_spacing=Pt(8))
        self._footer(s, source)
        return s

    def metric_cards(self, title, cards, source=''):
        """#12 Metric Cards — 3-4 accent-colored cards.
        cards: list of (letter, card_title, description, accent_color, light_bg)
               or (letter, card_title, description) — auto-colors from ACCENT_PAIRS.
        """
        s = self._ns()
        add_action_title(s, title)
        n = len(cards)
        card_w = (CW - Inches(0.2) * (n - 1)) / n
        card_g = Inches(0.2)
        for i, card in enumerate(cards):
            if len(card) == 5:
                letter, ctitle, desc, accent, light = card
            else:
                letter, ctitle, desc = card[:3]
                accent, light = NAVY, BG_GRAY
            cx = LM + (card_w + card_g) * i
            add_rect(s, cx, CONTENT_TOP + Inches(0.1), card_w, Inches(4.8), light)
            add_rect(s, cx, CONTENT_TOP + Inches(0.1), card_w, Inches(0.06), accent)
            add_oval(s, cx + card_w / 2 - Inches(0.225), CONTENT_TOP + Inches(0.3),
                     str(letter), bg=accent)
            add_text(s, cx + Inches(0.2), CONTENT_TOP + Inches(0.9), card_w - Inches(0.4), Inches(0.4),
                     ctitle, font_size=SUB_HEADER_SIZE, font_color=accent,
                     bold=True, alignment=PP_ALIGN.CENTER)
            add_hline(s, cx + Inches(0.4), CONTENT_TOP + Inches(1.4),
                      card_w - Inches(0.8), LINE_GRAY)
            add_text(s, cx + Inches(0.2), CONTENT_TOP + Inches(1.6), card_w - Inches(0.4), Inches(2.5),
                     desc if isinstance(desc, list) else desc,
                     font_size=BODY_SIZE, alignment=PP_ALIGN.LEFT, line_spacing=Pt(8))
        self._footer(s, source)
        return s

    def data_table(self, title, headers, rows, col_widths=None, source='',
                   bottom_bar=None):
        """#11 Data Table — header row + data rows with separators.
        headers: list[str], rows: list[list[str]], col_widths: list[Inches] or auto.
        """
        s = self._ns()
        add_action_title(s, title)
        n = len(headers)
        if col_widths is None:
            col_widths = [CW / n] * n
        hdr_y = CONTENT_TOP + Inches(0.1)
        cx = LM
        for hdr, cw in zip(headers, col_widths):
            add_text(s, cx, hdr_y, cw, Inches(0.4), hdr,
                     font_size=BODY_SIZE, font_color=MED_GRAY, bold=True)
            cx += cw
        add_hline(s, LM, hdr_y + Inches(0.45), CW, BLACK, Pt(1.0))
        # --- adaptive row height ---
        row_start_y = hdr_y + Inches(0.55)
        bottom_limit = (BOTTOM_BAR_Y - Inches(0.15)) if bottom_bar else (SOURCE_Y - Inches(0.1))
        avail_h = bottom_limit - row_start_y
        n_rows = len(rows)
        row_h = min(Inches(0.95), avail_h / n_rows) if n_rows > 0 else Inches(0.95)
        # shrink font when rows are compact
        row_font = SMALL_SIZE if row_h >= Inches(0.6) else Pt(10)
        for ri, row in enumerate(rows):
            ry = row_start_y + row_h * ri
            cx = LM
            for val, cw in zip(row, col_widths):
                add_text(s, cx, ry, cw, row_h, val, font_size=row_font)
                cx += cw
            add_hline(s, LM, ry + row_h, CW, LINE_GRAY)
        if bottom_bar:
            add_bottom_bar(s, bottom_bar[0], bottom_bar[1])
        self._footer(s, source)
        return s

    def table_insight(self, title, headers, rows, insights,
                      col_widths=None, insight_title='启示：',
                      source='', bottom_bar=None):
        """Table + right insight panel — McKinsey editorial layout.

        Left ~60%: data table with header row + horizontal-line separated rows.
                   Each row is list[str] matching headers.
                   Supports **bold** markup within cell text.
        Middle: double-chevron arrow icon bridging table → insight.
        Right ~32%: "启示：" title + decorative line + bullet insights.

        Parameters
        ----------
        headers : list[str] — column headers for the table.
        rows : list[list[str]] — each inner list maps to headers.
        insights : list[str] — insight bullet points (shown on the right panel).
        col_widths : list[Inches] or None — custom widths for table columns.
        insight_title : str — title for the right panel (default '启示：').
        """
        import re
        s = self._ns()
        add_action_title(s, title)

        # ── Layout geometry ──
        table_w = Inches(7.2)               # left table area width
        chevron_zone = Inches(0.7)           # middle zone for chevron icon
        insight_w = CW - table_w - chevron_zone  # right insight panel width
        table_x = LM
        chevron_x = LM + table_w             # chevron zone left edge
        insight_x = LM + table_w + chevron_zone

        # ── NO vertical separator line — chevron icon separates visually ──

        # ── Table: header row ──
        n_cols = len(headers)
        if col_widths is None:
            col_widths = [table_w / n_cols] * n_cols
        hdr_y = CONTENT_TOP + Inches(0.1)
        cx = table_x
        for hdr, cw in zip(headers, col_widths):
            add_text(s, cx, hdr_y, cw, Inches(0.4), hdr,
                     font_size=BODY_SIZE, font_color=BLACK, bold=True)
            cx += cw
        add_hline(s, table_x, hdr_y + Inches(0.45), table_w, BLACK, Pt(1.0))

        # ── Table: data rows ──
        row_start_y = hdr_y + Inches(0.55)
        bottom_limit = (BOTTOM_BAR_Y - Inches(0.15)) if bottom_bar else (SOURCE_Y - Inches(0.1))
        avail_h = bottom_limit - row_start_y
        n_rows = len(rows)
        row_h = min(Inches(1.55), avail_h / n_rows) if n_rows > 0 else Inches(1.55)
        row_font = BODY_SIZE if row_h >= Inches(1.0) else SMALL_SIZE

        for ri, row in enumerate(rows):
            ry = row_start_y + row_h * ri
            cx = table_x
            for ci, (val, cw) in enumerate(zip(row, col_widths)):
                # First column (label): bold, vertically centered
                if ci == 0:
                    add_text(s, cx, ry, cw, row_h, val,
                             font_size=SUB_HEADER_SIZE, font_color=BLACK, bold=True,
                             anchor=MSO_ANCHOR.MIDDLE)
                else:
                    # Support **bold** markup in cell text
                    lines = val if isinstance(val, list) else val.split('\n')
                    txBox = s.shapes.add_textbox(cx, ry, cw, row_h)
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    tf.auto_size = None
                    bodyPr = tf._txBody.find(qn('a:bodyPr'))
                    bodyPr.set('anchor', 'ctr')
                    for attr in ['lIns', 'tIns', 'rIns', 'bIns']:
                        bodyPr.set(attr, '45720')
                    for li, line in enumerate(lines):
                        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                        p.space_before = Pt(3) if li > 0 else Pt(0)
                        p.space_after = Pt(0)
                        p.line_spacing = Pt(row_font.pt * 1.35)
                        segments = re.split(r'(\*\*.*?\*\*)', line)
                        for seg in segments:
                            if seg.startswith('**') and seg.endswith('**'):
                                run = p.add_run()
                                run.text = seg[2:-2]
                                run.font.size = row_font
                                run.font.name = FONT_BODY
                                run.font.color.rgb = BLACK
                                run.font.bold = True
                                set_ea_font(run, FONT_EA)
                            elif seg:
                                run = p.add_run()
                                run.text = seg
                                run.font.size = row_font
                                run.font.name = FONT_BODY
                                run.font.color.rgb = DARK_GRAY
                                run.font.bold = False
                                set_ea_font(run, FONT_EA)
                cx += cw
            # Row separator line
            add_hline(s, table_x, ry + row_h, table_w, LINE_GRAY)

        # ── Middle: double-chevron arrow icon ──
        # Vertically centered in the content area, between table and insights
        content_mid_y = CONTENT_TOP + avail_h * 0.42
        chev_size = Inches(0.5)
        chev_shape = s.shapes.add_shape(
            MSO_SHAPE.CHEVRON,
            chevron_x + (chevron_zone - chev_size) // 2,
            content_mid_y,
            chev_size, chev_size)
        chev_shape.fill.solid()
        chev_shape.fill.fore_color.rgb = DARK_GRAY
        chev_shape.line.fill.background()
        _clean_shape(chev_shape)

        # ── Right insight panel ──
        n_insights = len(insights)
        if n_insights > 0:
            insight_area_top = CONTENT_TOP + Inches(0.1)

            # Gray background rectangle for the entire insight panel
            bg_pad = Inches(0.1)  # small padding around the bg
            bg_bottom = bottom_limit + Inches(0.05)
            add_rect(s, insight_x - bg_pad, insight_area_top - bg_pad,
                     insight_w + bg_pad * 2, bg_bottom - insight_area_top + bg_pad * 2,
                     BG_GRAY)

            # "启示：" title (no decorative line above)
            title_y = insight_area_top + Inches(0.1)
            add_text(s, insight_x, title_y, insight_w, Inches(0.45),
                     insight_title,
                     font_size=SUB_HEADER_SIZE, font_color=BLACK, bold=True)

            # Bullet insights — compact, with round bullet •
            bullet_start_y = title_y + Inches(0.55)
            bullet_spacing = Inches(0.15)  # gap between bullets
            # Calculate available height for all insights
            insight_bottom = bottom_limit
            insight_avail = insight_bottom - bullet_start_y
            block_h = (insight_avail - bullet_spacing * (n_insights - 1)) / n_insights if n_insights > 1 else insight_avail

            for ii, ins in enumerate(insights):
                by = bullet_start_y + ii * (block_h + bullet_spacing)

                # Round bullet •
                add_text(s, insight_x, by,
                         Inches(0.25), block_h, '•',
                         font_size=BODY_SIZE, font_color=BLACK, bold=True)
                # Insight text
                add_text(s, insight_x + Inches(0.3), by,
                         insight_w - Inches(0.35), block_h, ins,
                         font_size=BODY_SIZE, font_color=DARK_GRAY,
                         line_spacing=Pt(8))

        if bottom_bar:
            add_bottom_bar(s, bottom_bar[0], bottom_bar[1])
        self._footer(s, source)
        return s

    def scorecard(self, title, items, source=''):
        """#23 Scorecard — items with progress bars.
        items: list of (name, score_str, pct_float_0_to_1)
        """
        s = self._ns()
        add_action_title(s, title)
        headers = ['技术领域', '评分', '成熟度']
        add_text(s, LM, CONTENT_TOP + Inches(0.1), Inches(4.0), Inches(0.4),
                 headers[0], font_size=BODY_SIZE, font_color=MED_GRAY, bold=True)
        add_text(s, Inches(5.0), CONTENT_TOP + Inches(0.1), Inches(1.5), Inches(0.4),
                 headers[1], font_size=BODY_SIZE, font_color=MED_GRAY, bold=True)
        add_text(s, Inches(7.0), CONTENT_TOP + Inches(0.1), Inches(5.5), Inches(0.4),
                 headers[2], font_size=BODY_SIZE, font_color=MED_GRAY, bold=True)
        add_hline(s, LM, CONTENT_TOP + Inches(0.55), CW, BLACK, Pt(1.0))
        bar_max = Inches(5.0)
        ry = CONTENT_TOP + Inches(0.7)
        for name, score, pct in items:
            add_text(s, LM, ry, Inches(4.0), Inches(0.5), name, font_size=BODY_SIZE)
            add_text(s, Inches(5.0), ry, Inches(1.5), Inches(0.5), score,
                     font_size=BODY_SIZE, font_color=NAVY, bold=True)
            add_rect(s, Inches(7.0), ry + Inches(0.05), bar_max, Inches(0.4), BG_GRAY)
            bar_color = NAVY if pct >= 0.7 else ACCENT_ORANGE if pct >= 0.5 else ACCENT_RED
            add_rect(s, Inches(7.0), ry + Inches(0.05), Inches(5.0 * pct), Inches(0.4), bar_color)
            ry += Inches(0.55)
            add_hline(s, LM, ry, CW, LINE_GRAY)
            ry += Inches(0.1)
        self._footer(s, source)
        return s

    # ═══════════════════════════════════════════
    # FRAMEWORK LAYOUTS (#13, #15, #16, #18)
    # ═══════════════════════════════════════════

    def matrix_2x2(self, title, quadrants, axis_labels=None, source='',
                   bottom_bar=None):
        """#13 2×2 Matrix — four quadrants.
        quadrants: list of 4 (label, bg_color, description).
        axis_labels: (x_label, y_label) or None.
        """
        s = self._ns()
        add_action_title(s, title)
        grid_l = LM + Inches(1.8)
        grid_t = Inches(1.45)
        cell_w = Inches(4.5)
        cell_h = Inches(2.0)
        cell_gap = Inches(0.15)
        if axis_labels:
            add_text(s, LM, grid_t + cell_h - Inches(0.25), Inches(1.6), Inches(0.7),
                     axis_labels[1], font_size=BODY_SIZE, font_color=NAVY,
                     bold=True, alignment=PP_ALIGN.CENTER)
            add_text(s, grid_l + cell_w - Inches(0.45), grid_t + 2 * cell_h + cell_gap + Inches(0.03),
                     Inches(3.6), Inches(0.28),
                     axis_labels[0], font_size=BODY_SIZE, font_color=NAVY,
                     bold=True, alignment=PP_ALIGN.CENTER)
        for qi, (label, bg, desc) in enumerate(quadrants):
            row, col = qi // 2, qi % 2
            qx = grid_l + col * (cell_w + cell_gap)
            qy = grid_t + row * (cell_h + cell_gap)
            add_rect(s, qx, qy, cell_w, cell_h, bg)
            add_text(s, qx + Inches(0.2), qy + Inches(0.1), cell_w - Inches(0.4), Inches(0.32),
                     label, font_size=BODY_SIZE, font_color=NAVY, bold=True)
            add_text(s, qx + Inches(0.2), qy + Inches(0.46), cell_w - Inches(0.4), cell_h - Inches(0.56),
                     desc, font_size=Pt(11), font_color=DARK_GRAY)
        if bottom_bar:
            bar_y = Inches(6.26)
            add_rect(s, LM, bar_y, CW, Inches(0.56), BG_GRAY)
            add_text(s, LM + Inches(0.28), bar_y, Inches(1.5), Inches(0.56),
                     bottom_bar[0], font_size=BODY_SIZE, font_color=NAVY,
                     bold=True, anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, LM + Inches(1.95), bar_y, CW - Inches(2.2), Inches(0.56),
                     bottom_bar[1], font_size=Pt(13), font_color=DARK_GRAY,
                     anchor=MSO_ANCHOR.MIDDLE)
        self._footer(s, source)
        return s

    def pyramid(self, title, levels, source='', bottom_bar=None,
                detail_rows=None, detail_headers=None):
        """#15 Staircase Evolution — ascending steps with connected staircase outline.

        Parameters
        ----------
        title : str
            Action title (can be long, narrative-style).
        levels : list of (label, description, icon)
            Each level is one step. Displayed as ascending staircase L→R.
            icon can be:
            - **PNG file path** (recommended): absolute or relative path to a
              transparent-background white-line icon PNG file. The icon will be
              overlaid on the navy circle with slight inset (0.08") for clean fit.
              Example: '/path/to/assets/icons/icon_person_bust.png'
            - **Text/Unicode**: single char, Unicode symbol, or short text
              shown in white on navy circle. Example: '⚙', '1', 'A'
            If omitted, defaults to sequential number (1, 2, 3...).
        source : str
        bottom_bar : tuple(label, text) — optional summary bar (usually omitted).
        detail_rows : list of (row_label, [col1_texts, col2_texts, ...])
            Optional structured detail table below the staircase.
            Each row has a bold label on the left and one cell per level column.
            Each cell can be str or list[str].
            - **Single-line cells** are rendered as plain text (no bullet prefix).
              Use this for complete, narrative-style sentences.
            - **Multi-line cells** (list with 2+ items) are rendered with bullet
              prefix '•  ' for each line.
        detail_headers : list of str — optional column headers for detail table.

        Icon Guidelines
        ---------------
        For best results, generate PNG icons (200×200px, transparent background,
        white strokes ~6px) using PIL/Pillow ImageDraw. Store under assets/icons/.
        The engine automatically detects .png/.svg file paths and switches from
        text-in-circle to image-overlay mode.
        """
        s = self._ns()
        add_action_title(s, title)

        n = len(levels)

        # ── Staircase geometry ──
        # Steps ascend from left-bottom to right-top.
        # Platform lines extend rightward to connect with next vertical riser.
        col_w = Inches(3.6)       # content width per step
        gap = Inches(0.2)         # logical gap between columns (for content only)
        total_w = col_w * n + gap * (n - 1)
        x_start = LM + (CW - total_w) / 2  # center horizontally

        # Vertical: step platform Y positions (ascending)
        if detail_rows:
            platform_base = Inches(3.55)   # Y of the lowest step platform
            step_rise = Inches(0.85)       # vertical rise per step
        else:
            platform_base = Inches(3.8)
            step_rise = Inches(0.9)

        # Line thickness — use raw EMU for pixel-perfect alignment
        line_thick_emu = Emu(19050)  # ~1.5pt = 19050 EMU
        icon_size = Inches(0.5)

        # Collect platform coordinates for drawing
        # x_content_left, x_content_right = content area for labels/descriptions
        # platform_y = Y position of platform line
        platforms = []
        for i in range(n):
            cx = x_start + (col_w + gap) * i
            py = platform_base - step_rise * i
            platforms.append((cx, cx + col_w, py))

        # ── Draw connected staircase outline ──
        # Key: platform lines extend to the next column's left edge (no gap in lines).
        # Vertical risers connect from one platform's Y down to the previous platform's Y.
        for i in range(n):
            x_l, x_r, py = platforms[i]

            # Horizontal platform line:
            # - For non-last steps: extend from x_l to next step's x_l (bridging the gap)
            # - For last step: just x_l to x_r (normal width)
            if i < n - 1:
                next_x_l = platforms[i + 1][0]
                hline_length = next_x_l - x_l
            else:
                hline_length = x_r - x_l
            add_hline(s, x_l, py, hline_length, NAVY, Pt(1.5))

            # Vertical riser: from current platform UP to next (higher) platform
            if i < n - 1:
                next_py = platforms[i + 1][2]  # higher step's Y
                riser_x = platforms[i + 1][0]  # at the next step's left edge
                riser_h = py - next_py          # height from higher to lower
                # Riser goes from next_py down to py (top = next_py)
                # Center the thin rect on the x position
                add_rect(s, riser_x, next_py, line_thick_emu, riser_h + line_thick_emu, NAVY)

        # ── Draw content for each step ──
        for i, level_data in enumerate(levels):
            label = level_data[0]
            desc = level_data[1]
            icon_val = level_data[2] if len(level_data) > 2 else str(i + 1)
            x_l, x_r, py = platforms[i]

            # ── Navy circle icon (above platform, left side) ──
            icon_x = x_l + Inches(0.15)
            icon_y = py - icon_size - Inches(0.2)

            # Check if icon_val is a file path (PNG image) or text
            icon_is_image = isinstance(icon_val, str) and (
                icon_val.endswith('.png') or icon_val.endswith('.svg')
            ) and os.path.isfile(icon_val)

            if icon_is_image:
                # Draw navy circle background (no text)
                add_oval(s, icon_x, icon_y, '', size=icon_size)
                # Overlay the icon image centered on the circle
                # Inset the image slightly so it fits inside the circle
                inset = Inches(0.08)
                img_size = icon_size - inset * 2
                s.shapes.add_picture(
                    icon_val,
                    icon_x + inset, icon_y + inset,
                    img_size, img_size,
                )
            else:
                add_oval(s, icon_x, icon_y, str(icon_val), size=icon_size)

            # ── Stage label (bold, to the right of icon, same row) ──
            label_x = icon_x + icon_size + Inches(0.15)
            label_y = icon_y + Inches(0.02)
            label_w = x_r - label_x - Inches(0.1)
            add_text(s, label_x, label_y, label_w, Inches(0.45),
                     label, font_size=SUB_HEADER_SIZE, font_color=BLACK,
                     bold=True, alignment=PP_ALIGN.LEFT)

            # ── Description text (below platform line) ──
            desc_x = x_l + Inches(0.25)
            desc_y = py + Inches(0.15)
            desc_w = col_w - Inches(0.5)
            add_text(s, desc_x, desc_y, desc_w, Inches(0.7),
                     desc, font_size=BODY_SIZE, font_color=DARK_GRAY,
                     alignment=PP_ALIGN.LEFT)

        # ── Detail table (optional, below staircase) ──
        if detail_rows:
            table_top = Inches(4.65)
            # Top separator line
            add_hline(s, LM, table_top - Inches(0.08), CW, BLACK, Pt(0.75))

            label_col_w = Inches(1.6)  # left label column width
            data_col_w = (CW - label_col_w) / n  # equal-width data columns
            row_h = Inches(0.9)  # height per row

            for r, (row_label, cells) in enumerate(detail_rows):
                ry = table_top + row_h * r + Inches(0.05) * r
                # Row label (bold, left-aligned, vertically centered)
                add_text(s, LM, ry, label_col_w, row_h,
                         row_label, font_size=BODY_SIZE, font_color=BLACK,
                         bold=True, anchor=MSO_ANCHOR.MIDDLE)
                # Data cells
                for c, cell_text in enumerate(cells):
                    cell_x = LM + label_col_w + data_col_w * c
                    txt = cell_text if isinstance(cell_text, list) else [cell_text]
                    # If single line, show as plain text; if multi-line, add bullets
                    if len(txt) == 1:
                        formatted = txt
                    else:
                        formatted = ['•  ' + line for line in txt]
                    add_text(s, cell_x + Inches(0.1), ry, data_col_w - Inches(0.2), row_h,
                             formatted, font_size=SMALL_SIZE, font_color=DARK_GRAY,
                             anchor=MSO_ANCHOR.MIDDLE)
                # Row separator
                if r < len(detail_rows) - 1:
                    sep_y = ry + row_h
                    add_hline(s, LM, sep_y, CW, LINE_GRAY, Pt(0.5))

        if bottom_bar:
            add_bottom_bar(s, bottom_bar[0], bottom_bar[1])
        self._footer(s, source)
        return s

    def process_chevron(self, title, steps, source='', bottom_bar=None):
        """#16 Process Chevron — horizontal step flow with arrows.
        steps: list of (label, step_title, description).
        Supports 2–7 steps with dynamic sizing (Guard Rail Rule 10).
        """
        s = self._ns()
        add_action_title(s, title)
        n = len(steps)
        # ── Guard Rail Rule 10: dynamic step_w to prevent negative gap ──
        MIN_GAP = Inches(0.35)          # minimum arrow gap
        ARROW_ZONE = Inches(0.4)        # width reserved for "→" glyph
        PREFERRED_W = Inches(2.6)       # ideal step box width
        max_step_w = (CW - MIN_GAP * max(n - 1, 1)) / max(n, 1)
        step_w = min(PREFERRED_W, max_step_w)
        gap = (CW - step_w * n) / max(n - 1, 1)
        # adaptive font: shrink sub-header when boxes are narrow
        sub_font = SUB_HEADER_SIZE if step_w >= Inches(2.0) else BODY_SIZE
        desc_font = BODY_SIZE if step_w >= Inches(1.8) else SMALL_SIZE
        for i, (label, stitle, desc) in enumerate(steps):
            sx = LM + (step_w + gap) * i
            is_last = (i == n - 1)
            fill = NAVY if is_last else BG_GRAY
            tc = WHITE if is_last else NAVY
            add_rect(s, sx, Inches(1.5), step_w, Inches(1.0), fill)
            add_oval(s, sx + Inches(0.1), Inches(1.55), str(label),
                     bg=WHITE if is_last else NAVY,
                     fg=NAVY if is_last else WHITE)
            add_text(s, sx + Inches(0.65), Inches(1.55), step_w - Inches(0.8), Inches(0.9),
                     stitle, font_size=sub_font, font_color=tc,
                     bold=True, anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, sx + Inches(0.1), Inches(2.7), step_w - Inches(0.2), Inches(2.0),
                     desc, font_size=desc_font, alignment=PP_ALIGN.CENTER)
            if i < n - 1:
                arrow_w = max(gap - Inches(0.1), Inches(0.2))
                add_text(s, sx + step_w + Inches(0.05), Inches(1.7), arrow_w, Inches(0.5),
                         '→', font_size=Pt(24), font_color=NAVY, bold=True,
                         alignment=PP_ALIGN.CENTER)
        if bottom_bar:
            add_bottom_bar(s, bottom_bar[0], bottom_bar[1])
        self._footer(s, source)
        return s

    def temple(self, title, roof_text, pillar_names, foundation_text,
               pillar_colors=None, source=''):
        """#18 Temple / House Framework — roof + pillars + foundation.
        pillar_names: list[str], pillar_colors: list[RGBColor] or auto.
        """
        s = self._ns()
        add_action_title(s, title)
        # Roof
        add_rect(s, LM, CONTENT_TOP + Inches(0.1), CW, Inches(0.8), NAVY)
        add_text(s, LM + Inches(0.3), CONTENT_TOP + Inches(0.1), CW - Inches(0.6), Inches(0.8),
                 roof_text, font_size=SUB_HEADER_SIZE, font_color=WHITE, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE, alignment=PP_ALIGN.CENTER)
        # Pillars
        n = len(pillar_names)
        if pillar_colors is None:
            pillar_colors = [NAVY] * n
        ppw = (CW - Inches(0.2) * (n - 1)) / n
        ppg = Inches(0.2)
        for i, (name, color) in enumerate(zip(pillar_names, pillar_colors)):
            ppx = LM + (ppw + ppg) * i
            add_rect(s, ppx, Inches(2.5), ppw, Inches(2.8), BG_GRAY)
            add_rect(s, ppx, Inches(2.5), ppw, Inches(0.06), color)
            add_text(s, ppx + Inches(0.15), Inches(2.7), ppw - Inches(0.3), Inches(0.8),
                     name, font_size=BODY_SIZE, font_color=color, bold=True,
                     alignment=PP_ALIGN.CENTER)
        # Foundation
        add_rect(s, LM, Inches(5.5), CW, Inches(0.8), NAVY)
        add_text(s, LM + Inches(0.3), Inches(5.5), CW - Inches(0.6), Inches(0.8),
                 foundation_text, font_size=BODY_SIZE, font_color=WHITE, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE, alignment=PP_ALIGN.CENTER)
        self._footer(s, source)
        return s

    # ═══════════════════════════════════════════
    # COMPARISON LAYOUTS (#19, #20, #65)
    # ═══════════════════════════════════════════

    def side_by_side(self, title, options, source=''):
        """#19 Side-by-Side Comparison — two columns with navy headers.
        options: list of 2 (option_title, points:list[str]).
        """
        s = self._ns()
        add_action_title(s, title)
        cw_col = Inches(5.5)
        cg = Inches(0.733)
        for i, (otitle, pts) in enumerate(options):
            cx = LM + (cw_col + cg) * i
            add_rect(s, cx, CONTENT_TOP + Inches(0.1), cw_col, Inches(0.6), NAVY)
            add_text(s, cx + Inches(0.15), CONTENT_TOP + Inches(0.1), cw_col - Inches(0.3), Inches(0.6),
                     otitle, font_size=SUB_HEADER_SIZE, font_color=WHITE, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE, alignment=PP_ALIGN.CENTER)
            add_rect(s, cx, CONTENT_TOP + Inches(0.7), cw_col, Inches(4.2), BG_GRAY)
            add_text(s, cx + Inches(0.3), CONTENT_TOP + Inches(0.9), cw_col - Inches(0.6), Inches(3.8),
                     pts, font_size=BODY_SIZE, line_spacing=Pt(8))
        self._footer(s, source)
        return s

    def before_after(self, title, before_title, before_points,
                     after_title, after_points, source='',
                     corner_label='', bottom_bar=None,
                     left_summary='', right_summary='',
                     right_summary_color=None):
        """#20 Before/After — vertical divider with black circle arrow.

        Layout: left half + black vertical line with > circle + right half.
        No background color blocks — clean white layout with data rows.

        Parameters
        ----------
        title : str
            Slide action title.
        before_title : str
            Left section subtitle.
        before_points : list[dict] or list[str]
            Left section content. If list[dict], each dict has keys:
              - 'label': str — row label (bold black, left column)
              - 'brand1': str — first brand/description (gray small text)
              - 'val1': str — first value (red bold large text)
              - 'brand2': str — second brand (optional)
              - 'val2': str — second value (optional)
              - 'extra': str — supplementary note below val1 (optional, gray small text)
            If list[str], falls back to simple bullet points.
        after_title : str
            Right section subtitle.
        after_points : list[dict] or list[str]
            Right section content. If list[dict], each dict has keys:
              - 'title': str — numbered title (e.g. '1. xxx', bold black)
              - 'desc': str — description text (gray)
              - 'cases': list[tuple(str, str)] — (name, performance) pairs,
                performance shown in black bold + underline
            If list[str], falls back to simple bullet points.
        source : str
            Source footnote text.
        corner_label : str
            Optional dashed corner label (e.g. 'Part II > 退潮').
        bottom_bar : tuple(str, str) or None
            Optional bottom bar (label, text), e.g. ('关键洞察', '...').
        left_summary : str
            Optional summary text below left section (dark gray bold).
        right_summary : str
            Optional summary text below right section (red bold).
        right_summary_color : RGBColor or None
            Color for right summary. Defaults to ACCENT_RED.
        """
        from pptx.enum.dml import MSO_LINE_DASH_STYLE

        s = self._ns()
        add_action_title(s, title)

        if right_summary_color is None:
            right_summary_color = ACCENT_RED

        # ── Optional dashed corner label ──
        if corner_label:
            _corner_x = Inches(11.2)
            _corner_y = Inches(0.15)
            _corner_w = Inches(1.8)
            _corner_h = Inches(0.45)
            _corner_shape = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, _corner_x, _corner_y, _corner_w, _corner_h)
            _corner_shape.fill.background()
            _corner_shape.line.color.rgb = MED_GRAY
            _corner_shape.line.width = Pt(1.0)
            _corner_shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
            _clean_shape(_corner_shape)
            add_text(s, _corner_x, _corner_y, _corner_w, _corner_h,
                     corner_label, font_size=SMALL_SIZE, font_color=MED_GRAY,
                     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # ── Layout parameters ──
        ct = CONTENT_TOP + Inches(0.15)
        left_x = LM
        left_w = Inches(5.5)
        divider_x = left_x + left_w
        divider_w = Inches(0.733)
        right_x = divider_x + divider_w
        right_w = Inches(5.5)

        # ── Vertical divider line + black circle with > ──
        line_center_x = divider_x + divider_w / 2
        vline_top = ct
        vline_bottom = Inches(6.1)
        vline_h = vline_bottom - vline_top

        # Vertical line (thin rectangle)
        add_rect(s, line_center_x - Pt(0.5), vline_top, Pt(1.0), vline_h, BLACK)

        # Black circle with > arrow (text in oval text_frame, 0 margins, Arial)
        circle_size = Inches(0.5)
        circle_y = vline_top + (vline_h - circle_size) / 2
        circle_x = line_center_x - circle_size / 2
        circle = s.shapes.add_shape(
            MSO_SHAPE.OVAL, circle_x, circle_y, circle_size, circle_size)
        circle.fill.solid()
        circle.fill.fore_color.rgb = BLACK
        circle.line.fill.background()
        _clean_shape(circle)
        ctf = circle.text_frame
        ctf.paragraphs[0].text = '>'
        ctf.paragraphs[0].font.size = Pt(22)
        ctf.paragraphs[0].font.name = 'Arial'
        ctf.paragraphs[0].font.color.rgb = WHITE
        ctf.paragraphs[0].font.bold = True
        ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
        for run in ctf.paragraphs[0].runs:
            set_ea_font(run, 'Arial')
        cbodyPr = ctf._txBody.find(qn('a:bodyPr'))
        cbodyPr.set('anchor', 'ctr')
        for attr in ['lIns', 'tIns', 'rIns', 'bIns']:
            cbodyPr.set(attr, '0')

        # ══════════════════════════════════
        # LEFT HALF
        # ══════════════════════════════════
        add_text(s, left_x, ct, left_w, Inches(0.5),
                 before_title, font_size=Pt(16), font_color=BLACK, bold=True)

        if before_points and isinstance(before_points[0], dict):
            # ── Structured data rows ──
            table_top = ct + Inches(0.6)
            label_w = Inches(0.8)
            data_x = left_x + label_w + Inches(0.15)
            data_w = left_w - label_w - Inches(0.15)
            row_h = Inches(1.0)

            for ri, row in enumerate(before_points):
                ry = table_top + row_h * ri

                # Row label (bold black text, no background, vertically centered)
                add_text(s, left_x, ry, label_w, row_h,
                         row.get('label', ''), font_size=Pt(13),
                         font_color=BLACK, bold=True,
                         alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

                # First group: brand name (gray) + value (red bold)
                brand_y = ry + Inches(0.1)
                add_text(s, data_x, brand_y, Inches(2.2), Inches(0.22),
                         row.get('brand1', ''), font_size=Pt(11),
                         font_color=DARK_GRAY)
                add_text(s, data_x, brand_y + Inches(0.22), Inches(2.2), Inches(0.35),
                         row.get('val1', ''), font_size=Pt(18),
                         font_color=ACCENT_RED, bold=True)

                # Second group (optional)
                if row.get('brand2'):
                    col2_x = data_x + Inches(2.4)
                    add_text(s, col2_x, brand_y, Inches(2.0), Inches(0.22),
                             row['brand2'], font_size=Pt(11),
                             font_color=DARK_GRAY)
                    add_text(s, col2_x, brand_y + Inches(0.22),
                             Inches(2.0), Inches(0.35),
                             row.get('val2', ''), font_size=Pt(18),
                             font_color=ACCENT_RED, bold=True)

                # Extra note (optional, gray small text)
                if row.get('extra'):
                    add_text(s, data_x, brand_y + Inches(0.55),
                             data_w, Inches(0.2),
                             row['extra'], font_size=Pt(10),
                             font_color=MED_GRAY)

                # Row separator line
                if ri < len(before_points) - 1:
                    add_hline(s, left_x, ry + row_h, left_w, LINE_GRAY, Pt(0.5))

            # Left summary
            if left_summary:
                summary_y = table_top + row_h * len(before_points) + Inches(0.1)
                add_text(s, left_x, summary_y, left_w, Inches(0.35),
                         left_summary, font_size=Pt(12),
                         font_color=DARK_GRAY, bold=True)
        else:
            # ── Simple bullet points fallback ──
            add_text(s, left_x, ct + Inches(0.6), left_w, Inches(4.5),
                     before_points, font_size=SMALL_SIZE,
                     font_color=DARK_GRAY, line_spacing=Pt(8))
            if left_summary:
                add_text(s, left_x, Inches(5.8), left_w, Inches(0.35),
                         left_summary, font_size=Pt(12),
                         font_color=DARK_GRAY, bold=True)

        # ══════════════════════════════════
        # RIGHT HALF
        # ══════════════════════════════════
        add_text(s, right_x, ct, right_w, Inches(0.5),
                 after_title, font_size=Pt(16), font_color=BLACK, bold=True)

        if after_points and isinstance(after_points[0], dict):
            # ── Structured formula cards ──
            r_top = ct + Inches(0.65)
            fy = r_top
            for fi, f in enumerate(after_points):
                # Formula title (bold black)
                add_text(s, right_x, fy, right_w, Inches(0.35),
                         f.get('title', ''), font_size=Pt(14),
                         font_color=BLACK, bold=True)

                # Description (gray)
                add_text(s, right_x + Inches(0.15), fy + Inches(0.35),
                         right_w - Inches(0.15), Inches(0.3),
                         f.get('desc', ''), font_size=Pt(11),
                         font_color=DARK_GRAY)

                # Case data (name in gray bold, performance in black bold + underline)
                case_x = right_x + Inches(0.15)
                for ci, (name, perf) in enumerate(f.get('cases', [])):
                    cx = case_x + Inches(2.5) * ci
                    add_text(s, cx, fy + Inches(0.65), Inches(1.0), Inches(0.3),
                             name, font_size=Pt(11),
                             font_color=DARK_GRAY, bold=True)
                    perf_box = add_text(
                        s, cx + Inches(1.0), fy + Inches(0.6),
                        Inches(1.4), Inches(0.35),
                        perf, font_size=Pt(16),
                        font_color=BLACK, bold=True)
                    # Add underline
                    for p in perf_box.text_frame.paragraphs:
                        for r in p.runs:
                            r.font.underline = True

                fy += Inches(1.15)

                # Formula separator line
                if fi < len(after_points) - 1:
                    add_hline(s, right_x, fy - Inches(0.2), right_w,
                              LINE_GRAY, Pt(0.5))

            # Right summary
            if right_summary:
                add_text(s, right_x, fy, right_w, Inches(0.35),
                         right_summary, font_size=Pt(12),
                         font_color=right_summary_color, bold=True)
        else:
            # ── Simple bullet points fallback ──
            add_text(s, right_x, ct + Inches(0.65), right_w, Inches(4.5),
                     after_points, font_size=SMALL_SIZE,
                     font_color=DARK_GRAY, line_spacing=Pt(8))
            if right_summary:
                add_text(s, right_x, Inches(5.8), right_w, Inches(0.35),
                         right_summary, font_size=Pt(12),
                         font_color=right_summary_color, bold=True)

        # ── Bottom bar ──
        if bottom_bar:
            add_bottom_bar(s, bottom_bar[0], bottom_bar[1])

        self._footer(s, source)
        return s

    def swot(self, title, quadrants, source=''):
        """#65 SWOT Analysis — 2×2 colored grid.
        quadrants: list of 4 (label, accent_color, light_bg, points:list[str]).
        """
        s = self._ns()
        add_action_title(s, title)
        cell_w = CW / 2 - Inches(0.1)
        cell_h = Inches(2.3)
        grid_t = Inches(1.2)
        for qi, (label, accent, bg, pts) in enumerate(quadrants):
            row, col = qi // 2, qi % 2
            qx = LM + col * (cell_w + Inches(0.15))
            qy = grid_t + row * (cell_h + Inches(0.1))
            add_rect(s, qx, qy, cell_w, cell_h, bg)
            add_rect(s, qx, qy, cell_w, Inches(0.06), accent)
            add_text(s, qx + Inches(0.2), qy + Inches(0.15), cell_w - Inches(0.4), Inches(0.35),
                     label, font_size=EMPHASIS_SIZE, font_color=accent, bold=True)
            add_text(s, qx + Inches(0.2), qy + Inches(0.55), cell_w - Inches(0.4), cell_h - Inches(0.7),
                     pts, font_size=BODY_SIZE, font_color=DARK_GRAY, line_spacing=Pt(6))
        self._footer(s, source)
        return s

    # ═══════════════════════════════════════════
    # NARRATIVE LAYOUTS (#24, #25, #27, #28)
    # ═══════════════════════════════════════════

    def executive_summary(self, title, headline, items, source=''):
        """#24 Executive Summary — navy headline + numbered items.
        headline: str, items: list of (num, item_title, description).
        """
        s = self._ns()
        add_action_title(s, title)
        add_rect(s, LM, CONTENT_TOP + Inches(0.1), CW, Inches(1.0), NAVY)
        add_text(s, LM + Inches(0.3), CONTENT_TOP + Inches(0.1), CW - Inches(0.6), Inches(1.0),
                 headline, font_size=SUB_HEADER_SIZE, font_color=WHITE, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
        iy = CONTENT_TOP + Inches(1.5)
        for num, ititle, desc in items:
            add_oval(s, LM, iy, str(num))
            add_text(s, LM + Inches(0.6), iy, Inches(3.5), Inches(0.4),
                     ititle, font_size=BODY_SIZE, font_color=NAVY, bold=True)
            add_text(s, Inches(5.0), iy, Inches(7.5), Inches(0.4),
                     desc, font_size=BODY_SIZE)
            iy += Inches(0.6)
            add_hline(s, LM, iy, CW, LINE_GRAY)
            iy += Inches(0.3)
        self._footer(s, source)
        return s

    def key_takeaway(self, title, left_text, takeaways, source=''):
        """#25 Key Takeaway — left analysis + right gray panel.
        left_text: list[str], takeaways: list[str].
        """
        s = self._ns()
        add_action_title(s, title)
        left_w = Inches(7.5)
        add_text(s, LM, CONTENT_TOP + Inches(0.1), left_w, Inches(0.4),
                 '协同机制分析', font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True)
        add_hline(s, LM, CONTENT_TOP + Inches(0.6), left_w, LINE_GRAY)
        add_text(s, LM, CONTENT_TOP + Inches(0.8), left_w, Inches(4.0),
                 left_text, font_size=BODY_SIZE, line_spacing=Pt(4))
        tk_x = Inches(9.0)
        tk_w = Inches(3.5)
        add_rect(s, tk_x, CONTENT_TOP + Inches(0.1), tk_w, Inches(5.2), BG_GRAY)
        add_text(s, tk_x + Inches(0.2), CONTENT_TOP + Inches(0.3), tk_w - Inches(0.4), Inches(0.4),
                 'Key Takeaways', font_size=BODY_SIZE, font_color=NAVY, bold=True)
        add_hline(s, tk_x + Inches(0.2), CONTENT_TOP + Inches(0.8), tk_w - Inches(0.4), LINE_GRAY)
        add_text(s, tk_x + Inches(0.2), CONTENT_TOP + Inches(1.0), tk_w - Inches(0.4), Inches(4.0),
                 takeaways, font_size=BODY_SIZE, line_spacing=Pt(10))
        self._footer(s, source)
        return s

    def four_column(self, title, items, source=''):
        """#28 Four-Column Overview — 4 vertical cards.
        items: list of (num, col_title, description:str_or_list).
        """
        s = self._ns()
        add_action_title(s, title)
        n = len(items)
        col_w = (CW - Inches(0.2) * (n - 1)) / n
        col_g = Inches(0.2)
        for i, (num, ctitle, desc) in enumerate(items):
            cx = LM + (col_w + col_g) * i
            add_rect(s, cx, CONTENT_TOP + Inches(0.1), col_w, Inches(4.8), BG_GRAY)
            add_oval(s, cx + col_w / 2 - Inches(0.225), CONTENT_TOP + Inches(0.25), str(num))
            add_text(s, cx + Inches(0.15), CONTENT_TOP + Inches(0.9), col_w - Inches(0.3), Inches(0.4),
                     ctitle, font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True,
                     alignment=PP_ALIGN.CENTER)
            add_hline(s, cx + Inches(0.3), CONTENT_TOP + Inches(1.4), col_w - Inches(0.6), LINE_GRAY)
            add_text(s, cx + Inches(0.15), CONTENT_TOP + Inches(1.6), col_w - Inches(0.3), Inches(3.0),
                     desc, font_size=BODY_SIZE, alignment=PP_ALIGN.CENTER, line_spacing=Pt(8))
        self._footer(s, source)
        return s

    # ═══════════════════════════════════════════
    # TIMELINE & STEPS (#29, #30)
    # ═══════════════════════════════════════════

    def timeline(self, title, milestones, source='', bottom_bar=None):
        """#29 Timeline / Roadmap — horizontal line with milestone nodes.
        milestones: list of (label, description).
        """
        s = self._ns()
        add_action_title(s, title)
        line_w = Inches(10.7)
        add_hline(s, LM + Inches(0.5), Inches(3.0), line_w, LINE_GRAY, Pt(2))
        n = len(milestones)
        spacing = line_w / max(n - 1, 1)
        for i, (label, desc) in enumerate(milestones):
            mx = LM + Inches(0.5) + spacing * i
            is_last = (i == n - 1)
            add_oval(s, mx - Inches(0.225), Inches(2.775), str(i + 1),
                     bg=NAVY if is_last else ACCENT_BLUE)
            add_text(s, mx - Inches(1.0), Inches(2.0), Inches(2.0), Inches(0.5),
                     label, font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True,
                     alignment=PP_ALIGN.CENTER)
            add_text(s, mx - Inches(1.0), Inches(3.5), Inches(2.0), Inches(1.5),
                     desc, font_size=BODY_SIZE, alignment=PP_ALIGN.CENTER)
        if bottom_bar:
            add_bottom_bar(s, bottom_bar[0], bottom_bar[1])
        self._footer(s, source)
        return s

    def vertical_steps(self, title, steps, source='', bottom_bar=None):
        """#30 Vertical Steps — top-down numbered steps.
        steps: list of (num, step_title, description).
        """
        s = self._ns()
        add_action_title(s, title)
        # --- adaptive step spacing ---
        n_steps = len(steps)
        start_y = CONTENT_TOP + Inches(0.1)
        bottom_limit = (BOTTOM_BAR_Y - Inches(0.15)) if bottom_bar else (SOURCE_Y - Inches(0.1))
        avail = bottom_limit - start_y
        step_h = min(Inches(1.1), avail / n_steps) if n_steps > 0 else Inches(1.1)
        row_h = step_h * 0.55   # text row portion
        gap_h = step_h * 0.45   # gap + hline portion
        use_small = step_h < Inches(0.85)
        for i, (num, stitle, desc) in enumerate(steps):
            iy = start_y + step_h * i
            add_oval(s, LM, iy + Inches(0.05), str(num))
            add_text(s, LM + Inches(0.6), iy, Inches(3.0), row_h,
                     stitle, font_size=SUB_HEADER_SIZE if not use_small else BODY_SIZE,
                     font_color=NAVY, bold=True)
            add_text(s, Inches(4.5), iy, Inches(8.0), row_h, desc,
                     font_size=BODY_SIZE if not use_small else SMALL_SIZE)
            add_hline(s, LM, iy + row_h + gap_h * 0.5, CW, LINE_GRAY)
        if bottom_bar:
            add_bottom_bar(s, bottom_bar[0], bottom_bar[1], y=BOTTOM_BAR_Y)
        self._footer(s, source)
        return s

    # ═══════════════════════════════════════════
    # DATA LAYOUTS — continued (#10, #22)
    # ═══════════════════════════════════════════

    def three_stat(self, title, stats, detail_items=None, source=''):
        """#10 Three-Stat — three big numbers in a row.
        stats: list of 3 (number, label, is_navy:bool)
        """
        s = self._ns()
        add_action_title(s, title)
        sw2 = Inches(3.5)
        sg2 = (CW - sw2 * 3) / 2
        for i, (big, label, is_navy) in enumerate(stats):
            sx = LM + (sw2 + sg2) * i
            fill = NAVY if is_navy else BG_GRAY
            bc = WHITE if is_navy else NAVY
            sc = WHITE if is_navy else DARK_GRAY
            add_rect(s, sx, Inches(1.4), sw2, Inches(1.8), fill)
            add_text(s, sx + Inches(0.2), Inches(1.5), sw2 - Inches(0.4), Inches(0.7),
                     str(big), font_size=SECTION_TITLE_SIZE, font_color=bc, bold=True,
                     font_name=FONT_HEADER, alignment=PP_ALIGN.CENTER)
            add_text(s, sx + Inches(0.2), Inches(2.25), sw2 - Inches(0.4), Inches(0.6),
                     label, font_size=SMALL_SIZE, font_color=sc, alignment=PP_ALIGN.CENTER)
        if detail_items:
            add_text(s, LM, Inches(3.8), CW, Inches(2.5),
                     detail_items, font_size=BODY_SIZE, line_spacing=Pt(10))
        self._footer(s, source)
        return s

    def rag_status(self, title, headers, rows, source=''):
        """#22 RAG Status — table with red/amber/green status dots.
        headers: list[str], rows: list of (name, status_color, *values, note).
        """
        s = self._ns()
        add_action_title(s, title)
        n = len(headers)
        # Auto col widths
        col_widths = [CW / n] * n
        hx = LM
        for h, w in zip(headers, col_widths):
            add_text(s, hx, Inches(1.5), w, Inches(0.4), h,
                     font_size=BODY_SIZE, font_color=MED_GRAY, bold=True)
            hx += w
        add_hline(s, LM, Inches(2.0), CW, BLACK, Pt(1.0))
        ry = Inches(2.2)
        for row in rows:
            name = row[0]
            status_color = row[1]
            values = row[2:-1]
            note = row[-1]
            add_text(s, LM, ry, col_widths[0], Inches(0.6), name, font_size=BODY_SIZE)
            add_oval(s, LM + col_widths[0] + Inches(0.3), ry + Inches(0.05), '',
                     size=Inches(0.35), bg=status_color)
            cx = LM + col_widths[0] + col_widths[1]
            for vi, val in enumerate(values):
                add_text(s, cx, ry, col_widths[2 + vi], Inches(0.6), val, font_size=BODY_SIZE)
                cx += col_widths[2 + vi]
            add_text(s, cx, ry, col_widths[-1], Inches(0.6), note, font_size=BODY_SIZE)
            ry += Inches(0.7)
            add_hline(s, LM, ry, CW, LINE_GRAY)
            ry += Inches(0.15)
        self._footer(s, source)
        return s

    # ═══════════════════════════════════════════
    # STRUCTURE — continued (#7)
    # ═══════════════════════════════════════════

    def appendix_title(self, title, subtitle=''):
        """#7 Appendix Title — centered title with accent lines."""
        s = self._ns()
        add_rect(s, 0, 0, SW, Inches(0.05), NAVY)
        add_text(s, Inches(1), Inches(2.5), Inches(11.3), Inches(1.0),
                 title, font_size=Pt(36), font_color=NAVY, bold=True,
                 font_name=FONT_HEADER, alignment=PP_ALIGN.CENTER)
        add_hline(s, Inches(5.5), Inches(3.8), Inches(2.3), NAVY, Pt(1.5))
        if subtitle:
            add_text(s, Inches(1), Inches(4.2), Inches(11.3), Inches(0.5),
                     subtitle, font_size=BODY_SIZE, font_color=MED_GRAY,
                     alignment=PP_ALIGN.CENTER)
        add_page_number(s, self._page, self.total)
        return s

    # ═══════════════════════════════════════════
    # FRAMEWORK — continued (#17, #31, #32)
    # ═══════════════════════════════════════════

    def venn(self, title, circles, overlap_label='', right_text=None, source=''):
        """#17 Venn Diagram — ⚠️ RETIRED (已废弃).
        This layout has been retired per design review.
        circles: list of (label, points:list[str], x, y, w, h) positioned rects.
        overlap_label: text for overlap zone.
        right_text: list[str] explanation on the right side.
        """
        s = self._ns()
        add_action_title(s, title)
        for i, circ in enumerate(circles):
            label, points = circ[0], circ[1]
            cx, cy, cw, ch = circ[2], circ[3], circ[4], circ[5]
            is_overlap = (i == len(circles) - 1 and overlap_label)
            fill = NAVY if is_overlap else BG_GRAY
            fc = WHITE if is_overlap else NAVY
            add_rect(s, cx, cy, cw, ch, fill)
            add_text(s, cx + Inches(0.2), cy + Inches(0.2), cw - Inches(0.4), Inches(0.4),
                     label, font_size=SUB_HEADER_SIZE, font_color=fc, bold=True)
            if points:
                add_text(s, cx + Inches(0.2), cy + Inches(0.7), cw - Inches(0.4), ch - Inches(0.9),
                         points, font_size=BODY_SIZE, font_color=fc if is_overlap else DARK_GRAY)
        if overlap_label and len(circles) > 2:
            pass  # overlap already handled as last circle
        if right_text:
            add_text(s, Inches(9.0), Inches(2.0), Inches(3.5), Inches(4.0),
                     right_text, font_size=BODY_SIZE, line_spacing=Pt(8))
        self._footer(s, source)
        return s

    def cycle(self, title, phases, right_panel=None, source=''):
        """#31 Cycle Diagram — rectangular nodes with arrows in a loop.
        phases: list of (label, x_inches, y_inches) — positioned boxes.
        right_panel: (panel_title, points:list[str]) or None.
        """
        s = self._ns()
        add_action_title(s, title)
        box_w, box_h = Inches(2.2), Inches(1.2)
        for i, (label, px, py) in enumerate(phases):
            fill = NAVY if i == 0 else BG_GRAY
            tc = WHITE if i == 0 else NAVY
            add_rect(s, LM + Inches(px), Inches(py), box_w, box_h, fill)
            add_text(s, LM + Inches(px) + Inches(0.1), Inches(py) + Inches(0.1),
                     box_w - Inches(0.2), box_h - Inches(0.2), label,
                     font_size=SUB_HEADER_SIZE, font_color=tc, bold=True,
                     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Arrows between phases (simplified directional)
        arrows = ['→', '↓', '←', '↑']
        n = len(phases)
        if n >= 4:
            positions = [
                (LM + Inches(phases[0][1]) + box_w + Inches(0.2), Inches(phases[0][2]) + Inches(0.3)),
                (LM + Inches(phases[1][1]) + box_w / 2 - Inches(0.15), Inches(phases[1][2]) + box_h + Inches(0.1)),
                (LM + Inches(phases[2][1]) - Inches(0.5), Inches(phases[2][2]) + Inches(0.3)),
                (LM + Inches(phases[3][1]) + box_w / 2 - Inches(0.15), Inches(phases[3][2]) - Inches(0.5)),
            ]
            for i, (ax, ay) in enumerate(positions[:n]):
                add_text(s, ax, ay, Inches(0.6), Inches(0.5), arrows[i % 4],
                         font_size=Pt(24), font_color=NAVY, alignment=PP_ALIGN.CENTER)
        if right_panel:
            pt, pp = right_panel
            add_rect(s, Inches(8.5), Inches(1.5), Inches(4.0), Inches(5.0), BG_GRAY)
            add_text(s, Inches(8.8), Inches(1.7), Inches(3.4), Inches(0.4),
                     pt, font_size=BODY_SIZE, font_color=NAVY, bold=True)
            add_text(s, Inches(8.8), Inches(2.3), Inches(3.4), Inches(3.5),
                     pp, font_size=BODY_SIZE, line_spacing=Pt(10))
        self._footer(s, source)
        return s

    def funnel(self, title, stages, source=''):
        """#32 Funnel — ⚠️ RETIRED (已废弃).
        This layout has been retired per design review.
        stages: list of (name, count_label, pct_float).
        """
        s = self._ns()
        add_action_title(s, title)
        max_w = Inches(8.0)
        fy = Inches(1.6)
        for i, (name, count, pct) in enumerate(stages):
            w = int(max_w * pct)
            fx = Inches(6.666) - w // 2
            fill = NAVY if i == 0 else BG_GRAY
            tc = WHITE if i == 0 else NAVY
            add_rect(s, fx, fy, w, Inches(1.0), fill)
            add_text(s, fx + Inches(0.2), fy, w - Inches(0.4), Inches(1.0), name,
                     font_size=SUB_HEADER_SIZE, font_color=tc, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE, alignment=PP_ALIGN.CENTER)
            add_text(s, fx + w + Inches(0.3), fy + Inches(0.2), Inches(3.0), Inches(0.5),
                     f'{count} ({int(pct * 100)}%)',
                     font_size=BODY_SIZE, font_color=NAVY, bold=True)
            fy += Inches(1.2)
        self._footer(s, source)
        return s

    # ═══════════════════════════════════════════
    # NARRATIVE — continued (#21, #26, #27, #33-#35)
    # ═══════════════════════════════════════════

    def pros_cons(self, title, pros_title, pros, cons_title, cons,
                  conclusion=None, source=''):
        """#21 Pros/Cons — two-column layout.
        pros, cons: list[str].
        conclusion: (label, text) or None.
        """
        s = self._ns()
        add_action_title(s, title)
        hw = Inches(5.5)
        add_text(s, LM, Inches(1.5), hw, Inches(0.4),
                 pros_title, font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True)
        add_hline(s, LM, Inches(2.0), hw, NAVY)
        add_text(s, LM, Inches(2.2), hw, Inches(2.5),
                 pros, font_size=BODY_SIZE, line_spacing=Pt(10))
        cx = LM + hw + Inches(0.733)
        add_text(s, cx, Inches(1.5), hw, Inches(0.4),
                 cons_title, font_size=SUB_HEADER_SIZE, font_color=DARK_GRAY, bold=True)
        add_hline(s, cx, Inches(2.0), hw, DARK_GRAY)
        add_text(s, cx, Inches(2.2), hw, Inches(2.5),
                 cons, font_size=BODY_SIZE, line_spacing=Pt(10))
        if conclusion:
            label, text = conclusion
            add_rect(s, LM, Inches(5.2), CW, Inches(1.5), BG_GRAY)
            add_text(s, LM + Inches(0.3), Inches(5.3), Inches(1.5), Inches(0.4),
                     label, font_size=BODY_SIZE, font_color=NAVY, bold=True)
            add_text(s, LM + Inches(0.3), Inches(5.8), CW - Inches(0.6), Inches(0.6),
                     text, font_size=BODY_SIZE)
        self._footer(s, source)
        return s

    def quote(self, quote_text, attribution=''):
        """#26 Quote Slide — centered quote with accent lines."""
        s = self._ns()
        add_rect(s, 0, 0, SW, Inches(0.05), NAVY)
        add_hline(s, Inches(5.5), Inches(2.0), Inches(2.3), NAVY, Pt(1.5))
        add_text(s, Inches(1.5), Inches(2.5), Inches(10.3), Inches(2.5),
                 quote_text, font_size=Pt(24), font_color=DARK_GRAY,
                 alignment=PP_ALIGN.CENTER)
        add_hline(s, Inches(5.5), Inches(5.3), Inches(2.3), NAVY, Pt(1.5))
        if attribution:
            add_text(s, Inches(1.5), Inches(5.6), Inches(10.3), Inches(0.5),
                     attribution, font_size=BODY_SIZE, font_color=MED_GRAY,
                     alignment=PP_ALIGN.CENTER)
        add_page_number(s, self._page, self.total)
        return s

    def two_column_text(self, title, columns, source=''):
        """#27 Two-Column Text — lettered columns with bullet lists.
        columns: list of 2 (letter, col_title, points:list[str]).
        """
        s = self._ns()
        add_action_title(s, title)
        cw_col = Inches(5.5)
        cg = Inches(0.733)
        for i, (letter, ctitle, points) in enumerate(columns):
            cx = LM + (cw_col + cg) * i
            add_oval(s, cx, Inches(1.5), letter)
            add_text(s, cx + Inches(0.6), Inches(1.5), cw_col - Inches(0.6), Inches(0.4),
                     ctitle, font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True)
            add_hline(s, cx, Inches(2.0), cw_col, LINE_GRAY)
            add_text(s, cx, Inches(2.2), cw_col, Inches(4.0),
                     points, font_size=BODY_SIZE, line_spacing=Pt(10))
        self._footer(s, source)
        return s

    def meet_the_team(self, title, members, source=''):
        """#33 Meet the Team — profile cards in a row.
        members: list of (name, role, bio:str_or_list).
        """
        s = self._ns()
        add_action_title(s, title)
        n = len(members)
        cw_card = (CW - Inches(0.2) * (n - 1)) / n
        cg = Inches(0.2)
        for i, (name, role, bio) in enumerate(members):
            cx = LM + (cw_card + cg) * i
            add_rect(s, cx, Inches(1.5), cw_card, Inches(5.0), BG_GRAY)
            add_oval(s, cx + cw_card / 2 - Inches(0.5), Inches(1.7), name[0],
                     size=Inches(1.0))
            add_text(s, cx + Inches(0.15), Inches(2.9), cw_card - Inches(0.3), Inches(0.4),
                     name, font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True,
                     alignment=PP_ALIGN.CENTER)
            add_text(s, cx + Inches(0.15), Inches(3.4), cw_card - Inches(0.3), Inches(0.4),
                     role, font_size=BODY_SIZE, font_color=MED_GRAY,
                     alignment=PP_ALIGN.CENTER)
            add_hline(s, cx + Inches(0.3), Inches(3.9), cw_card - Inches(0.6), LINE_GRAY)
            bio_items = bio.split('\n') if isinstance(bio, str) else bio
            add_text(s, cx + Inches(0.15), Inches(4.1), cw_card - Inches(0.3), Inches(2.0),
                     bio_items, font_size=BODY_SIZE, line_spacing=Pt(8),
                     alignment=PP_ALIGN.CENTER)
        self._footer(s, source)
        return s

    def case_study(self, title, sections, result_box=None, source=''):
        """#34 Case Study — S/A/R or custom sections.
        sections: list of (letter, section_title, description).
        result_box: (label, text) or None.
        """
        s = self._ns()
        add_action_title(s, title)
        n = len(sections)
        sw_sec = (CW - Inches(0.2) * (n - 1)) / n
        sg = Inches(0.2)
        for i, (letter, stitle, desc) in enumerate(sections):
            sx = LM + (sw_sec + sg) * i
            is_last = (i == n - 1)
            fill = NAVY if is_last else BG_GRAY
            tc = WHITE if is_last else NAVY
            dc = WHITE if is_last else DARK_GRAY
            add_rect(s, sx, Inches(1.5), sw_sec, Inches(3.0), fill)
            add_oval(s, sx + Inches(0.15), Inches(1.65), letter,
                     bg=WHITE if is_last else NAVY,
                     fg=NAVY if is_last else WHITE)
            add_text(s, sx + Inches(0.15), Inches(2.2), sw_sec - Inches(0.3), Inches(0.8),
                     stitle, font_size=BODY_SIZE, font_color=tc, bold=True,
                     alignment=PP_ALIGN.CENTER)
            desc_items = desc.split('\n') if isinstance(desc, str) else desc
            add_text(s, sx + Inches(0.15), Inches(3.1), sw_sec - Inches(0.3), Inches(1.0),
                     desc_items, font_size=BODY_SIZE, font_color=dc,
                     alignment=PP_ALIGN.CENTER)
        if result_box:
            label, text = result_box
            add_rect(s, LM, Inches(5.0), CW, Inches(1.5), BG_GRAY)
            add_text(s, LM + Inches(0.3), Inches(5.1), Inches(1.5), Inches(0.4),
                     label, font_size=BODY_SIZE, font_color=NAVY, bold=True)
            add_text(s, LM + Inches(0.3), Inches(5.6), CW - Inches(0.6), Inches(0.6),
                     text, font_size=BODY_SIZE)
        self._footer(s, source)
        return s

    def action_items(self, title, actions, source=''):
        """#35 Action Items — cards with timeline + owner.
        actions: list of (action_title, timeline, description, owner).
        """
        s = self._ns()
        add_action_title(s, title)
        n = len(actions)
        cw_card = (CW - Inches(0.2) * (n - 1)) / n
        cg = Inches(0.2)
        for i, (atitle, timeline, desc, owner) in enumerate(actions):
            cx = LM + (cw_card + cg) * i
            add_rect(s, cx, Inches(1.5), cw_card, Inches(0.6), NAVY)
            add_text(s, cx + Inches(0.15), Inches(1.5), cw_card - Inches(0.3), Inches(0.6),
                     atitle, font_size=BODY_SIZE, font_color=WHITE, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE, alignment=PP_ALIGN.CENTER)
            add_rect(s, cx, Inches(2.1), cw_card, Inches(0.4), BG_GRAY)
            add_text(s, cx + Inches(0.15), Inches(2.1), cw_card - Inches(0.3), Inches(0.4),
                     timeline, font_size=BODY_SIZE, font_color=NAVY, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE, alignment=PP_ALIGN.CENTER)
            desc_items = desc.split('\n') if isinstance(desc, str) else desc
            add_text(s, cx + Inches(0.15), Inches(2.7), cw_card - Inches(0.3), Inches(2.0),
                     desc_items, font_size=BODY_SIZE, line_spacing=Pt(8),
                     alignment=PP_ALIGN.CENTER)
            add_hline(s, cx + Inches(0.3), Inches(4.9), cw_card - Inches(0.6), LINE_GRAY)
            add_text(s, cx + Inches(0.15), Inches(5.1), cw_card - Inches(0.3), Inches(0.4),
                     f'负责人：{owner}', font_size=BODY_SIZE, font_color=MED_GRAY,
                     alignment=PP_ALIGN.CENTER)
        self._footer(s, source)
        return s

    # ═══════════════════════════════════════════
    # ADVANCED (#40, #54, #61, #67)
    # ═══════════════════════════════════════════

    def content_right_image(self, title, subtitle, bullets, takeaway='',
                            image_label='Image', source=''):
        """#40 Content + Right Image."""
        s = self._ns()
        add_action_title(s, title)
        left_w = Inches(6.5)
        add_text(s, LM, CONTENT_TOP, left_w, Inches(0.4),
                 subtitle, font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True)
        add_text(s, LM, CONTENT_TOP + Inches(0.5), left_w, Inches(3.5),
                 bullets, font_size=BODY_SIZE, line_spacing=Pt(10))
        if takeaway:
            add_rect(s, LM, Inches(5.5), left_w, Inches(0.8), BG_GRAY)
            add_text(s, LM + Inches(0.3), Inches(5.5), left_w - Inches(0.6), Inches(0.8),
                     takeaway, font_size=BODY_SIZE, font_color=NAVY, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE)
        img_x = LM + left_w + Inches(0.3)
        img_w = CW - left_w - Inches(0.3)
        add_image_placeholder(s, img_x, CONTENT_TOP, img_w, Inches(4.2), image_label)
        self._footer(s, source)
        return s

    def checklist(self, title, columns, col_widths, rows, status_map=None,
                  source='', bottom_bar=None):
        """#61 Checklist / Status table.
        columns: list[str] header labels.
        col_widths: list[Inches].
        rows: list of tuples — last element is status key.
        status_map: dict of status_key → (label, color, bg_color).
        """
        s = self._ns()
        add_action_title(s, title)
        if status_map is None:
            status_map = {
                'active':  ('→ 活跃', ACCENT_GREEN, LIGHT_GREEN),
                'risk':    ('△ 困难', ACCENT_ORANGE, LIGHT_ORANGE),
                'pending': ('○ 早期', MED_GRAY, BG_GRAY),
                'done':    ('✓ 完成', ACCENT_BLUE, LIGHT_BLUE),
            }
        hx = LM
        hy = CONTENT_TOP + Inches(0.1)
        for label, w in zip(columns, col_widths):
            add_text(s, hx, hy, w, Inches(0.35), label,
                     font_size=SMALL_SIZE, font_color=NAVY, bold=True)
            hx += w
        add_hline(s, LM, hy + Inches(0.35), CW, BLACK, Pt(0.75))
        # Dynamic row height: fit all rows between header and bottom area
        bottom_limit = (BOTTOM_BAR_Y - Inches(0.1)) if bottom_bar else (SOURCE_Y - Inches(0.05))
        available_h = bottom_limit - (hy + Inches(0.5))
        row_h_raw = available_h / max(len(rows), 1)
        row_h = min(Inches(0.85), row_h_raw)  # cap at 0.85" max
        # Use smaller font when rows are tight
        row_font = SMALL_SIZE if row_h < Inches(0.65) else BODY_SIZE
        data_widths = col_widths[:-1]  # last col is status
        status_w = col_widths[-1]
        for i, row in enumerate(rows):
            *data_vals, status_key = row
            ry = hy + Inches(0.5) + i * row_h
            st_label, st_color, st_bg = status_map.get(status_key, ('?', MED_GRAY, BG_GRAY))
            if i % 2 == 0:
                add_rect(s, LM, ry, CW, row_h, RGBColor(0xFA, 0xFA, 0xFA))
            rx = LM
            for val, w in zip(data_vals, data_widths):
                add_text(s, rx, ry, w, row_h, val, font_size=row_font,
                         anchor=MSO_ANCHOR.MIDDLE)
                rx += w
            add_rect(s, rx + Inches(0.1), ry + Inches(0.12), status_w - Inches(0.2), row_h - Inches(0.24), st_bg)
            add_text(s, rx + Inches(0.1), ry, status_w - Inches(0.2), row_h, st_label,
                     font_size=SMALL_SIZE, font_color=st_color, bold=True,
                     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if bottom_bar:
            add_bottom_bar(s, bottom_bar[0], bottom_bar[1])
        self._footer(s, source)
        return s

    def value_chain(self, title, stages, source='', bottom_bar=None):
        """#67 Value Chain / Horizontal Flow — stages with arrows.
        stages: list of (stage_title, description, accent_color).
        Stages fill the full content width; height fills available vertical space.
        """
        s = self._ns()
        add_action_title(s, title)
        n = len(stages)
        arrow_w = Inches(0.35)
        # Calculate stage width to fill entire content area
        stage_w = (CW - arrow_w * (n - 1)) / n
        start_x = LM
        stage_y = CONTENT_TOP + Inches(0.1)
        # Stage height fills down to bottom_bar or source area
        stage_h = (BOTTOM_BAR_Y - Inches(0.15) - stage_y) if bottom_bar else (SOURCE_Y - Inches(0.15) - stage_y)
        desc_h = stage_h - Inches(0.9)  # space for title row + padding
        for i, (stitle, desc, color) in enumerate(stages):
            sx = start_x + i * (stage_w + arrow_w)
            add_rect(s, sx, stage_y, stage_w, stage_h, WHITE)
            add_rect(s, sx, stage_y, stage_w, Inches(0.06), color)
            add_oval(s, sx + Inches(0.1), stage_y + Inches(0.2), str(i + 1),
                     size=Inches(0.4), bg=color)
            add_text(s, sx + Inches(0.6), stage_y + Inches(0.2), stage_w - Inches(0.7), Inches(0.4),
                     stitle, font_size=BODY_SIZE, font_color=color, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, sx + Inches(0.1), stage_y + Inches(0.8), stage_w - Inches(0.2), desc_h,
                     desc, font_size=SMALL_SIZE, alignment=PP_ALIGN.CENTER)
            if i < n - 1:
                ax = sx + stage_w + Inches(0.02)
                add_text(s, ax, stage_y + stage_h / 2 - Inches(0.2), arrow_w - Inches(0.05), Inches(0.4),
                         '→', font_size=Pt(20), font_color=LINE_GRAY, alignment=PP_ALIGN.CENTER)
        if bottom_bar:
            add_bottom_bar(s, bottom_bar[0], bottom_bar[1])
        self._footer(s, source)
        return s

    # ═══════════════════════════════════════════
    # CHART LAYOUTS (#37, #38, #39)
    # ═══════════════════════════════════════════

    def grouped_bar(self, title, categories, series, data, max_val=None,
                    y_ticks=None, summary=None, source=''):
        """#37 Grouped Bar Chart — vertical bars grouped by category.
        categories: list[str] x-labels. series: list of (name, color).
        data: list[list[int]] — data[cat_idx][series_idx].
        summary: (label, text) or None.
        """
        s = self._ns()
        add_action_title(s, title)
        # Layout: summary sinks to bottom, chart fills space above
        sum_h = Inches(0.7)
        sum_y = SOURCE_Y - sum_h - Inches(0.05)
        cb = sum_y - Inches(0.4) if summary else SOURCE_Y - Inches(0.3)
        cl = LM + Inches(0.8); ct = Inches(1.6); ch = cb - ct
        cr = Inches(11.5); cww = cr - cl
        if max_val is None:
            max_val = max(max(row) for row in data) * 1.15
        # Chart sub-title + legend
        add_text(s, cl, Inches(1.2), Inches(5.0), Inches(0.3),
                 title, font_size=Pt(13), font_color=DARK_GRAY, bold=True)
        for ci, (sname, scolor) in enumerate(series):
            lx = LM + Inches(8.0) + Inches(1.8) * ci
            add_rect(s, lx, Inches(1.25), Inches(0.2), Inches(0.2), scolor)
            add_text(s, lx + Inches(0.3), Inches(1.2), Inches(1.2), Inches(0.3),
                     sname, font_size=SMALL_SIZE, font_color=DARK_GRAY)
        # Y-axis ticks
        if y_ticks:
            for tk in y_ticks:
                ty = cb - int(ch * tk / max_val)
                add_text(s, LM, ty - Inches(0.15), Inches(0.7), Inches(0.3),
                         str(tk), font_size=FOOTNOTE_SIZE, font_color=MED_GRAY,
                         alignment=PP_ALIGN.RIGHT)
                if tk > 0:
                    add_hline(s, cl, ty, cww, LINE_GRAY, Pt(0.25))
        add_hline(s, cl, cb, cww, BLACK, Pt(0.5))
        # Bars
        nd = len(categories); nc = len(series)
        gw = cww / nd; bw = Inches(0.35); bgp = Inches(0.05)
        gbw = bw * nc + bgp * (nc - 1)
        for di, cat in enumerate(categories):
            gx = cl + gw * di + (gw - gbw) / 2
            for ci in range(nc):
                val = data[di][ci]
                bh = int(ch * val / max_val)
                bx = gx + (bw + bgp) * ci; by = cb - bh
                if val > 0:
                    add_rect(s, bx, by, bw, bh, series[ci][1])
                    if val >= 50:
                        add_text(s, bx - Inches(0.05), by - Inches(0.25),
                                 bw + Inches(0.1), Inches(0.25), str(val),
                                 font_size=FOOTNOTE_SIZE, font_color=DARK_GRAY,
                                 alignment=PP_ALIGN.CENTER)
            add_text(s, cl + gw * di, cb + Inches(0.05), gw, Inches(0.3),
                     cat, font_size=BODY_SIZE, font_color=DARK_GRAY,
                     alignment=PP_ALIGN.CENTER)
        if summary:
            add_rect(s, LM, sum_y, CW, sum_h, BG_GRAY)
            add_text(s, LM + Inches(0.3), sum_y, Inches(1.5), sum_h,
                     summary[0], font_size=BODY_SIZE, font_color=NAVY, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, LM + Inches(2.0), sum_y, CW - Inches(2.3), sum_h,
                     summary[1], font_size=BODY_SIZE, font_color=DARK_GRAY,
                     anchor=MSO_ANCHOR.MIDDLE)
        self._footer(s, source)
        return s

    def stacked_bar(self, title, periods, series, data, summary=None, source=''):
        """#38 Stacked Bar Chart — 100% stacked vertical bars.
        periods: list[str] x-labels. series: list of (name, color).
        data: list[list[int]] — percentages, data[period_idx][series_idx].
        summary: (label, text) or None.
        """
        s = self._ns()
        add_action_title(s, title)
        # Layout: summary sinks to bottom, chart fills space above
        sum_h = Inches(0.7)
        sum_y = SOURCE_Y - sum_h - Inches(0.05)
        cb = sum_y - Inches(0.4) if summary else SOURCE_Y - Inches(0.3)
        cl = LM + Inches(0.8); ct = Inches(2.0); ch = cb - ct
        cr = Inches(11.5); cww = cr - cl; np_ = len(periods)
        bw = Inches(1.2); sbs = cww / np_
        # Sub-title + legend (adaptive spacing to prevent overflow)
        legend_y = Inches(1.55)
        add_text(s, cl, Inches(1.2), Inches(5.0), Inches(0.3),
                 title, font_size=Pt(13), font_color=DARK_GRAY, bold=True)
        n_series = len(series)
        legend_total_w = cr - LM - Inches(5.0)
        legend_spacing = legend_total_w / max(n_series, 1)
        for ci, (sname, scolor) in enumerate(series):
            lx = LM + Inches(5.2) + legend_spacing * ci
            add_rect(s, lx, legend_y + Inches(0.05), Inches(0.2), Inches(0.2), scolor)
            add_text(s, lx + Inches(0.3), legend_y, min(Inches(1.6), legend_spacing - Inches(0.4)), Inches(0.3),
                     sname, font_size=SMALL_SIZE, font_color=DARK_GRAY)
        # Y-axis
        for tk in [0, 25, 50, 75, 100]:
            ty = cb - int(ch * tk / 100)
            add_text(s, LM, ty - Inches(0.15), Inches(0.7), Inches(0.3),
                     f'{tk}%', font_size=FOOTNOTE_SIZE, font_color=MED_GRAY,
                     alignment=PP_ALIGN.RIGHT)
            if tk > 0:
                add_hline(s, cl, ty, cww, LINE_GRAY, Pt(0.25))
        add_hline(s, cl, cb, cww, BLACK, Pt(0.5))
        for pi, period in enumerate(periods):
            bx = cl + sbs * pi + (sbs - bw) / 2; cum = 0
            for ci in range(len(series)):
                val = data[pi][ci]; sh_ = int(ch * val / 100)
                sy = cb - int(ch * (cum + val) / 100)
                if val > 0:
                    add_rect(s, bx, sy, bw, sh_, series[ci][1])
                    if sh_ >= Inches(0.4):
                        lc = WHITE if ci < 2 else DARK_GRAY
                        add_text(s, bx, sy, bw, sh_, f'{val}%',
                                 font_size=Pt(11), font_color=lc, bold=True,
                                 alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                cum += val
            add_text(s, cl + sbs * pi, cb + Inches(0.05), sbs, Inches(0.3),
                     period, font_size=BODY_SIZE, font_color=DARK_GRAY,
                     alignment=PP_ALIGN.CENTER)
        if summary:
            add_rect(s, LM, sum_y, CW, sum_h, BG_GRAY)
            add_text(s, LM + Inches(0.3), sum_y, Inches(1.5), sum_h,
                     summary[0], font_size=BODY_SIZE, font_color=NAVY, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, LM + Inches(2.0), sum_y, CW - Inches(2.3), sum_h,
                     summary[1], font_size=BODY_SIZE, font_color=DARK_GRAY,
                     anchor=MSO_ANCHOR.MIDDLE)
        self._footer(s, source)
        return s

    def horizontal_bar(self, title, items, summary=None, source=''):
        """#39 Horizontal Bar Chart — labeled bars with percentage.
        items: list of (name, pct_int_0_to_100, bar_color).
        summary: (label, text) or None.
        """
        s = self._ns()
        add_action_title(s, title)
        lbl_w = Inches(2.0); bar_x = LM + Inches(2.2); bar_mw = Inches(7.5)
        val_x = bar_x + bar_mw + Inches(0.2)
        rh = Inches(0.65); sy = Inches(1.6)
        for i, (name, val, bcolor) in enumerate(items):
            ry = sy + rh * i
            bw = int(bar_mw * val / 100)
            tc = NAVY if i == 0 else DARK_GRAY
            add_text(s, LM, ry, lbl_w, rh, name, font_size=BODY_SIZE,
                     font_color=tc, bold=(i == 0), anchor=MSO_ANCHOR.MIDDLE)
            add_rect(s, bar_x, ry + Inches(0.12), bar_mw, Inches(0.4), BG_GRAY)
            add_rect(s, bar_x, ry + Inches(0.12), bw, Inches(0.4), bcolor)
            add_text(s, val_x, ry, Inches(1.0), rh, f'{val}%',
                     font_size=BODY_SIZE, font_color=tc, bold=(i == 0),
                     anchor=MSO_ANCHOR.MIDDLE)
            if i < len(items) - 1:
                add_hline(s, LM, ry + rh, bar_mw + Inches(2.5), LINE_GRAY, Pt(0.25))
        if summary:
            sum_h = Inches(0.7)
            sum_y = SOURCE_Y - sum_h - Inches(0.05)
            add_rect(s, LM, sum_y, CW, sum_h, BG_GRAY)
            add_text(s, LM + Inches(0.3), sum_y, Inches(1.5), sum_h,
                     summary[0], font_size=BODY_SIZE, font_color=NAVY, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, LM + Inches(2.0), sum_y, CW - Inches(2.3), sum_h,
                     summary[1], font_size=BODY_SIZE, font_color=DARK_GRAY,
                     anchor=MSO_ANCHOR.MIDDLE)
        self._footer(s, source)
        return s

    # ═══════════════════════════════════════════
    # IMAGE LAYOUTS (#42, #43, #44, #45, #46, #47)
    # ═══════════════════════════════════════════

    def three_images(self, title, items, source=''):
        """#42 Three Images — three image+caption columns.
        items: list of (caption_title, description, image_label).
        """
        s = self._ns()
        add_action_title(s, title)
        col_w = Inches(3.7); gap = Inches(0.35); img_h = Inches(2.5); ty = Inches(1.3)
        for i, (cap_title, desc, img_label) in enumerate(items):
            cx = LM + i * (col_w + gap)
            add_image_placeholder(s, cx, ty, col_w, img_h, img_label)
            add_text(s, cx, ty + img_h + Inches(0.15), col_w, Inches(0.35),
                     cap_title, font_size=EMPHASIS_SIZE, font_color=NAVY, bold=True)
            add_text(s, cx, ty + img_h + Inches(0.55), col_w, Inches(1.0),
                     desc, font_size=BODY_SIZE, font_color=DARK_GRAY)
        self._footer(s, source)
        return s

    def image_four_points(self, title, image_label, points, source=''):
        """#43 Image + 4 Points — center image with 4 corner cards.
        points: list of 4 (point_title, description, accent_color).
        """
        s = self._ns()
        add_action_title(s, title)
        img_w = Inches(5.0); img_h = Inches(2.4)
        img_x = LM + (CW - img_w) / 2
        add_image_placeholder(s, img_x, Inches(2.8), img_w, img_h, image_label)
        card_w = Inches(5.2); card_h = Inches(0.7)
        positions = [
            (LM + Inches(0.5), Inches(1.3)),
            (LM + CW - card_w - Inches(0.5), Inches(1.3)),
            (LM + Inches(0.5), Inches(5.5)),
            (LM + CW - card_w - Inches(0.5), Inches(5.5)),
        ]
        accents = [ACCENT_BLUE, ACCENT_GREEN, ACCENT_ORANGE, ACCENT_RED]
        for i, (pt_title, desc, *rest) in enumerate(points):
            color = rest[0] if rest else accents[i % 4]
            px, py = positions[i]
            add_oval(s, px, py + Inches(0.08), str(i + 1), bg=color)
            add_text(s, px + Inches(0.55), py, Inches(1.5), Inches(0.35),
                     pt_title, font_size=EMPHASIS_SIZE, font_color=color, bold=True)
            add_text(s, px + Inches(0.55), py + Inches(0.35),
                     card_w - Inches(0.55), Inches(0.35),
                     desc, font_size=BODY_SIZE, font_color=DARK_GRAY)
        self._footer(s, source)
        return s

    def full_width_image(self, title, image_label, overlay_text='', attribution='',
                         source=''):
        """#44 Full-Width Image — edge-to-edge image with text overlay."""
        s = self._ns()
        add_image_placeholder(s, 0, 0, SW, Inches(6.5), image_label)
        # Semi-transparent overlay bar
        overlay = add_rect(s, 0, Inches(4.0), SW, Inches(2.0), NAVY)
        fe = overlay._element.find(qn('p:spPr'))
        if fe is not None:
            sf = fe.find(qn('a:solidFill'))
            if sf is not None:
                sr = sf.find(qn('a:srgbClr'))
                if sr is not None:
                    alpha = sr.makeelement(qn('a:alpha'), {'val': '70000'})
                    sr.append(alpha)
        if overlay_text:
            add_text(s, LM, Inches(4.1), CW, Inches(0.8), overlay_text,
                     font_size=SECTION_TITLE_SIZE, font_color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if attribution:
            add_text(s, LM, Inches(4.9), CW, Inches(0.4), attribution,
                     font_size=BODY_SIZE, font_color=LINE_GRAY,
                     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if source:
            add_source(s, source)
        add_page_number(s, self._page, self.total)
        return s

    def case_study_image(self, title, sections, image_label, kpis=None, source=''):
        """#45 Case Study with Image — left text sections + right image + KPIs.
        sections: list of (label, text, accent_color).
        kpis: list of (value, label) or None.
        """
        s = self._ns()
        add_action_title(s, title)
        lw = Inches(6.5)
        ty = Inches(1.4)
        for i, (label, text, color) in enumerate(sections):
            add_rect(s, LM, ty, Inches(0.06), Inches(1.0), color)
            add_text(s, LM + Inches(0.2), ty, Inches(1.2), Inches(0.3),
                     label, font_size=EMPHASIS_SIZE, font_color=color, bold=True)
            add_text(s, LM + Inches(0.2), ty + Inches(0.35), lw - Inches(0.3), Inches(0.65),
                     text, font_size=BODY_SIZE, font_color=DARK_GRAY)
            ty += Inches(1.3)
        rx = LM + lw + Inches(0.3); rw = CW - lw - Inches(0.3)
        add_image_placeholder(s, rx, Inches(1.4), rw, Inches(2.5), image_label)
        if kpis:
            kpw = rw / len(kpis)
            for i, (val, label) in enumerate(kpis):
                kx = rx + i * kpw
                add_rect(s, kx, Inches(4.2), kpw - Inches(0.1), Inches(1.2), BG_GRAY)
                add_text(s, kx, Inches(4.25), kpw - Inches(0.1), Inches(0.6), val,
                         font_size=SECTION_TITLE_SIZE, font_color=NAVY, bold=True,
                         alignment=PP_ALIGN.CENTER, font_name=FONT_HEADER)
                add_text(s, kx, Inches(4.85), kpw - Inches(0.1), Inches(0.4), label,
                         font_size=SMALL_SIZE, font_color=MED_GRAY,
                         alignment=PP_ALIGN.CENTER)
        self._footer(s, source)
        return s

    def quote_bg_image(self, image_label, quote_text, attribution='', source=''):
        """#46 Quote with Background Image — image top + quote bottom."""
        s = self._ns()
        add_image_placeholder(s, 0, 0, SW, Inches(3.5), image_label)
        add_rect(s, 0, Inches(3.5), SW, Inches(4.0), WHITE)
        ln_x = LM + Inches(1.0); ln_w = CW - Inches(2.0)
        add_hline(s, ln_x, Inches(3.9), ln_w, NAVY, Pt(1.0))
        add_text(s, LM + Inches(1.5), Inches(4.1), CW - Inches(3.0), Inches(1.4),
                 quote_text, font_size=Pt(24), font_color=NAVY, bold=True,
                 alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if attribution:
            add_text(s, LM + Inches(1.5), Inches(5.5), CW - Inches(3.0), Inches(0.4),
                     attribution, font_size=BODY_SIZE, font_color=MED_GRAY,
                     alignment=PP_ALIGN.CENTER)
        add_hline(s, ln_x, Inches(6.0), ln_w, NAVY, Pt(1.0))
        if source:
            add_source(s, source)
        add_page_number(s, self._page, self.total)
        return s

    def goals_illustration(self, title, goals, image_label, source=''):
        """#47 Goals with Illustration — left numbered goals + right image.
        goals: list of (goal_title, description, accent_color).
        """
        s = self._ns()
        add_action_title(s, title)
        lw = Inches(6.5)
        ty = Inches(1.4)
        for i, (gtitle, desc, color) in enumerate(goals):
            add_rect(s, LM, ty, Inches(0.06), Inches(0.8), color)
            add_oval(s, LM + Inches(0.25), ty + Inches(0.15), str(i + 1), bg=color)
            add_text(s, LM + Inches(0.8), ty, lw - Inches(1.0), Inches(0.35),
                     gtitle, font_size=EMPHASIS_SIZE, font_color=color, bold=True)
            add_text(s, LM + Inches(0.8), ty + Inches(0.35), lw - Inches(1.0), Inches(0.45),
                     desc, font_size=BODY_SIZE, font_color=DARK_GRAY)
            ty += Inches(1.2)
        rx = LM + lw + Inches(0.3); rw = CW - lw - Inches(0.3)
        add_image_placeholder(s, rx, Inches(1.4), rw, Inches(4.8), image_label)
        self._footer(s, source)
        return s

    # ═══════════════════════════════════════════
    # ADVANCED CHARTS (#48, #49, #50, #51, #52, #53, #54, #55, #56, #64, #70)
    # ═══════════════════════════════════════════

    def donut(self, title, segments, center_label='', center_sub='',
              legend_x=None, summary=None, source=''):
        """#48 Donut Chart — BLOCK_ARC ring segments.
        segments: list of (pct_float, color, label).
        """
        s = self._ns()
        add_action_title(s, title)
        cx = LM + Inches(3.05); cy = Inches(3.58); outer_r = Inches(1.58)
        thin_ring_adj = 5200
        add_block_arc(s, cx, cy, outer_r, 0, 359.8, LINE_GRAY, inner_ratio=thin_ring_adj)
        dsa = 0
        for pct, color, _ in segments:
            sweep = pct * 360
            add_block_arc(s, cx, cy, outer_r, dsa, sweep, color, inner_ratio=thin_ring_adj)
            dsa += sweep
        if center_label:
            add_text(s, cx - Inches(0.75), cy - Inches(0.28), Inches(1.5), Inches(0.45),
                     center_label, font_size=Pt(24), font_color=NAVY, bold=True,
                     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                     font_name=FONT_HEADER)
        if center_sub:
            add_text(s, cx - Inches(0.9), cy + Inches(0.12), Inches(1.8), Inches(0.26),
                     center_sub, font_size=SMALL_SIZE, font_color=MED_GRAY,
                     alignment=PP_ALIGN.CENTER)
        lgx = legend_x or (LM + Inches(6.4))
        for i, (pct, color, label) in enumerate(segments):
            ly = Inches(1.85) + i * Inches(0.78)
            add_rect(s, lgx, ly + Inches(0.04), Inches(0.28), Inches(0.28), color)
            add_text(s, lgx + Inches(0.42), ly, Inches(3.2), Inches(0.34),
                     f'{label}  {int(pct * 100)}%',
                     font_size=Pt(15), font_color=DARK_GRAY, bold=True)
        if summary:
            add_rect(s, lgx, Inches(5.32), Inches(4.85), Inches(0.76), BG_GRAY)
            add_text(s, lgx + Inches(0.18), Inches(5.32), Inches(4.45), Inches(0.76),
                     summary, font_size=BODY_SIZE, font_color=NAVY, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE)
        self._footer(s, source)
        return s

    def waterfall(self, title, items, max_val=None, legend_items=None,
                  summary=None, source=''):
        """#49 Waterfall Chart — bridge from start to end.
        items: list of (label, value, type) — type: 'base'|'up'|'down'.
        """
        s = self._ns()
        add_action_title(s, title)
        # Chart sub-title + legend
        add_text(s, LM, Inches(1.2), Inches(6.0), Inches(0.3),
                 title, font_size=Pt(13), font_color=DARK_GRAY, bold=True)
        if legend_items:
            for li, (ln, lc) in enumerate(legend_items):
                lx = LM + Inches(8.0) + li * Inches(1.3)
                add_rect(s, lx, Inches(1.25), Inches(0.2), Inches(0.2), lc)
                add_text(s, lx + Inches(0.25), Inches(1.2), Inches(0.8), Inches(0.3),
                         ln, font_size=Pt(11), font_color=DARK_GRAY)
        # Layout: summary sinks to bottom, chart fills space above
        sum_h = Inches(0.6)
        sum_y = SOURCE_Y - sum_h - Inches(0.05)  # summary flush to source
        cb = sum_y - Inches(0.4) if summary else SOURCE_Y - Inches(0.3)
        cl = LM + Inches(0.3); ct = Inches(3.1); ch = cb - ct
        if max_val is None:
            max_val = max(abs(v) for _, v, _ in items) * 1.3
        bw = Inches(1.2); gp = Inches(0.4)
        running = 0; prev_top = 0
        for i, (label, val, typ) in enumerate(items):
            bx = cl + i * (bw + gp)
            if typ == 'base':
                bh = int(ch * val / max_val); bt = cb - bh; color = NAVY; running = val
            elif typ == 'up':
                bh = int(ch * val / max_val); bt = cb - int(ch * running / max_val) - bh
                color = ACCENT_GREEN; running += val
            else:
                bh = int(ch * abs(val) / max_val); bt = cb - int(ch * running / max_val)
                color = ACCENT_RED; running += val
            add_rect(s, bx, bt, bw, bh, color)
            if i > 0:
                add_hline(s, bx - gp, prev_top, gp + Inches(0.05), LINE_GRAY, Pt(0.75))
            prev_top = bt if typ != 'down' else bt + bh
            vs = f'+{val}' if val > 0 and typ != 'base' else str(val)
            add_text(s, bx, bt - Inches(0.35), bw, Inches(0.3), vs,
                     font_size=BODY_SIZE, font_color=DARK_GRAY, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name=FONT_HEADER)
            add_text(s, bx, cb + Inches(0.05), bw, Inches(0.3), label,
                     font_size=Pt(11), font_color=MED_GRAY, alignment=PP_ALIGN.CENTER)
        add_hline(s, cl, cb, Inches(11.5), LINE_GRAY, Pt(0.5))
        if summary:
            add_rect(s, LM, sum_y, CW, sum_h, BG_GRAY)
            add_text(s, LM + Inches(0.3), sum_y, CW - Inches(0.6), sum_h,
                     summary, font_size=BODY_SIZE, font_color=NAVY, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE)
        self._footer(s, source)
        return s

    def line_chart(self, title, x_labels, y_labels, values, legend_label='',
                   summary=None, source=''):
        """#50 Line Chart — single line with dot approximation.
        x_labels: list[str], y_labels: list[str], values: list[float] 0.0-1.0 normalized.
        """
        s = self._ns()
        add_action_title(s, title)
        # Layout: summary sinks to bottom, chart fills space above
        sum_h = Inches(0.6)
        sum_y = SOURCE_Y - sum_h - Inches(0.05)
        cb_val = sum_y - Inches(0.4) if summary else SOURCE_Y - Inches(0.3)
        cl = LM + Inches(0.8); cr = LM + CW - Inches(1.5); cw_ = cr - cl
        ct = Inches(2.5); cb = cb_val; ch = cb - ct
        # Sub-title + legend
        add_text(s, cl, Inches(1.2), Inches(5.0), Inches(0.3),
                 title, font_size=Pt(13), font_color=DARK_GRAY, bold=True)
        if legend_label:
            add_rect(s, LM + Inches(9.0), Inches(1.25), Inches(0.3), Inches(0.15), NAVY)
            add_text(s, LM + Inches(9.4), Inches(1.2), Inches(1.5), Inches(0.3),
                     legend_label, font_size=Pt(11), font_color=NAVY, bold=True)
        # Y-axis
        for i, yl in enumerate(y_labels):
            yy = cb - int(ch * i / (len(y_labels) - 1))
            add_text(s, LM, yy - Inches(0.12), Inches(0.7), Inches(0.24), yl,
                     font_size=FOOTNOTE_SIZE, font_color=MED_GRAY,
                     alignment=PP_ALIGN.RIGHT)
            if i > 0:
                add_hline(s, cl, yy, cw_, RGBColor(0xE8, 0xE8, 0xE8), Pt(0.25))
        # X-axis
        npt = len(x_labels)
        for i, xl in enumerate(x_labels):
            xx = cl + int(cw_ * i / (npt - 1))
            add_text(s, xx - Inches(0.3), cb + Inches(0.05), Inches(0.6), Inches(0.25),
                     xl, font_size=FOOTNOTE_SIZE, font_color=MED_GRAY,
                     alignment=PP_ALIGN.CENTER)
        # Line segments (thick rect approximation)
        for j in range(len(values) - 1):
            x1 = cl + int(cw_ * j / (npt - 1)); y1 = cb - int(ch * values[j])
            x2 = cl + int(cw_ * (j + 1) / (npt - 1)); y2 = cb - int(ch * values[j + 1])
            sw_ = x2 - x1; sy = min(y1, y2); sh_ = max(abs(y2 - y1), int(Pt(3)))
            add_rect(s, x1, sy, sw_, sh_, NAVY)
        add_hline(s, cl, cb, cw_, BLACK, Pt(0.5))
        if summary:
            add_rect(s, LM, sum_y, CW, sum_h, BG_GRAY)
            add_text(s, LM + Inches(0.3), sum_y, CW - Inches(0.6), sum_h,
                     summary, font_size=BODY_SIZE, font_color=NAVY, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE)
        self._footer(s, source)
        return s

    def pareto(self, title, items, max_val=None, summary=None, source=''):
        """#51 Pareto Chart — descending bars with value/pct labels.
        items: list of (label, value).
        """
        s = self._ns()
        add_action_title(s, title)
        total = sum(v for _, v in items)
        if max_val is None:
            max_val = max(v for _, v in items) * 1.15
        # Layout: summary sinks to bottom
        sum_h = Inches(0.6)
        sum_y = SOURCE_Y - sum_h - Inches(0.05)
        cb = sum_y - Inches(0.4) if summary else SOURCE_Y - Inches(0.3)
        cl = LM + Inches(1.0); ct = Inches(2.7); ch = cb - ct
        cw_ = Inches(9.0)
        n = len(items); gap = Inches(0.15)
        bw = (cw_ - gap * (n - 1)) / n
        for i, (label, val) in enumerate(items):
            bx = cl + i * (bw + gap)
            bh = int(ch * val / max_val); bt = cb - bh
            add_rect(s, bx, bt, bw, bh, NAVY)
            add_text(s, bx, bt - Inches(0.25), bw, Inches(0.22),
                     f'{val}', font_size=Pt(11), font_color=DARK_GRAY, bold=True,
                     alignment=PP_ALIGN.CENTER)
            pct_ = val / total if total else 0
            add_text(s, bx, bt - Inches(0.48), bw, Inches(0.22),
                     f'{int(pct_ * 100)}%', font_size=FOOTNOTE_SIZE, font_color=MED_GRAY,
                     alignment=PP_ALIGN.CENTER)
            add_text(s, bx, cb + Inches(0.05), bw, Inches(0.35), label,
                     font_size=Pt(10), font_color=MED_GRAY, alignment=PP_ALIGN.CENTER)
        add_hline(s, cl, cb, cw_, BLACK, Pt(0.5))
        if summary:
            add_rect(s, LM, sum_y, CW, sum_h, BG_GRAY)
            add_text(s, LM + Inches(0.3), sum_y, CW - Inches(0.6), sum_h,
                     summary, font_size=BODY_SIZE, font_color=NAVY, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE)
        self._footer(s, source)
        return s

    def kpi_tracker(self, title, kpis, summary=None, source=''):
        """#52 KPI Tracker — progress bars with status dots.
        kpis: list of (name, pct_float, detail, status_key).
        status_key: 'on'|'risk'|'off'.
        """
        s = self._ns()
        add_action_title(s, title)
        sc_map = {
            'on': (ACCENT_GREEN, '达标'),
            'risk': (ACCENT_ORANGE, '关注'),
            'off': (ACCENT_RED, '滞后'),
        }
        # Headers
        hy = Inches(1.3)
        add_text(s, LM, hy, Inches(3.5), Inches(0.35), 'KPI指标',
                 font_size=SMALL_SIZE, font_color=MED_GRAY, bold=True)
        add_text(s, LM + Inches(3.5), hy, Inches(6.0), Inches(0.35), '进度',
                 font_size=SMALL_SIZE, font_color=MED_GRAY, bold=True)
        add_text(s, LM + Inches(9.5), hy, Inches(1.2), Inches(0.35), '达成率',
                 font_size=SMALL_SIZE, font_color=MED_GRAY, bold=True,
                 alignment=PP_ALIGN.CENTER)
        add_text(s, LM + Inches(10.7), hy, Inches(1.0), Inches(0.35), '状态',
                 font_size=SMALL_SIZE, font_color=MED_GRAY, bold=True,
                 alignment=PP_ALIGN.CENTER)
        add_hline(s, LM, hy + Inches(0.35), CW, BLACK, Pt(0.75))
        bar_x = LM + Inches(3.5); bar_mw = Inches(5.8); bar_h = Inches(0.25)
        rh = Inches(0.8)
        for i, (name, pct_, detail, status) in enumerate(kpis):
            ry = Inches(1.9) + i * rh
            add_text(s, LM, ry, Inches(3.3), rh, name,
                     font_size=BODY_SIZE, font_color=DARK_GRAY, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE)
            add_rect(s, bar_x, ry + (rh - bar_h) / 2, bar_mw, bar_h, BG_GRAY)
            color, label = sc_map.get(status, (MED_GRAY, '?'))
            fw = int(bar_mw * min(pct_, 1.0))
            add_rect(s, bar_x, ry + (rh - bar_h) / 2, fw, bar_h, color)
            add_text(s, LM + Inches(9.5), ry, Inches(1.2), rh,
                     f'{int(min(pct_, 1.0) * 100)}%',
                     font_size=EMPHASIS_SIZE, font_color=DARK_GRAY, bold=True,
                     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            dsz = Inches(0.15)
            add_rect(s, LM + Inches(10.8), ry + (rh - dsz) / 2, dsz, dsz, color)
            add_text(s, LM + Inches(11.0), ry, Inches(0.7), rh, label,
                     font_size=Pt(11), font_color=color, anchor=MSO_ANCHOR.MIDDLE)
            if i < len(kpis) - 1:
                add_hline(s, LM, ry + rh, CW, LINE_GRAY, Pt(0.25))
        if summary:
            add_rect(s, LM, Inches(6.0), CW, Inches(0.7), BG_GRAY)
            add_text(s, LM + Inches(0.3), Inches(6.0), Inches(1.5), Inches(0.7),
                     '总结', font_size=BODY_SIZE, font_color=NAVY, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, LM + Inches(2.0), Inches(6.0), CW - Inches(2.3), Inches(0.7),
                     summary, font_size=BODY_SIZE, font_color=DARK_GRAY,
                     anchor=MSO_ANCHOR.MIDDLE)
        self._footer(s, source)
        return s

    def bubble(self, title, bubbles, x_label='', y_label='',
               legend_items=None, summary=None, source=''):
        """#53 Bubble / Scatter — positioned circles on XY plane.
        bubbles: list of (x_pct, y_pct, size_inches, label, color).
        """
        s = self._ns()
        add_action_title(s, title)
        cl = LM + Inches(1.2); cb = Inches(5.5); ct = Inches(1.6)
        cw_ = Inches(9.0); ch = cb - ct
        # Legend
        if legend_items:
            for li, (ln, lc) in enumerate(legend_items):
                lx = LM + Inches(7.2) + li * Inches(0.85)
                add_rect(s, lx, Inches(1.25), Inches(0.15), Inches(0.15), lc)
                add_text(s, lx + Inches(0.2), Inches(1.2), Inches(0.55), Inches(0.3),
                         ln, font_size=Pt(10), font_color=DARK_GRAY)
        add_hline(s, cl, cb, cw_, BLACK, Pt(0.5))
        add_rect(s, cl, ct, Pt(0.5), ch, BLACK)
        if x_label:
            add_text(s, cl + cw_ // 2 - Inches(1.0), cb + Inches(0.15), Inches(2.0), Inches(0.3),
                     x_label, font_size=Pt(11), font_color=MED_GRAY,
                     alignment=PP_ALIGN.CENTER)
        if y_label:
            add_text(s, LM, ct + ch // 2 - Inches(0.5), Inches(1.0), Inches(1.0),
                     y_label, font_size=Pt(11), font_color=MED_GRAY,
                     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        for xp, yp, sz, lb, color in bubbles:
            bx = cl + int(cw_ * xp) - Inches(sz / 2)
            by = cb - int(ch * yp) - Inches(sz / 2)
            oval = s.shapes.add_shape(MSO_SHAPE.OVAL, bx, by, Inches(sz), Inches(sz))
            oval.fill.solid(); oval.fill.fore_color.rgb = color
            oval.line.fill.background(); _clean_shape(oval)
            tc = WHITE if color != LINE_GRAY else DARK_GRAY
            add_text(s, bx, by + Inches(sz * 0.2), Inches(sz), Inches(sz * 0.6),
                     lb, font_size=FOOTNOTE_SIZE, font_color=tc, bold=True,
                     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if summary:
            add_rect(s, LM, Inches(6.0), CW, Inches(0.6), BG_GRAY)
            add_text(s, LM + Inches(0.3), Inches(6.0), CW - Inches(0.6), Inches(0.6),
                     summary, font_size=BODY_SIZE, font_color=NAVY, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE)
        self._footer(s, source)
        return s

    def risk_matrix(self, title, grid_colors, grid_lights, risks,
                    y_labels=None, x_labels=None, notes=None, source=''):
        """#54 Risk Matrix — 3×3 heatmap grid with risk labels.
        grid_colors: 3×3 list[list[RGBColor]] dot colors.
        grid_lights: 3×3 list[list[RGBColor]] cell backgrounds.
        risks: list of (row, col, name).
        notes: list[str] or None for bottom panel.
        """
        s = self._ns()
        add_action_title(s, title)
        gl = LM + Inches(1.8); gt = Inches(1.5)
        gcw = Inches(3.0); gch = Inches(1.1)
        if y_labels is None:
            y_labels = ['高概率', '中概率', '低概率']
        if x_labels is None:
            x_labels = ['低影响', '中影响', '高影响']
        for r in range(3):
            add_text(s, LM, gt + r * gch, Inches(1.6), gch, y_labels[r],
                     font_size=Pt(13), font_color=DARK_GRAY, bold=True,
                     alignment=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        for c in range(3):
            add_text(s, gl + c * gcw, gt - Inches(0.35), gcw, Inches(0.3),
                     x_labels[c], font_size=Pt(13), font_color=DARK_GRAY, bold=True,
                     alignment=PP_ALIGN.CENTER)
        for r in range(3):
            for c in range(3):
                add_rect(s, gl + c * gcw, gt + r * gch,
                         gcw - Inches(0.05), gch - Inches(0.05), grid_lights[r][c])
                add_rect(s, gl + c * gcw + Inches(0.1), gt + r * gch + Inches(0.1),
                         Inches(0.2), Inches(0.2), grid_colors[r][c])
        for r, c, name in risks:
            add_text(s, gl + c * gcw + Inches(0.4), gt + r * gch + Inches(0.25),
                     gcw - Inches(0.6), Inches(0.6), name,
                     font_size=Pt(13), font_color=DARK_GRAY, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE)
        if notes:
            add_rect(s, LM, Inches(5.1), CW, Inches(1.6), BG_GRAY)
            add_text(s, LM + Inches(0.3), Inches(5.15), Inches(1.5), Inches(0.3),
                     '应对措施', font_size=BODY_SIZE, font_color=NAVY, bold=True)
            add_text(s, LM + Inches(0.3), Inches(5.5), CW - Inches(0.6), Inches(1.1),
                     notes, font_size=BODY_SIZE, font_color=DARK_GRAY, line_spacing=Pt(4))
        self._footer(s, source)
        return s

    def gauge(self, title, score, benchmarks=None, source=''):
        """#55 Gauge — ⚠️ RETIRED (已废弃).
        This layout has been retired per design review.
        score: int 0-100.
        benchmarks: list of (label, value_str, color) shown below gauge.
        """
        s = self._ns()
        add_action_title(s, title)
        gcx = LM + CW // 2; gcy = Inches(4.0); gr = Inches(2.2)
        segs = [(0.40, ACCENT_RED), (0.30, ACCENT_ORANGE), (0.30, ACCENT_GREEN)]
        ppt_cum = 270
        for gpct, gcolor in segs:
            gsweep = gpct * 180
            ppt_start = ppt_cum % 360
            ppt_end = (ppt_cum + gsweep) % 360
            gir = gr - Pt(10)
            inner_ratio = int((gir / gr) * 50000)
            d = gr * 2
            sh = s.shapes.add_shape(MSO_SHAPE.BLOCK_ARC, gcx - gr, gcy - gr, d, d)
            sh.fill.solid(); sh.fill.fore_color.rgb = gcolor
            sh.line.fill.background(); _clean_shape(sh)
            sp = sh._element.find(qn('p:spPr'))
            prstGeom = sp.find(qn('a:prstGeom'))
            if prstGeom is not None:
                avLst = prstGeom.find(qn('a:avLst'))
                if avLst is None:
                    avLst = prstGeom.makeelement(qn('a:avLst'), {})
                    prstGeom.append(avLst)
                for gd in avLst.findall(qn('a:gd')):
                    avLst.remove(gd)
                gd1 = avLst.makeelement(qn('a:gd'),
                                        {'name': 'adj1', 'fmla': f'val {int(ppt_start * 60000)}'})
                gd2 = avLst.makeelement(qn('a:gd'),
                                        {'name': 'adj2', 'fmla': f'val {int(ppt_end * 60000)}'})
                gd3 = avLst.makeelement(qn('a:gd'),
                                        {'name': 'adj3', 'fmla': f'val {inner_ratio}'})
                avLst.append(gd1); avLst.append(gd2); avLst.append(gd3)
            ppt_cum += gsweep
        # Center score circle
        cir_sz = Inches(2.0)
        o = s.shapes.add_shape(MSO_SHAPE.OVAL, gcx - cir_sz // 2, gcy - cir_sz // 2,
                               cir_sz, cir_sz)
        o.fill.solid(); o.fill.fore_color.rgb = WHITE
        o.line.fill.background(); _clean_shape(o)
        add_text(s, gcx - Inches(0.8), gcy - Inches(0.6), Inches(1.6), Inches(0.6),
                 str(score), font_size=COVER_TITLE_SIZE, font_color=NAVY, bold=True,
                 alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font_name=FONT_HEADER)
        add_text(s, gcx - Inches(0.5), gcy + Inches(0.1), Inches(1.0), Inches(0.3),
                 '/ 100', font_size=BODY_SIZE, font_color=MED_GRAY,
                 alignment=PP_ALIGN.CENTER)
        if benchmarks:
            bsx = LM + Inches(0.5); bry = Inches(5.3)
            for i, (label, val, color) in enumerate(benchmarks):
                bx = bsx + i * Inches(2.5)
                add_rect(s, bx, bry, Inches(0.08), Inches(0.6), color)
                add_text(s, bx + Inches(0.2), bry, Inches(2.0), Inches(0.3),
                         label, font_size=SMALL_SIZE, font_color=MED_GRAY)
                add_text(s, bx + Inches(0.2), bry + Inches(0.3), Inches(2.0), Inches(0.3),
                         val, font_size=Pt(22), font_color=color, bold=True,
                         font_name=FONT_HEADER)
        self._footer(s, source)
        return s

    def harvey_ball_table(self, title, criteria, options, scores,
                          legend_text=None, summary=None, source=''):
        """#56 Harvey Ball Table — matrix with Harvey Ball indicators.
        criteria: list[str] row labels. options: list[str] column headers.
        scores: list[list[int]] — scores[row][col], 0-4.
        """
        s = self._ns()
        add_action_title(s, title)
        c1w = Inches(2.8); colw = Inches(2.5); rh = Inches(0.6)
        tl = LM; ty = Inches(1.3)
        add_text(s, tl, ty, c1w, rh, '评估维度',
                 font_size=Pt(13), font_color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        for j, opt in enumerate(options):
            add_text(s, tl + c1w + j * colw, ty, colw, rh, opt,
                     font_size=Pt(13), font_color=NAVY, bold=True,
                     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        total_w = c1w + len(options) * colw
        add_hline(s, tl, ty + rh, total_w, BLACK, Pt(0.75))
        for i, cr in enumerate(criteria):
            ry = ty + rh + Inches(0.05) + i * rh
            add_text(s, tl, ry, c1w, rh, cr,
                     font_size=BODY_SIZE, font_color=DARK_GRAY, anchor=MSO_ANCHOR.MIDDLE)
            for j in range(len(options)):
                bx = tl + c1w + j * colw + (colw - Inches(0.35)) // 2
                by = ry + (rh - Inches(0.35)) // 2
                draw_harvey_ball(s, bx, by, scores[i][j])
            if i < len(criteria) - 1:
                add_hline(s, tl, ry + rh, total_w, LINE_GRAY, Pt(0.25))
        # Legend
        lgy = ty + rh + len(criteria) * rh + Inches(0.3)
        add_hline(s, tl, lgy - Inches(0.1), total_w, BLACK, Pt(0.5))
        if legend_text:
            lx = tl
            for item in legend_text:
                add_text(s, lx, lgy, Inches(2.2), Inches(0.3), item,
                         font_size=Pt(11), font_color=MED_GRAY)
                lx += Inches(2.4)
        if summary:
            add_rect(s, LM, Inches(5.8), CW, Inches(0.8), BG_GRAY)
            add_text(s, LM + Inches(0.3), Inches(5.8), CW - Inches(0.6), Inches(0.8),
                     summary, font_size=BODY_SIZE, font_color=NAVY, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE)
        self._footer(s, source)
        return s

    def pie(self, title, segments, legend_x=None, summary=None, source=''):
        """#64 Pie Chart — BLOCK_ARC with inner_ratio=0 (solid).
        segments: list of (pct_float, color, label, sub_label).
        """
        s = self._ns()
        add_action_title(s, title)
        pcx = LM + Inches(3.0); pcy = Inches(3.5); pr = Inches(1.8)
        psa = 0
        for pct, color, _, _ in segments:
            sweep = pct * 360
            add_block_arc(s, pcx, pcy, pr, psa, sweep, color,
                          inner_ratio=0, ring_width=pr)
            psa += sweep
        lgx = legend_x or (LM + Inches(7.0))
        for i, (pct, color, label, sub_label) in enumerate(segments):
            ly = Inches(1.8) + i * Inches(0.9)
            add_rect(s, lgx, ly + Inches(0.05), Inches(0.3), Inches(0.3), color)
            add_text(s, lgx + Inches(0.45), ly, Inches(3.5), Inches(0.3),
                     label, font_size=EMPHASIS_SIZE, font_color=DARK_GRAY, bold=True)
            if sub_label:
                add_text(s, lgx + Inches(0.45), ly + Inches(0.3), Inches(3.5), Inches(0.3),
                         sub_label, font_size=Pt(13), font_color=MED_GRAY)
        if summary:
            add_rect(s, LM, Inches(5.6), CW, Inches(0.8), BG_GRAY)
            add_text(s, LM + Inches(0.3), Inches(5.6), CW - Inches(0.6), Inches(0.8),
                     summary, font_size=BODY_SIZE, font_color=NAVY, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE)
        self._footer(s, source)
        return s

    def stacked_area(self, title, years, series_data, max_val=None,
                     summary=None, source=''):
        """#70 Stacked Area Chart — stacked columns for area approximation.
        years: list[str] x-labels.
        series_data: list of (name, values:list[int], color).
        """
        s = self._ns()
        add_action_title(s, title)
        # Layout: summary sinks to bottom, chart fills space above
        sum_h = Inches(0.7)
        sum_y = SOURCE_Y - sum_h - Inches(0.05)
        cb = sum_y - Inches(0.35) if summary else SOURCE_Y - Inches(0.3)
        cl = LM + Inches(1.0); ct = Inches(2.5)
        cw_ = Inches(9.5); ch = cb - ct
        if max_val is None:
            max_val = max(sum(sd[1][yi] for sd in series_data)
                         for yi in range(len(years))) * 1.15
        # Sub-title + legend
        add_text(s, cl, Inches(1.2), Inches(5.0), Inches(0.3),
                 title, font_size=Pt(13), font_color=DARK_GRAY, bold=True)
        for i, (name, _, color) in enumerate(series_data):
            lgx = LM + Inches(8.0) + i * Inches(1.5)
            add_rect(s, lgx, Inches(1.25), Inches(0.25), Inches(0.2), color)
            add_text(s, lgx + Inches(0.35), Inches(1.2), Inches(1.0), Inches(0.3),
                     name, font_size=Pt(11), font_color=DARK_GRAY)
        # Y-axis
        for i in range(5):
            val = max_val * i / 4; yy = cb - int(ch * i / 4)
            add_text(s, LM, yy - Inches(0.1), Inches(0.8), Inches(0.2),
                     f'¥{int(val)}', font_size=FOOTNOTE_SIZE, font_color=MED_GRAY,
                     alignment=PP_ALIGN.RIGHT)
            if i > 0:
                add_hline(s, cl, yy, cw_, RGBColor(0xE8, 0xE8, 0xE8), Pt(0.25))
        npt = len(years); col_w = cw_ // npt
        for yi in range(npt):
            cum = 0
            for si, (name, values, color) in enumerate(series_data):
                val = values[yi]
                bh = int(ch * val / max_val); base_h = int(ch * cum / max_val)
                bx = cl + int(cw_ * yi / npt); by = cb - base_h - bh
                add_rect(s, bx + Inches(0.05), by, col_w - Inches(0.1), bh, color)
                cum += val
            total = sum(sd[1][yi] for sd in series_data)
            th = int(ch * total / max_val)
            add_text(s, cl + int(cw_ * yi / npt), cb - th - Inches(0.25),
                     col_w, Inches(0.2), f'¥{total}',
                     font_size=Pt(10), font_color=DARK_GRAY, bold=True,
                     alignment=PP_ALIGN.CENTER)
            add_text(s, cl + int(cw_ * yi / npt), cb + Inches(0.05),
                     col_w, Inches(0.2), years[yi],
                     font_size=Pt(10), font_color=MED_GRAY, alignment=PP_ALIGN.CENTER)
        add_hline(s, cl, cb, cw_, BLACK, Pt(0.5))
        if summary:
            add_rect(s, LM, sum_y, CW, sum_h, BG_GRAY)
            add_text(s, LM + Inches(0.3), sum_y, Inches(1.5), sum_h,
                     '趋势分析', font_size=BODY_SIZE, font_color=NAVY, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, LM + Inches(2.0), sum_y, CW - Inches(2.3), sum_h,
                     summary, font_size=BODY_SIZE, font_color=DARK_GRAY,
                     anchor=MSO_ANCHOR.MIDDLE)
        self._footer(s, source)
        return s

    # ═══════════════════════════════════════════
    # DASHBOARDS & SPECIAL (#57, #58, #59, #60, #62, #63, #66, #68, #69)
    # ═══════════════════════════════════════════

    def dashboard_kpi_chart(self, title, kpi_cards, chart_data=None,
                            summary=None, source=''):
        """#57 Dashboard KPI + Chart — top KPI cards + bottom mini chart.
        kpi_cards: list of (value, label, detail, accent_color).
        chart_data: dict with 'labels','actual','target','max_val','legend'.
        """
        s = self._ns()
        add_action_title(s, title)
        n = len(kpi_cards)
        dkw = CW / n - Inches(0.15); dkh = Inches(1.1); dky = Inches(1.2)
        for i, (val, label, detail, color) in enumerate(kpi_cards):
            cx = LM + i * (dkw + Inches(0.15))
            add_rect(s, cx, dky, dkw, dkh, WHITE)
            add_rect(s, cx, dky, dkw, Inches(0.06), color)
            add_text(s, cx + Inches(0.2), dky + Inches(0.15), dkw - Inches(0.4), Inches(0.45),
                     val, font_size=Pt(24), font_color=color, bold=True,
                     font_name=FONT_HEADER)
            add_text(s, cx + Inches(0.2), dky + Inches(0.6), Inches(1.5), Inches(0.2),
                     label, font_size=Pt(11), font_color=MED_GRAY)
            add_text(s, cx + Inches(1.8), dky + Inches(0.6), dkw - Inches(2.0), Inches(0.2),
                     detail, font_size=Pt(10), font_color=ACCENT_GREEN,
                     alignment=PP_ALIGN.RIGHT)
        if chart_data:
            dcy = Inches(2.6); dch = Inches(2.5)
            dcl = LM + Inches(0.5); dcb = dcy + dch
            labels = chart_data.get('labels', [])
            actual = chart_data.get('actual', [])
            target = chart_data.get('target', [])
            mv = chart_data.get('max_val', max(actual + target) * 1.15)
            dbw = Inches(0.6); dpg = Inches(0.15)
            dgw = dbw * 2 + dpg; dmgp = Inches(0.5)
            for i in range(len(labels)):
                gx = dcl + i * (dgw + dmgp)
                for j, (vals, color) in enumerate([(actual, NAVY), (target, BG_GRAY)]):
                    bx = gx + j * (dbw + dpg); val = vals[i]
                    bh = int(dch * val / mv); bt = dcb - bh
                    add_rect(s, bx, bt, dbw, bh, color)
                add_text(s, gx, dcb + Inches(0.03), dgw, Inches(0.2),
                         labels[i], font_size=FOOTNOTE_SIZE, font_color=MED_GRAY,
                         alignment=PP_ALIGN.CENTER)
            add_hline(s, dcl, dcb, Inches(10.5), LINE_GRAY, Pt(0.5))
            legend = chart_data.get('legend', [('实际', NAVY), ('目标', BG_GRAY)])
            for li, (ln, lc) in enumerate(legend):
                lx = LM + Inches(9.0) + li * Inches(1.3)
                add_rect(s, lx, dcy, Inches(0.3), Inches(0.15), lc)
                add_text(s, lx + Inches(0.4), dcy - Inches(0.02), Inches(0.8), Inches(0.2),
                         ln, font_size=FOOTNOTE_SIZE, font_color=DARK_GRAY)
        if summary:
            add_rect(s, LM, Inches(5.6), CW, Inches(0.9), BG_GRAY)
            add_text(s, LM + Inches(0.3), Inches(5.6), Inches(1.5), Inches(0.9),
                     '关键发现', font_size=BODY_SIZE, font_color=NAVY, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE)
            summary_items = summary if isinstance(summary, list) else [summary]
            add_text(s, LM + Inches(2.0), Inches(5.6), CW - Inches(2.3), Inches(0.9),
                     summary_items, font_size=BODY_SIZE, font_color=DARK_GRAY,
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=Pt(4))
        self._footer(s, source)
        return s

    def dashboard_table_chart(self, title, table_data, chart_data=None,
                              factoids=None, source=''):
        """#58 Dashboard Table + Chart — left table + right mini chart + bottom facts.
        table_data: dict with 'headers','col_widths','rows'.
        chart_data: dict with 'title','items':(name, value, max_val).
        factoids: list of (value, label, color).
        """
        s = self._ns()
        add_action_title(s, title)
        td = table_data or {}
        headers = td.get('headers', [])
        col_widths = td.get('col_widths', [])
        rows = td.get('rows', [])
        dlw = sum(col_widths) if col_widths else Inches(6.2)
        dty = Inches(1.3)
        dhx = LM
        for ci, (cn, cw) in enumerate(zip(headers, col_widths)):
            add_text(s, dhx, dty, cw, Inches(0.3), cn,
                     font_size=SMALL_SIZE, font_color=NAVY, bold=True)
            dhx += cw
        add_hline(s, LM, dty + Inches(0.3), dlw, BLACK, Pt(0.75))
        for ri, row in enumerate(rows):
            ry = dty + Inches(0.4) + ri * Inches(0.5); rx = LM
            for ci, (val, cw) in enumerate(zip(row, col_widths)):
                fc = DARK_GRAY
                if ci == 2 and '+' in str(val):
                    fc = ACCENT_GREEN
                add_text(s, rx, ry, cw, Inches(0.4), str(val),
                         font_size=BODY_SIZE, font_color=fc, anchor=MSO_ANCHOR.MIDDLE)
                rx += cw
            if ri < len(rows) - 1:
                add_hline(s, LM, ry + Inches(0.45), dlw, LINE_GRAY, Pt(0.25))
        if chart_data:
            dcx = LM + dlw + Inches(0.5); dcrw = CW - dlw - Inches(0.5)
            add_text(s, dcx, Inches(1.3), dcrw, Inches(0.3),
                     chart_data.get('title', ''),
                     font_size=SMALL_SIZE, font_color=NAVY, bold=True)
            add_hline(s, dcx, Inches(1.6), dcrw, BLACK, Pt(0.5))
            items = chart_data.get('items', [])
            bmw = Inches(3.5)
            mx = max((v for _, v in items), default=1)
            for i, (rg, rev) in enumerate(items):
                by = Inches(1.75) + i * Inches(0.45)
                bw = int(bmw * rev / mx)
                add_text(s, dcx, by, Inches(0.8), Inches(0.3), rg,
                         font_size=Pt(11), font_color=MED_GRAY, anchor=MSO_ANCHOR.MIDDLE)
                add_rect(s, dcx + Inches(0.9), by, bw, Inches(0.3),
                         NAVY if i == 0 else ACCENT_BLUE)
        if factoids:
            dfy = Inches(5.2)
            n = len(factoids)
            dfcw = CW / n - Inches(0.15); dfch = Inches(1.0)
            for i, (val, label, color) in enumerate(factoids):
                fx = LM + i * (dfcw + Inches(0.15))
                add_rect(s, fx, dfy, dfcw, dfch, BG_GRAY)
                add_rect(s, fx, dfy, Inches(0.06), dfch, color)
                add_text(s, fx + Inches(0.2), dfy + Inches(0.1), dfcw - Inches(0.3), Inches(0.5),
                         val, font_size=Pt(24), font_color=color, bold=True)
                add_text(s, fx + Inches(0.2), dfy + Inches(0.6), dfcw - Inches(0.3), Inches(0.3),
                         label, font_size=Pt(11), font_color=MED_GRAY)
        self._footer(s, source)
        return s

    def stakeholder_map(self, title, quadrants, x_label='影响力 →',
                        y_label='关注度 ↑', summary=None, source=''):
        """#59 Stakeholder Map — 2×2 quadrant with stakeholder lists.
        quadrants: list of 4 (label_cn, label_en, bg_color, members:list[str]).
        """
        s = self._ns()
        add_action_title(s, title)
        sgl = LM + Inches(2.0); sgt = Inches(1.5)
        sgcw = Inches(4.5); sgch = Inches(1.8)
        for qi, (lcn, len_, bgc, members) in enumerate(quadrants):
            r, c = qi // 2, qi % 2
            qx = sgl + c * sgcw; qy = sgt + r * sgch
            add_rect(s, qx, qy, sgcw - Inches(0.05), sgch - Inches(0.05), bgc)
            add_text(s, qx + Inches(0.15), qy + Inches(0.1), sgcw - Inches(0.3), Inches(0.35),
                     f'{lcn} ({len_})', font_size=Pt(13), font_color=NAVY, bold=True)
            for ni, name in enumerate(members):
                add_oval(s, qx + Inches(0.2), qy + Inches(0.55) + ni * Inches(0.4),
                         str(ni + 1), size=Inches(0.3), bg=NAVY)
                add_text(s, qx + Inches(0.6), qy + Inches(0.5) + ni * Inches(0.4),
                         Inches(2.5), Inches(0.35), name,
                         font_size=BODY_SIZE, font_color=DARK_GRAY)
        add_text(s, LM, sgt + sgch - Inches(0.3), Inches(1.8), Inches(0.6),
                 y_label.replace('↑', '\n↑'), font_size=SMALL_SIZE, font_color=MED_GRAY,
                 alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, sgl + sgcw - Inches(0.5), sgt + 2 * sgch + Inches(0.1),
                 Inches(2.5), Inches(0.3), x_label,
                 font_size=SMALL_SIZE, font_color=MED_GRAY, alignment=PP_ALIGN.CENTER)
        if summary:
            add_rect(s, LM, Inches(5.5), CW, Inches(0.8), BG_GRAY)
            add_text(s, LM + Inches(0.3), Inches(5.5), CW - Inches(0.6), Inches(0.8),
                     summary, font_size=BODY_SIZE, font_color=NAVY, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE)
        self._footer(s, source)
        return s

    def decision_tree(self, title, root, branches, right_panel=None, source=''):
        """#60 Decision Tree — root → L1 → L2 hierarchy with connector lines.
        root: (label,).
        branches: list of (L1_title, L1_metric, L1_color, children:list[(name, metric)]).
        right_panel: (panel_title, points:list[str]) or None.
        """
        s = self._ns()
        add_action_title(s, title)
        L0x = LM + Inches(0.3); L0y = Inches(2.5); L0w = Inches(2.2); L0h = Inches(1.2)
        add_rect(s, L0x, L0y, L0w, L0h, NAVY)
        add_text(s, L0x, L0y, L0w, L0h, root[0],
                 font_size=EMPHASIS_SIZE, font_color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        L1x = L0x + L0w + Inches(0.6); L1w = Inches(2.0); L1h = Inches(0.9)
        L1_boxes = []
        colors_L1 = [ACCENT_BLUE, ACCENT_ORANGE, ACCENT_GREEN, ACCENT_RED]
        for i, (l1_title, l1_metric, l1_color, children) in enumerate(branches):
            L1y = Inches(1.5) + i * Inches(2.2)
            L1_boxes.append((L1x, L1y, L1w, L1h))
            # Connector lines
            conn_start = L0y + L0h // 2
            conn_end = L1y + L1h // 2
            mid_x = L0x + L0w + (L1x - L0x - L0w) // 2
            add_hline(s, L0x + L0w, conn_start, mid_x - L0x - L0w, LINE_GRAY, Pt(1))
            vy_top = min(conn_start, conn_end)
            vy_h = abs(conn_end - conn_start)
            if vy_h > 0:
                add_rect(s, mid_x, vy_top, Pt(1), vy_h, LINE_GRAY)
            add_hline(s, mid_x, conn_end, L1x - mid_x, LINE_GRAY, Pt(1))
            add_rect(s, L1x, L1y, L1w, L1h, l1_color)
            add_text(s, L1x, L1y, L1w, Inches(0.5), l1_title,
                     font_size=BODY_SIZE, font_color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, L1x, L1y + Inches(0.45), L1w, Inches(0.4), l1_metric,
                     font_size=SUB_HEADER_SIZE, font_color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            # L2 children
            L2x = L1x + L1w + Inches(0.6); L2w = Inches(2.0); L2h = Inches(0.65)
            for li, (c_name, c_metric) in enumerate(children):
                L2y = L1y - Inches(0.3) + li * Inches(0.8)
                c_start = L1y + L1h // 2; c_end = L2y + L2h // 2
                mid_x2 = L1x + L1w + (L2x - L1x - L1w) // 2
                add_hline(s, L1x + L1w, c_start, mid_x2 - L1x - L1w, LINE_GRAY, Pt(0.5))
                vy2_top = min(c_start, c_end); vy2_h = abs(c_end - c_start)
                if vy2_h > 0:
                    add_rect(s, mid_x2, vy2_top, Pt(0.5), vy2_h, LINE_GRAY)
                add_hline(s, mid_x2, c_end, L2x - mid_x2, LINE_GRAY, Pt(0.5))
                add_rect(s, L2x, L2y, L2w, L2h, BG_GRAY)
                add_text(s, L2x + Inches(0.1), L2y, L2w * 0.6, L2h,
                         c_name, font_size=SMALL_SIZE, font_color=DARK_GRAY,
                         anchor=MSO_ANCHOR.MIDDLE)
                add_text(s, L2x + L2w * 0.6, L2y, L2w * 0.4, L2h,
                         c_metric, font_size=BODY_SIZE, font_color=NAVY, bold=True,
                         anchor=MSO_ANCHOR.MIDDLE, alignment=PP_ALIGN.CENTER)
        if right_panel:
            pt, pp = right_panel
            L3x = L1x + L1w + Inches(0.6) + Inches(2.0) + Inches(0.4)
            L3w = CW - (L3x - LM)
            if L3w > Inches(0.5):
                add_rect(s, L3x, Inches(1.2), L3w, Inches(4.5), BG_GRAY)
                add_text(s, L3x + Inches(0.15), Inches(1.3), L3w - Inches(0.3), Inches(0.3),
                         pt, font_size=BODY_SIZE, font_color=NAVY, bold=True)
                add_text(s, L3x + Inches(0.15), Inches(1.7), L3w - Inches(0.3), Inches(3.5),
                         pp, font_size=SMALL_SIZE, font_color=DARK_GRAY,
                         line_spacing=Pt(8))
        self._footer(s, source)
        return s

    def metric_comparison(self, title, metrics, source=''):
        """#62 Metric Comparison — before/after row cards with delta badges.
        metrics: list of (label, before_val, after_val, delta_str).
        """
        s = self._ns()
        add_action_title(s, title)
        rh = Inches(0.95); mcw = Inches(4.0)
        bx = LM + Inches(0.5); ax = LM + Inches(6.5)
        dx = ax + mcw + Inches(0.3)
        add_text(s, bx, Inches(1.3), mcw, Inches(0.3), '之前',
                 font_size=Pt(13), font_color=MED_GRAY, alignment=PP_ALIGN.CENTER)
        add_text(s, ax, Inches(1.3), mcw, Inches(0.3), '之后',
                 font_size=Pt(13), font_color=MED_GRAY, alignment=PP_ALIGN.CENTER)
        add_text(s, dx, Inches(1.3), Inches(1.5), Inches(0.3), '变化',
                 font_size=Pt(13), font_color=MED_GRAY, alignment=PP_ALIGN.CENTER)
        for i, (label, before, after, delta) in enumerate(metrics):
            ry = Inches(1.8) + i * rh
            add_rect(s, bx, ry, mcw, rh - Inches(0.1), BG_GRAY)
            add_text(s, bx + Inches(0.2), ry, Inches(1.5), rh - Inches(0.1), label,
                     font_size=SMALL_SIZE, font_color=MED_GRAY, anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, bx + Inches(1.8), ry, Inches(2.0), rh - Inches(0.1), before,
                     font_size=Pt(22), font_color=DARK_GRAY, bold=True,
                     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, bx + mcw + Inches(0.1), ry, Inches(1.5), rh - Inches(0.1), '→',
                     font_size=Pt(24), font_color=LINE_GRAY,
                     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            add_rect(s, ax, ry, mcw, rh - Inches(0.1), LIGHT_BLUE)
            add_text(s, ax + Inches(0.2), ry, Inches(1.5), rh - Inches(0.1), label,
                     font_size=SMALL_SIZE, font_color=ACCENT_BLUE,
                     anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, ax + Inches(1.8), ry, Inches(2.0), rh - Inches(0.1), after,
                     font_size=Pt(22), font_color=NAVY, bold=True,
                     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            is_pos = delta.startswith('+')
            dc = ACCENT_GREEN if is_pos else ACCENT_RED
            add_rect(s, dx + Inches(0.1), ry + Inches(0.15),
                     Inches(1.2), rh - Inches(0.35), dc)
            add_text(s, dx + Inches(0.1), ry + Inches(0.15),
                     Inches(1.2), rh - Inches(0.35), delta,
                     font_size=EMPHASIS_SIZE, font_color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._footer(s, source)
        return s

    def icon_grid(self, title, items, cols=3, source=''):
        """#63 Icon Grid — grid of icon cards.
        items: list of (item_title, description, accent_color).
        """
        s = self._ns()
        add_action_title(s, title)
        rows_ = math.ceil(len(items) / cols)
        celw = CW / cols - Inches(0.15)
        celh = Inches(2.2)
        ty = Inches(1.3)
        for i, (ititle, desc, color) in enumerate(items):
            col = i % cols; row = i // cols
            ix = LM + col * (celw + Inches(0.15))
            iy = ty + row * (celh + Inches(0.1))
            add_rect(s, ix, iy, celw, celh, WHITE)
            add_rect(s, ix, iy, celw, Inches(0.06), color)
            isz = Inches(0.6)
            o = s.shapes.add_shape(MSO_SHAPE.OVAL, ix + Inches(0.3), iy + Inches(0.25),
                                   isz, isz)
            o.fill.solid(); o.fill.fore_color.rgb = color
            o.line.fill.background(); _clean_shape(o)
            add_text(s, ix + Inches(0.3), iy + Inches(0.25), isz, isz, ititle[0],
                     font_size=SUB_HEADER_SIZE, font_color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, ix + Inches(1.1), iy + Inches(0.25), celw - Inches(1.3), Inches(0.4),
                     ititle, font_size=EMPHASIS_SIZE, font_color=color, bold=True,
                     anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, ix + Inches(0.3), iy + Inches(1.0), celw - Inches(0.6), Inches(1.0),
                     desc, font_size=BODY_SIZE, font_color=DARK_GRAY)
        self._footer(s, source)
        return s

    def agenda(self, title, headers, items, footer_text='', source=''):
        """#66 Agenda — table-style meeting agenda.
        headers: list of (label, width).
        items: list of (*values, item_type) — type: 'key'|'normal'|'break'.
        """
        s = self._ns()
        add_action_title(s, title)
        hy = Inches(1.3)
        hx = LM
        for label, w in headers:
            add_text(s, hx, hy, w, Inches(0.35), label,
                     font_size=Pt(13), font_color=NAVY, bold=True)
            hx += w
        total_w = sum(w for _, w in headers)
        add_hline(s, LM, hy + Inches(0.35), total_w, BLACK, Pt(0.75))
        rh = Inches(0.55)
        status_w = headers[-1][1] if headers else Inches(1.5)
        data_headers = headers[:-1]
        for i, item in enumerate(items):
            *vals, itype = item
            ry = Inches(1.8) + i * rh
            if itype == 'break':
                add_rect(s, LM, ry, total_w, rh, BG_GRAY)
            elif itype == 'key':
                add_rect(s, LM, ry, Inches(0.06), rh, ACCENT_BLUE)
            rx = LM
            for vi, (_, w) in enumerate(data_headers):
                val = vals[vi] if vi < len(vals) else ''
                fc = MED_GRAY if itype == 'break' else DARK_GRAY
                bld = (itype == 'key')
                add_text(s, rx, ry, w, rh, val,
                         font_size=BODY_SIZE, font_color=fc, bold=bld,
                         anchor=MSO_ANCHOR.MIDDLE)
                rx += w
            if itype == 'key':
                add_rect(s, rx + Inches(0.1), ry + Inches(0.12),
                         Inches(1.0), rh - Inches(0.24), LIGHT_BLUE)
                add_text(s, rx + Inches(0.1), ry, Inches(1.0), rh, '★ 重点',
                         font_size=Pt(11), font_color=ACCENT_BLUE, bold=True,
                         alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            if i < len(items) - 1:
                add_hline(s, LM, ry + rh, total_w, LINE_GRAY, Pt(0.25))
        if footer_text:
            add_text(s, LM, Inches(6.1), CW, Inches(0.3), footer_text,
                     font_size=Pt(10), font_color=MED_GRAY)
        self._footer(s, source)
        return s

    def two_col_image_grid(self, title, items, source=''):
        """#68 Two-Column Image + Text Grid — 2×2 image-text cards.
        items: list of (card_title, description, accent_color, image_label).
        """
        s = self._ns()
        add_action_title(s, title)
        gcw = CW / 2 - Inches(0.15); gch = Inches(2.2)
        imgw = Inches(2.8); imgh = Inches(1.8); ty = Inches(1.3)
        for i, (ctitle, desc, color, img_label) in enumerate(items):
            col = i % 2; row = i // 2
            cx = LM + col * (gcw + Inches(0.15))
            cy = ty + row * (gch + Inches(0.1))
            add_rect(s, cx, cy, gcw, gch, WHITE)
            add_image_placeholder(s, cx + Inches(0.15), cy + Inches(0.15), imgw, imgh,
                                  img_label)
            tx = cx + imgw + Inches(0.3); tw = gcw - imgw - Inches(0.45)
            add_rect(s, tx - Inches(0.05), cy, Inches(0.06), gch, color)
            add_text(s, tx + Inches(0.1), cy + Inches(0.2), tw, Inches(0.35),
                     ctitle, font_size=EMPHASIS_SIZE, font_color=color, bold=True)
            add_text(s, tx + Inches(0.1), cy + Inches(0.6), tw, Inches(1.2),
                     desc, font_size=BODY_SIZE, font_color=DARK_GRAY)
        self._footer(s, source)
        return s

    def numbered_list_panel(self, title, items, panel=None, source=''):
        """#69 Numbered List + Side Panel — left numbered list + right accent panel.
        items: list of (item_title, description).
        panel: dict with 'subtitle','big_number','big_label','metrics':list[(label,value)].
        """
        s = self._ns()
        add_action_title(s, title)
        nlw = Inches(7.5)
        nty = Inches(1.3)

        # Dynamic row height: fill the available vertical space evenly
        n_items = len(items)
        avail_h = Inches(4.8)  # match panel height
        row_h = avail_h / max(n_items, 1)
        # Clamp: min 0.85", max 1.8"
        row_h = max(Inches(0.85), min(row_h, Inches(1.8)))
        # Title takes ~0.35", rest is for description
        title_h = Inches(0.3)
        desc_h = row_h - title_h - Inches(0.15)  # 0.15" gap + separator

        for i, (ititle, desc) in enumerate(items):
            ry = nty + i * row_h
            add_oval(s, LM, ry + Inches(0.05), str(i + 1), bg=NAVY)
            add_text(s, LM + Inches(0.6), ry, nlw - Inches(0.6), title_h,
                     ititle, font_size=Pt(15), font_color=NAVY, bold=True)
            add_text(s, LM + Inches(0.6), ry + title_h + Inches(0.05),
                     nlw - Inches(0.6), desc_h,
                     desc, font_size=BODY_SIZE, font_color=DARK_GRAY)
            if i < len(items) - 1:
                sep_y = ry + row_h - Inches(0.05)
                add_hline(s, LM + Inches(0.6), sep_y,
                          nlw - Inches(0.8), LINE_GRAY, Pt(0.25))
        if panel:
            nrx = LM + nlw + Inches(0.3); nrw = CW - nlw - Inches(0.3)
            add_rect(s, nrx, Inches(1.3), nrw, Inches(4.8), NAVY)
            if 'subtitle' in panel:
                add_text(s, nrx + Inches(0.3), Inches(1.6), nrw - Inches(0.6), Inches(0.3),
                         panel['subtitle'], font_size=BODY_SIZE,
                         font_color=RGBColor(0xCC, 0xCC, 0xCC))
            if 'big_number' in panel:
                add_text(s, nrx + Inches(0.3), Inches(1.8), nrw - Inches(0.6), Inches(0.6),
                         panel['big_number'], font_size=Pt(36), font_color=WHITE,
                         bold=True, alignment=PP_ALIGN.CENTER)
            if 'big_label' in panel:
                add_text(s, nrx + Inches(0.3), Inches(2.5), nrw - Inches(0.6), Inches(0.3),
                         panel['big_label'], font_size=Pt(13),
                         font_color=RGBColor(0xCC, 0xCC, 0xCC),
                         alignment=PP_ALIGN.CENTER)
            add_hline(s, nrx + Inches(0.3), Inches(3.1), nrw - Inches(0.6),
                      RGBColor(0x33, 0x44, 0x55), Pt(0.5))
            metrics = panel.get('metrics', [])
            for mi, (label, val) in enumerate(metrics):
                my = Inches(3.4) + mi * Inches(0.7)
                add_text(s, nrx + Inches(0.3), my, nrw - Inches(0.6), Inches(0.3),
                         label, font_size=Pt(11),
                         font_color=RGBColor(0xAA, 0xAA, 0xAA))
                add_text(s, nrx + Inches(0.3), my + Inches(0.3), nrw - Inches(0.6), Inches(0.3),
                         val, font_size=SUB_HEADER_SIZE, font_color=WHITE, bold=True)
        self._footer(s, source)
        return s

    # ═══════════════════════════════════════════
    # MULTI-BAR-PANEL CHART (#71 — McKinsey editorial bar chart panels)
    # ═══════════════════════════════════════════

    def multi_bar_panel(self, title, panels, connectors=None, footnotes=None, source=''):
        """#71 Multi-Bar-Panel Chart — 2-3 side-by-side bar chart panels with
        CAGR trend arrows and value labels on each bar.
        McKinsey / BCG editorial chart style.

        Parameters
        ----------
        title : str
            Action title (top of slide).
        panels : list[dict]
            Each panel dict has:
              - 'title': str — panel headline (supports **bold** markup for key numbers)
              - 'unit': str — Y-axis unit label, e.g. '万人'
              - 'legend': str — legend text, e.g. '15-64岁人口数量'
              - 'categories': list[str] — X-axis labels (years)
              - 'values': list[int|float] — bar heights
              - 'bar_color': RGBColor (optional, default NAVY)
              - 'cagr': list[dict] (optional) — CAGR annotations, each dict:
                  {'rate': str, 'start': int, 'end': int}  # index into categories
              - 'highlight_idx': list[int] (optional) — indices of bars to highlight
              - 'highlight_color': RGBColor (optional, default ACCENT_BLUE)
              - 'value_format': str (optional) — f-string format, e.g. '{:,.0f}'
        connectors : list[str] or None
            (Kept for API compat — ignored in v2. Panels are flush side-by-side.)
        footnotes : list[str] or None
            Footnote lines below the chart area.
        source : str
            Source attribution.
        """
        import re
        s = self._ns()
        add_action_title(s, title)

        n_panels = len(panels)

        # --- Layout: panels flush side-by-side, no connector circles ---
        # CRITICAL: all coordinates must be int (EMU) to avoid XML float errors
        _I = lambda v: int(v)  # shorthand for int coercion
        panel_gap = Inches(0.25)               # gap between panels
        total_gap = panel_gap * max(n_panels - 1, 0)
        pw = _I((CW - total_gap) / n_panels)   # panel width

        title_area_top = CONTENT_TOP
        title_area_h = Inches(0.65)
        unit_row_y = title_area_top + title_area_h + Inches(0.02)
        cagr_row_y = unit_row_y + Inches(0.28)      # CAGR rate text
        cagr_arrow_y = cagr_row_y + Inches(0.22)    # CAGR arrow line
        chart_top = cagr_arrow_y + Inches(0.08)      # bars start
        chart_bot = Inches(5.85)                     # bars bottom
        chart_h = chart_bot - chart_top
        xaxis_y = chart_bot + Inches(0.02)           # X-axis label y

        footnote_y = Inches(6.25)

        for pi, panel in enumerate(panels):
            px = _I(LM + pi * (pw + panel_gap))

            bar_color = panel.get('bar_color', NAVY)
            hl_color = panel.get('highlight_color', ACCENT_BLUE)
            hl_idx = set(panel.get('highlight_idx', []))
            values = panel['values']
            categories = panel['categories']
            n_bars = len(values)
            max_val = max(values) if values else 1
            vfmt = panel.get('value_format', '{:,.0f}')

            # --- Panel title with number prefix: "1. xxx **bold** xxx" ---
            ptitle = panel['title']
            txBox = s.shapes.add_textbox(px, title_area_top, pw, title_area_h)
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.auto_size = None
            bodyPr = tf._txBody.find(qn('a:bodyPr'))
            bodyPr.set('anchor', 't')
            for attr in ['lIns', 'tIns', 'rIns', 'bIns']:
                bodyPr.set(attr, '36576')
            p = tf.paragraphs[0]
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            p.line_spacing = 1.05

            # Add number prefix: "1. " in bold
            num_run = p.add_run()
            num_run.text = f'{pi + 1}. '
            num_run.font.size = SUB_HEADER_SIZE
            num_run.font.name = FONT_BODY
            num_run.font.color.rgb = BLACK
            num_run.font.bold = True
            set_ea_font(num_run, FONT_EA)

            # Parse **bold** segments in title
            segments = re.split(r'(\*\*.*?\*\*)', ptitle)
            for seg in segments:
                if seg.startswith('**') and seg.endswith('**'):
                    run = p.add_run()
                    run.text = seg[2:-2]
                    run.font.size = SUB_HEADER_SIZE
                    run.font.name = FONT_BODY
                    run.font.color.rgb = BLACK
                    run.font.bold = True
                    set_ea_font(run, FONT_EA)
                elif seg:
                    run = p.add_run()
                    run.text = seg
                    run.font.size = SUB_HEADER_SIZE
                    run.font.name = FONT_BODY
                    run.font.color.rgb = BLACK
                    run.font.bold = False
                    set_ea_font(run, FONT_EA)

            # --- Unit + legend row (plain text, no color square) ---
            unit = panel.get('unit', '')
            legend = panel.get('legend', '')
            unit_legend_text = unit
            if legend:
                unit_legend_text = f'{unit}          {legend}'
            if unit_legend_text:
                add_text(s, px, unit_row_y, pw, Inches(0.25),
                         unit_legend_text, font_size=SMALL_SIZE, font_color=MED_GRAY)

            # --- Pre-calculate bar geometry (needed for both CAGR arrows and bars) ---
            cagr_list = panel.get('cagr', [])
            bar_area_w = _I(pw - Inches(0.15))
            bar_x_start = _I(px + Inches(0.05))
            bar_spacing = _I(bar_area_w / n_bars)
            bar_w = _I(bar_spacing * 0.72)       # wider bars, tighter gaps
            bar_gap_w = _I(bar_spacing * 0.28)

            scale_min = 0
            scale_range = max_val - scale_min if max_val > scale_min else 1
            usable_h = _I(chart_h * 0.82)  # leave headroom for value labels

            # --- CAGR trend annotations (slope follows bar tops) ---
            import math
            for ci, cagr in enumerate(cagr_list):
                rate = cagr['rate']
                c_start = cagr['start']
                c_end = cagr['end']
                rate_color = ACCENT_RED if rate.lstrip().startswith('-') else ACCENT_GREEN

                # Compute bar-top Y for start and end bars
                def _bar_top_y(idx):
                    v = values[idx]
                    ratio = (v - scale_min) / scale_range if scale_range > 0 else 0
                    bh = _I(usable_h * ratio)
                    if bh < Inches(0.05):
                        bh = Inches(0.05)
                    return _I(chart_bot - bh)

                start_top_y = _bar_top_y(c_start)
                end_top_y = _bar_top_y(c_end)

                # Arrow line X: from center of start bar to center of end bar
                sx_center = _I(bar_x_start + bar_spacing * c_start + bar_gap_w / 2 + bar_w * 0.5)
                ex_center = _I(bar_x_start + bar_spacing * c_end + bar_gap_w / 2 + bar_w * 0.5)

                # Arrow line Y: offset above bar tops
                arrow_y_offset = Inches(0.27)  # gap above bar top (closer to bars)
                sy = _I(start_top_y - arrow_y_offset)
                ey = _I(end_top_y - arrow_y_offset)

                # Rate label: centered horizontally, at the higher (min-Y) point minus offset
                mid_x = _I((sx_center + ex_center) / 2)
                label_ref_y = _I(min(sy, ey) - Inches(0.28))
                rate_label_w = Inches(1.8)
                add_text(s, _I(mid_x - rate_label_w / 2), label_ref_y,
                         rate_label_w, Inches(0.25), rate,
                         font_size=EMPHASIS_SIZE, font_color=rate_color, bold=True,
                         alignment=PP_ALIGN.CENTER)

                # Draw sloped arrow using a single RIGHT_ARROW shape
                # (no tail notch, arrowhead and line unified)
                dx = ex_center - sx_center
                dy = ey - sy  # positive dy = line goes down, negative = goes up
                line_len = _I(math.sqrt(dx * dx + dy * dy))
                if line_len > 0:
                    angle_rad = math.atan2(dy, dx)
                    angle_deg = math.degrees(angle_rad)
                    # Arrow shape: height controls head spread; shaft thinness via adj1
                    # We want shaft ~1.5px but need enough height for visible arrowhead
                    arrow_h = Inches(0.13)  # overall height (arrowhead spread)
                    # Place at center of the line (rotation is around shape center)
                    cx = _I((sx_center + ex_center) / 2)
                    cy = _I((sy + ey) / 2)
                    arrow_shape = s.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        _I(cx - line_len / 2), _I(cy - arrow_h / 2),
                        line_len, arrow_h)
                    arrow_shape.fill.solid()
                    arrow_shape.fill.fore_color.rgb = rate_color
                    arrow_shape.line.fill.background()
                    _clean_shape(arrow_shape)
                    arrow_shape.rotation = angle_deg
                    # Adjust arrow proportions via XML adj handles:
                    # RIGHT_ARROW adj handles:
                    #   adj1 = shaft width as % of shape height (50000 = 50%)
                    #   adj2 = arrowhead length as % of shape width
                    # Target: shaft ~1.5px thin, arrowhead 1.5x longer
                    from lxml import etree
                    avLst = arrow_shape._element.find(
                        './/{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
                    if avLst is not None:
                        for child in list(avLst):
                            avLst.remove(child)
                        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                        # adj1: shaft width ratio — very thin shaft (~1.5px)
                        # shape height is 0.13" ≈ 9.36pt, 1.5px ≈ 1.125pt
                        # ratio = 1.125 / 9.36 ≈ 12%, so val ~12000
                        gd1 = etree.SubElement(avLst, f'{{{ns}}}gd')
                        gd1.set('name', 'adj1')
                        gd1.set('fmla', 'val 12000')  # ~12% = ultra-thin shaft
                        # adj2: arrowhead length — 1.5x default
                        # default arrowhead length ≈ 50000 (50% of height as width)
                        # 1.5x → 75000
                        gd2 = etree.SubElement(avLst, f'{{{ns}}}gd')
                        gd2.set('name', 'adj2')
                        gd2.set('fmla', 'val 75000')  # 1.5x longer arrowhead

            # --- Bar chart ---

            for bi, val in enumerate(values):
                bx = _I(bar_x_start + bar_spacing * bi + bar_gap_w / 2)
                ratio = (val - scale_min) / scale_range if scale_range > 0 else 0
                bh = _I(usable_h * ratio)
                if bh < Inches(0.05):
                    bh = Inches(0.05)
                by = _I(chart_bot - bh)
                bc = hl_color if bi in hl_idx else bar_color
                add_rect(s, bx, by, bar_w, bh, bc)

                # Value label — tight above bar
                val_text = vfmt.format(val)
                val_font = SMALL_SIZE
                label_h = Inches(0.22)
                add_text(s, _I(bx - Inches(0.08)), _I(by - label_h - Inches(0.02)),
                         _I(bar_w + Inches(0.16)), label_h,
                         val_text, font_size=val_font, font_color=BLACK, bold=False,
                         alignment=PP_ALIGN.CENTER)

            # --- X-axis labels ---
            for bi, cat in enumerate(categories):
                bx = _I(bar_x_start + bar_spacing * bi + bar_gap_w / 2)
                add_text(s, _I(bx - Inches(0.05)), xaxis_y,
                         _I(bar_w + Inches(0.1)), Inches(0.22),
                         cat, font_size=SMALL_SIZE, font_color=DARK_GRAY,
                         alignment=PP_ALIGN.CENTER)

        # --- Footnotes ---
        if footnotes:
            for fi, fn in enumerate(footnotes):
                add_text(s, LM, footnote_y + Inches(0.18) * fi, CW, Inches(0.18),
                         fn, font_size=FOOTNOTE_SIZE, font_color=MED_GRAY)

        self._footer(s, source)
        return s

    # ═══════════════════════════════════════════
    # SAVE
    # ═══════════════════════════════════════════

    def save(self, outpath):
        """Save presentation and run full XML cleanup."""
        outdir = os.path.dirname(outpath)
        if outdir:
            os.makedirs(outdir, exist_ok=True)
        self.prs.save(outpath)
        full_cleanup(outpath)
        print(f"✅ Saved: {outpath} ({self._page} slides, {os.path.getsize(outpath):,} bytes)")
        return outpath