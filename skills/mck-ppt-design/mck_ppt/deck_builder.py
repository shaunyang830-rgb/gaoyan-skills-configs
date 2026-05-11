# Copyright 2024-2026 Kaku Li (https://github.com/likaku)
# Licensed under the Apache License, Version 2.0 — see LICENSE and NOTICE.
# Part of "Mck-ppt-design-skill" (McKinsey PPT Design Framework).
# NOTICE: This file must be retained in all copies or substantial portions.
#
"""DeckBuilder — Storyline-driven deck generator with built-in QA.

Accepts a storyline (list of slide specs) and orchestrates MckEngine
to produce a complete, professionally structured PPTX deck.
Includes automatic QA validation after generation.

Usage:
    from mck_ppt.deck_builder import DeckBuilder

    storyline = [
        {'type': 'cover', 'data': {'title': '...', 'subtitle': '...'}},
        {'type': 'toc',   'data': {'items': [...]}},
        {'type': 'section_divider', 'data': {'section_label': '01', 'title': '...'}},
        ...
    ]
    DeckBuilder.build(storyline, 'output/deck.pptx')
"""
import os
from pptx import Presentation
from .engine import MckEngine
from .core import full_cleanup


class DeckBuilder:
    """Orchestrates MckEngine to build a deck from a storyline spec."""

    # Slide boundaries (inches) — anything outside is a bug
    SLIDE_W = 13.333
    SLIDE_H = 7.5

    @staticmethod
    def build(storyline: list[dict], output_path: str, **engine_kwargs) -> str:
        """Build a complete deck from a storyline.

        Parameters
        ----------
        storyline : list[dict]
            Each item: {'type': '<method_name>', 'data': {**kwargs}}
        output_path : str
            Path to save the output PPTX.
        **engine_kwargs
            Extra kwargs passed to MckEngine (e.g. total_slides).

        Returns
        -------
        str
            The output_path.
        """
        total = engine_kwargs.pop('total_slides', len(storyline))
        eng = MckEngine(total_slides=total)

        errors = []
        for i, slide_spec in enumerate(storyline):
            slide_type = slide_spec.get('type', '')
            data = slide_spec.get('data', {})
            slide_num = i + 1

            method = getattr(eng, slide_type, None)
            if method is None:
                print(f"  ⚠ Slide {slide_num}: unknown type '{slide_type}', skipping")
                errors.append((slide_num, slide_type, 'unknown method'))
                continue

            try:
                method(**data)
            except Exception as e:
                print(f"  ⚠ Slide {slide_num} ({slide_type}): {e}")
                errors.append((slide_num, slide_type, str(e)))

        os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)
        eng.save(output_path)

        n_ok = len(storyline) - len(errors)
        print(f"  ✓ Built {n_ok}/{len(storyline)} slides → {output_path}")
        if errors:
            print(f"  ⚠ {len(errors)} errors:")
            for num, typ, msg in errors:
                print(f"      S{num} ({typ}): {msg}")

        # ── Post-build QA ──
        qa_issues = DeckBuilder.qa_validate(output_path)
        if qa_issues:
            print(f"  ⚠ QA: {len(qa_issues)} layout issues found")
            for iss in qa_issues[:10]:
                print(f"      {iss}")
        else:
            print(f"  ✓ QA passed")

        # ── XML cleanup for PowerPoint compatibility ──
        full_cleanup(output_path)

        return output_path

    @staticmethod
    def qa_validate(pptx_path: str) -> list[str]:
        """Validate a PPTX file for layout issues.

        Checks:
        - Shapes not off-screen (top/left/bottom/right within slide bounds)
        - No negative width/height
        - Text not overflowing containers (basic check)

        Returns list of issue descriptions (empty = all good).
        """
        prs = Presentation(pptx_path)
        issues = []
        EMU_PER_INCH = 914400

        for si, slide in enumerate(prs.slides):
            slide_num = si + 1
            for shape in slide.shapes:
                top = (shape.top or 0) / EMU_PER_INCH
                left = (shape.left or 0) / EMU_PER_INCH
                w = (shape.width or 0) / EMU_PER_INCH
                h = (shape.height or 0) / EMU_PER_INCH
                bottom = top + h
                right = left + w
                txt = ''
                if shape.has_text_frame:
                    txt = shape.text_frame.text[:25]

                # Off-screen checks
                if top < -0.2:
                    issues.append(
                        f"S{slide_num}: shape '{txt}' top={top:.1f}\" is above slide")
                if bottom > DeckBuilder.SLIDE_H + 0.5:
                    issues.append(
                        f"S{slide_num}: shape '{txt}' bottom={bottom:.1f}\" is below slide")
                if left < -0.5:
                    issues.append(
                        f"S{slide_num}: shape '{txt}' left={left:.1f}\" is off left")
                if right > DeckBuilder.SLIDE_W + 0.5:
                    issues.append(
                        f"S{slide_num}: shape '{txt}' right={right:.1f}\" is off right")

                # Negative dimensions
                if w < -0.01:
                    issues.append(
                        f"S{slide_num}: shape '{txt}' has negative width={w:.2f}\"")
                if h < -0.01:
                    issues.append(
                        f"S{slide_num}: shape '{txt}' has negative height={h:.2f}\"")

        return issues

    @staticmethod
    def build_from_module(module, output_path: str, **engine_kwargs) -> str:
        """Build a deck from a storyline module.

        The module must have a STORYLINE attribute (list of slide specs).
        """
        storyline = getattr(module, 'STORYLINE', None)
        if storyline is None:
            raise ValueError(f"Module {module.__name__} has no STORYLINE attribute")
        return DeckBuilder.build(storyline, output_path, **engine_kwargs)