# Copyright 2024-2026 Kaku Li (https://github.com/likaku)
# Licensed under the Apache License, Version 2.0 — see LICENSE and NOTICE.
# Part of "Mck-ppt-design-skill" (McKinsey PPT Design Framework).
# NOTICE: This file must be retained in all copies or substantial portions.
#
"""
Post-generation Review + Auto-fix Pipeline for mck_ppt.

Four-stage flow:
  1. Page Brief — structure content into page_objective / one_message / mece_buckets
  2. Dual QA — Narrative QA (content) + Layout QA (geometry, via qa.py)
  3. Auto-fix — priority chain: 去冗余 → 统一语言 → 压缩句式 → 重构层级 → 字号微调
     (NO layout/布局 changes — user constraint)
  4. Gate — 0 ERROR = pass; otherwise iterate (max N rounds)

Usage:
    from mck_ppt.review import SlideReviewer
    reviewer = SlideReviewer("output/deck.pptx")
    report = reviewer.run()          # read-only audit
    report.print_summary()

    # Auto-fix (mutates the file in-place, up to max_rounds)
    from mck_ppt.review import AutoFixPipeline
    result = AutoFixPipeline("output/deck.pptx").run(max_rounds=3)
    print(result)
"""

from __future__ import annotations

import copy
import math
import os
import re
from dataclasses import dataclass, field
from typing import List, Optional, Tuple, Dict

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn

from .qa import (
    PptQA, QAReport, QAIssue, Severity,
    _estimate_text_height, _bbox,
    SW, SH, LM, CW, CONTENT_TOP, SOURCE_Y,
    OVERFLOW_TOLERANCE, TEXT_OVERFLOW_LINE_RATIO,
    MIN_FONT_SIZE,
)
from .constants import (
    BODY_SIZE, SMALL_SIZE, FOOTNOTE_SIZE, ACTION_TITLE_SIZE,
    SUB_HEADER_SIZE, EMPHASIS_SIZE,
)


# ═══════════════════════════════════════════════════════════════
# Content density limits (per-layout experience caps)
# ═══════════════════════════════════════════════════════════════

# Max characters per text box before we flag as "too dense"
# Keyed by approximate box height in inches (rounded down)
CHAR_DENSITY_LIMITS = {
    0.2: 15,   # tiny label boxes
    0.3: 30,
    0.4: 50,
    0.5: 70,
    0.6: 90,
    0.8: 130,
    1.0: 180,
    1.5: 300,
    2.0: 450,
    3.0: 700,
    5.0: 1200,
}

# Title length warning thresholds
ACTION_TITLE_MAX_CHARS = 45
SUBTITLE_MAX_CHARS = 30

# Font size guard rails for autofix
TITLE_MIN_PT = 20   # action titles cannot shrink below this
BODY_MIN_PT = 11    # body text floor
SMALL_MIN_PT = 9    # small/footnote floor


# ═══════════════════════════════════════════════════════════════
# Narrative QA Issue
# ═══════════════════════════════════════════════════════════════

@dataclass
class NarrativeIssue:
    slide_num: int
    severity: str       # ERROR / WARNING / INFO
    category: str       # e.g. "no_action_title", "lang_mix", "density"
    message: str
    shape_name: str = ""
    suggestion: str = ""


@dataclass
class NarrativeReport:
    total_slides: int
    issues: List[NarrativeIssue] = field(default_factory=list)

    @property
    def errors(self):
        return [i for i in self.issues if i.severity == Severity.ERROR]

    @property
    def warnings(self):
        return [i for i in self.issues if i.severity == Severity.WARNING]

    def print_summary(self):
        print(f"\n{'─'*60}")
        print(f"  Narrative QA — {self.total_slides} slides")
        print(f"  Errors: {len(self.errors)}  Warnings: {len(self.warnings)}")
        print(f"{'─'*60}")
        for iss in self.issues:
            icon = {"ERROR": "❌", "WARNING": "⚠️ ", "INFO": "ℹ️ "}[iss.severity]
            print(f"  S{iss.slide_num} {icon} [{iss.category}] {iss.message}")
            if iss.suggestion:
                print(f"       → {iss.suggestion}")
        print()


# ═══════════════════════════════════════════════════════════════
# NarrativeReviewer — content-level checks
# ═══════════════════════════════════════════════════════════════

class NarrativeReviewer:
    """Checks each slide for narrative quality issues."""

    def __init__(self, filepath: str):
        self.filepath = filepath
        self.prs = Presentation(filepath)

    def run(self) -> NarrativeReport:
        issues: List[NarrativeIssue] = []
        for idx, slide in enumerate(self.prs.slides, 1):
            issues.extend(self._check_slide(idx, slide))
        return NarrativeReport(
            total_slides=len(self.prs.slides),
            issues=issues,
        )

    def _check_slide(self, num: int, slide) -> List[NarrativeIssue]:
        issues = []
        shapes = list(slide.shapes)
        all_text = self._collect_text(shapes)

        # 1. Check for mixed Chinese/English where unnecessary
        issues.extend(self._check_lang_mix(num, shapes))

        # 2. Check text density per box
        issues.extend(self._check_density(num, shapes))

        # 3. Check title length
        issues.extend(self._check_title_length(num, shapes))

        return issues

    def _collect_text(self, shapes) -> str:
        parts = []
        for s in shapes:
            if s.has_text_frame:
                parts.append(s.text_frame.text)
        return "\n".join(parts)

    # ── Mixed language check ──────────────────────────────────
    _EN_JARGON_RE = re.compile(
        r'\b(selling motion|business acumen|true leader|value proposition|'
        r'key takeaway|stakeholder|playbook|deal review|pipeline|'
        r'account lead|solution architect|domain expert|customer success|'
        r'partner ecosystem|extended capability|value realization|'
        r'industry coe|orchestrator|value architect)\b',
        re.IGNORECASE,
    )

    def _check_lang_mix(self, num: int, shapes) -> List[NarrativeIssue]:
        issues = []
        for s in shapes:
            if not s.has_text_frame:
                continue
            text = s.text_frame.text
            # Only flag if the text is predominantly Chinese (>30% CJK)
            cjk = sum(1 for c in text if '\u4e00' <= c <= '\u9fff')
            if len(text) == 0 or cjk / max(len(text), 1) < 0.2:
                continue
            matches = self._EN_JARGON_RE.findall(text)
            if matches:
                issues.append(NarrativeIssue(
                    slide_num=num,
                    severity=Severity.INFO,
                    category="lang_mix",
                    message=f"中英混用术语: {', '.join(set(m.lower() for m in matches))}",
                    shape_name=getattr(s, "name", ""),
                    suggestion="仅供参考，英文专有名词可保留原文",
                ))
        return issues

    # ── Text density per box ──────────────────────────────────
    def _check_density(self, num: int, shapes) -> List[NarrativeIssue]:
        issues = []
        for s in shapes:
            if not s.has_text_frame or not s.text_frame.text.strip():
                continue
            if s.height is None or s.height <= 0:
                continue
            box_h_in = s.height / 914400
            text = s.text_frame.text.strip()
            char_count = len(text)

            # Find applicable limit
            limit = None
            for h_thresh in sorted(CHAR_DENSITY_LIMITS.keys()):
                if box_h_in >= h_thresh:
                    limit = CHAR_DENSITY_LIMITS[h_thresh]
            if limit and char_count > limit * 1.3:
                issues.append(NarrativeIssue(
                    slide_num=num,
                    severity=Severity.WARNING,
                    category="density",
                    message=f"文本密度过高: {char_count}字 / {box_h_in:.2f}\" 高度框 (建议≤{limit})",
                    shape_name=getattr(s, "name", ""),
                    suggestion="考虑精简文案或拆分内容",
                ))
        return issues

    # ── Title length ──────────────────────────────────────────
    def _check_title_length(self, num: int, shapes) -> List[NarrativeIssue]:
        issues = []
        for s in shapes:
            if not s.has_text_frame:
                continue
            for para in s.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size.pt >= 20:
                        text = para.text.strip()
                        if len(text) > ACTION_TITLE_MAX_CHARS:
                            issues.append(NarrativeIssue(
                                slide_num=num,
                                severity=Severity.WARNING,
                                category="title_long",
                                message=f"标题过长: {len(text)}字 (建议≤{ACTION_TITLE_MAX_CHARS})",
                                shape_name=getattr(s, "name", ""),
                                suggestion="缩短标题，把补充信息移到副标题或正文",
                            ))
                        break  # one check per shape
                break
        return issues


# ═══════════════════════════════════════════════════════════════
# SlideReviewer — combined Narrative + Layout QA
# ═══════════════════════════════════════════════════════════════

@dataclass
class CombinedReport:
    filepath: str
    layout_report: QAReport
    narrative_report: NarrativeReport

    @property
    def total_errors(self) -> int:
        return len(self.layout_report.errors) + len(self.narrative_report.errors)

    @property
    def passed(self) -> bool:
        return len(self.layout_report.errors) == 0

    def print_summary(self):
        print(f"\n{'='*70}")
        print(f"  Combined Review: {os.path.basename(self.filepath)}")
        print(f"  Layout score: {self.layout_report.overall_score}/100")
        print(f"  Layout errors: {len(self.layout_report.errors)}  "
              f"warnings: {len(self.layout_report.warnings)}")
        print(f"  Narrative errors: {len(self.narrative_report.errors)}  "
              f"warnings: {len(self.narrative_report.warnings)}")
        print(f"  Gate: {'✅ PASS' if self.passed else '❌ FAIL'}")
        print(f"{'='*70}")
        self.layout_report.print_summary()
        self.narrative_report.print_summary()


class SlideReviewer:
    """Run dual QA (Layout + Narrative) on a generated pptx."""

    def __init__(self, filepath: str):
        self.filepath = filepath

    def run(self) -> CombinedReport:
        layout_report = PptQA(self.filepath).run()
        narrative_report = NarrativeReviewer(self.filepath).run()
        return CombinedReport(
            filepath=self.filepath,
            layout_report=layout_report,
            narrative_report=narrative_report,
        )


# ═══════════════════════════════════════════════════════════════
# AutoFixPipeline — fix text overflow by priority chain
#
# Priority (no layout changes, per user constraint):
#   1. 去冗余 — remove redundant phrasing, weak info
#   2. 统一语言 — replace unnecessary English jargon with Chinese
#   3. 压缩句式 — shorten sentences
#   4. 重构层级 — restructure bullets (within same box)
#   5. 字号微调 — shrink font within guard rails
#
# Works directly on the .pptx shapes (post-generation fix).
# ═══════════════════════════════════════════════════════════════

# Mapping of English jargon → Chinese replacement
# DISABLED: English domain terms (selling motion, business acumen, playbook,
# deal review, etc.) are intentional professional vocabulary — do NOT translate.
_LANG_REPLACEMENTS = {}

# Redundancy patterns: (regex, replacement_or_empty)
_REDUNDANCY_PATTERNS = [
    # Remove weak hedging
    (re.compile(r'(?:从某种意义上说|在一定程度上|可以说是)'), ''),
    # Remove trailing "等" after a clear list
    (re.compile(r'([^、]+、[^、]+)等(?=[。，])'), r'\1'),
    # Collapse "进行XX" → "XX"
    (re.compile(r'进行([\u4e00-\u9fff]{2,4})'), r'\1'),
    # Remove "的话"
    (re.compile(r'的话(?=[，。；])'), ''),
]

# Sentence compression: simplify common verbose patterns
_COMPRESSION_PATTERNS = [
    # "因为A所以B" → "A → B"  (for bullet-style text)
    (re.compile(r'因为(.{4,20})(?:，|,)所以(.{4,20})'), r'\1 → \2'),
    # "不仅...而且..." → combine
    (re.compile(r'不仅(.{3,15})(?:，|,)而且(.{3,15})'), r'\1，且\2'),
    # "通过...来实现..." simplify
    (re.compile(r'通过(.{3,15})来实现(.{3,15})'), r'\1实现\2'),
]


class AutoFixPipeline:
    """Fix text overflow issues in a generated .pptx, in-place.

    Only modifies text content and font sizes.
    Does NOT change layout, shape positions, or box dimensions.
    """

    def __init__(self, filepath: str):
        self.filepath = filepath

    def run(self, max_rounds: int = 3, verbose: bool = True) -> CombinedReport:
        """Iterate: fix → re-check → fix ... up to max_rounds.
        After overflow fixes converge, harmonize peer font groups.
        """
        # Phase 1: Fix text overflow errors iteratively
        for round_num in range(1, max_rounds + 1):
            if verbose:
                print(f"\n{'━'*60}")
                print(f"  Auto-fix round {round_num}/{max_rounds}")
                print(f"{'━'*60}")

            # Run QA
            report = SlideReviewer(self.filepath).run()

            layout_errors = report.layout_report.errors
            text_overflow_errors = [
                e for e in layout_errors if e.category == "text_overflow"
            ]

            if not text_overflow_errors:
                if verbose:
                    print("  ✅ No text overflow errors.")
                break

            if verbose:
                print(f"  Found {len(text_overflow_errors)} text overflow error(s).")

            # Apply fixes
            fixes_applied = self._apply_fixes(text_overflow_errors, verbose)

            if not fixes_applied:
                if verbose:
                    print("  ⚠️  No more fixes applicable. Stopping iteration.")
                break

        # Phase 2: Harmonize peer font groups
        if verbose:
            print(f"\n{'━'*60}")
            print(f"  Peer font harmonization")
            print(f"{'━'*60}")
        harmonized = self._harmonize_peer_fonts(verbose)

        # Phase 3: Final check
        report = SlideReviewer(self.filepath).run()
        if verbose and report.passed:
            print("  ✅ Pipeline PASSED.")
        return report

    def _apply_fixes(self, overflow_errors: List[QAIssue], verbose: bool) -> int:
        """Apply the priority chain to fix overflows. Returns number of fixes applied."""
        prs = Presentation(self.filepath)
        fixes_count = 0

        for err in overflow_errors:
            slide_idx = err.slide_num - 1
            if slide_idx >= len(prs.slides):
                continue
            slide = prs.slides[slide_idx]
            shape = self._find_shape(slide, err.shape_name)
            if shape is None or not shape.has_text_frame:
                continue

            tf = shape.text_frame
            original_text = tf.text
            box_h = shape.height
            box_w = shape.width

            # Priority 1: Remove redundancy
            fixed = self._fix_redundancy(tf)
            if fixed and self._check_fits(tf, box_w, box_h):
                fixes_count += 1
                if verbose:
                    print(f"    S{err.slide_num} [{err.shape_name}] ✂️  去冗余 → 通过")
                continue

            # Priority 2: Unify language (replace English jargon)
            fixed = self._fix_language(tf)
            if fixed and self._check_fits(tf, box_w, box_h):
                fixes_count += 1
                if verbose:
                    print(f"    S{err.slide_num} [{err.shape_name}] 🔤 统一语言 → 通过")
                continue

            # Priority 3: Compress sentences
            fixed = self._fix_compress(tf)
            if fixed and self._check_fits(tf, box_w, box_h):
                fixes_count += 1
                if verbose:
                    print(f"    S{err.slide_num} [{err.shape_name}] 📐 压缩句式 → 通过")
                continue

            # Priority 4: Restructure (trim bullets, shorten)
            fixed = self._fix_restructure(tf)
            if fixed and self._check_fits(tf, box_w, box_h):
                fixes_count += 1
                if verbose:
                    print(f"    S{err.slide_num} [{err.shape_name}] 🔄 重构层级 → 通过")
                continue

            # Priority 5: Font size micro-adjust
            fixed = self._fix_font_size(tf, box_w, box_h)
            if fixed:
                fixes_count += 1
                if verbose:
                    print(f"    S{err.slide_num} [{err.shape_name}] 🔠 字号微调 → 通过")
                continue

            if verbose:
                print(f"    S{err.slide_num} [{err.shape_name}] ⚠️  所有策略未能解决溢出")

        if fixes_count > 0:
            prs.save(self.filepath)
            # Re-run full_cleanup to sanitize
            try:
                from .core import full_cleanup
                full_cleanup(self.filepath)
            except Exception:
                pass

        return fixes_count

    # ── Peer font harmonization ──────────────────────────────
    PEER_Y_TOLERANCE = Emu(18288)  # 0.02" — match qa.py

    def _harmonize_peer_fonts(self, verbose: bool) -> int:
        """After overflow fixes, ensure peer groups (same Y position)
        share the same font size.  Uses min(sizes) so no new overflow.
        Returns number of shapes adjusted.
        """
        prs = Presentation(self.filepath)
        total_adjusted = 0

        for slide_idx, slide in enumerate(prs.slides):
            # Collect text shapes with their resolved font info
            entries = []  # (top_emu, shape, para_font_size_pt, run_font_size_pt)
            for s in slide.shapes:
                if not s.has_text_frame or not s.text_frame.text.strip():
                    continue
                if s.top is None:
                    continue
                for para in s.text_frame.paragraphs:
                    if not para.text.strip():
                        continue
                    p_size = para.font.size.pt if para.font.size else None
                    r_size = None
                    for run in para.runs:
                        if run.font.size:
                            r_size = run.font.size.pt
                            break
                    effective = r_size or p_size
                    if effective is not None:
                        entries.append((s.top, s, p_size, r_size, effective))
                    break  # first paragraph only

            if len(entries) < 3:
                continue

            # Group by Y position
            entries.sort(key=lambda e: e[0])
            groups: list[list] = []
            for entry in entries:
                placed = False
                for g in groups:
                    if abs(entry[0] - g[0][0]) <= self.PEER_Y_TOLERANCE:
                        g.append(entry)
                        placed = True
                        break
                if not placed:
                    groups.append([entry])

            # Harmonize each peer group (≥3 members)
            for g in groups:
                if len(g) < 3:
                    continue
                sizes = [e[4] for e in g]
                if len(set(sizes)) <= 1:
                    continue  # already uniform

                target_size = min(sizes)  # use smallest to avoid new overflows
                for top_emu, shape, p_size, r_size, eff_size in g:
                    if eff_size == target_size:
                        continue
                    # Adjust paragraph-level font
                    for para in shape.text_frame.paragraphs:
                        if para.font.size and para.font.size.pt != target_size:
                            para.font.size = Pt(target_size)
                            total_adjusted += 1
                        for run in para.runs:
                            if run.font.size and run.font.size.pt != target_size:
                                run.font.size = Pt(target_size)
                        break  # first paragraph only

                if verbose:
                    names = [getattr(e[1], 'name', '?') for e in g]
                    texts = [e[1].text_frame.text.strip()[:20] for e in g]
                    print(f"    S{slide_idx+1} y≈{g[0][0]/914400:.2f}\" "
                          f"统一为 {target_size}pt "
                          f"({', '.join(f'{t}' for t in texts)})")

        if total_adjusted > 0:
            prs.save(self.filepath)
            try:
                from .core import full_cleanup
                full_cleanup(self.filepath)
            except Exception:
                pass
            if verbose:
                print(f"  已调整 {total_adjusted} 个 shape 的字号以保持同级一致")
        else:
            if verbose:
                print("  ✅ 所有同级 peer group 字号已一致")

        return total_adjusted

    # ── Shape finder ──────────────────────────────────────────
    def _find_shape(self, slide, shape_name: str):
        for s in slide.shapes:
            if getattr(s, "name", "") == shape_name:
                return s
        return None

    # ── Fit checker ───────────────────────────────────────────
    def _check_fits(self, tf, box_w, box_h) -> bool:
        est = _estimate_text_height(tf, box_w)
        return est <= box_h * TEXT_OVERFLOW_LINE_RATIO

    # ── Priority 1: Remove redundancy ─────────────────────────
    def _fix_redundancy(self, tf) -> bool:
        changed = False
        for para in tf.paragraphs:
            for run in para.runs:
                original = run.text
                text = original
                for pattern, replacement in _REDUNDANCY_PATTERNS:
                    text = pattern.sub(replacement, text)
                if text != original:
                    run.text = text
                    changed = True
        return changed

    # ── Priority 2: Unify language ────────────────────────────
    def _fix_language(self, tf) -> bool:
        changed = False
        for para in tf.paragraphs:
            for run in para.runs:
                text = run.text
                original = text
                for en, zh in _LANG_REPLACEMENTS.items():
                    text = re.sub(re.escape(en), zh, text, flags=re.IGNORECASE)
                if text != original:
                    run.text = text
                    changed = True
        return changed

    # ── Priority 3: Compress sentences ────────────────────────
    def _fix_compress(self, tf) -> bool:
        changed = False
        for para in tf.paragraphs:
            for run in para.runs:
                original = run.text
                text = original
                for pattern, replacement in _COMPRESSION_PATTERNS:
                    text = pattern.sub(replacement, text)
                # Also: strip trailing punctuation repetition
                text = re.sub(r'[。]{2,}', '。', text)
                text = re.sub(r'[，]{2,}', '，', text)
                if text != original:
                    run.text = text
                    changed = True
        return changed

    # ── Priority 4: Restructure (within same box) ─────────────
    def _fix_restructure(self, tf) -> bool:
        """Trim each paragraph to fit better:
        - Remove trailing explanatory clauses after semicolons
        - If a bullet has >2 sub-clauses, keep only the first two
        """
        changed = False
        for para in tf.paragraphs:
            for run in para.runs:
                text = run.text
                original = text

                # If text has semicolons with 3+ clauses, trim to first 2
                if '；' in text:
                    parts = text.split('；')
                    if len(parts) > 2:
                        text = '；'.join(parts[:2])
                        changed = True

                # If text has commas with many clauses (>3), trim
                if text.count('，') > 3:
                    parts = text.split('，')
                    if len(parts) > 4:
                        text = '，'.join(parts[:3]) + '。'
                        changed = True

                if text != original:
                    run.text = text
        return changed

    # ── Priority 5: Font size micro-adjust ────────────────────
    def _fix_font_size(self, tf, box_w, box_h) -> bool:
        """Shrink font by 1pt at a time, respecting guard rails.
        Handles both run-level and paragraph-level font sizes.
        """
        changed = False
        for attempt in range(4):  # max 4pt reduction
            if self._check_fits(tf, box_w, box_h):
                return changed

            shrunk_any = False
            for para in tf.paragraphs:
                # Check paragraph-level font size
                if para.font.size:
                    current_pt = para.font.size.pt
                    if current_pt >= 20:
                        floor = TITLE_MIN_PT
                    elif current_pt >= 13:
                        floor = BODY_MIN_PT
                    else:
                        floor = SMALL_MIN_PT
                    if current_pt > floor:
                        para.font.size = Pt(current_pt - 1)
                        shrunk_any = True
                        changed = True

                # Also check run-level
                for run in para.runs:
                    if run.font.size:
                        current_pt = run.font.size.pt
                        if current_pt >= 20:
                            floor = TITLE_MIN_PT
                        elif current_pt >= 13:
                            floor = BODY_MIN_PT
                        else:
                            floor = SMALL_MIN_PT
                        if current_pt > floor:
                            run.font.size = Pt(current_pt - 1)
                            shrunk_any = True
                            changed = True

            if not shrunk_any:
                break  # all fonts already at floor

        return changed


# ═══════════════════════════════════════════════════════════════
# CLI entry
# ═══════════════════════════════════════════════════════════════

def review(filepath: str) -> CombinedReport:
    """Read-only dual review."""
    return SlideReviewer(filepath).run()


def autofix(filepath: str, max_rounds: int = 3) -> CombinedReport:
    """Review + auto-fix pipeline."""
    return AutoFixPipeline(filepath).run(max_rounds=max_rounds)


if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("Usage: python -m mck_ppt.review <path.pptx> [--fix]")
        sys.exit(1)
    path = sys.argv[1]
    if "--fix" in sys.argv:
        result = autofix(path)
    else:
        result = review(path)
    result.print_summary()
    sys.exit(0 if result.passed else 1)