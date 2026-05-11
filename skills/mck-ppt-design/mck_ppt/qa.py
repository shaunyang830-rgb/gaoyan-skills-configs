# Copyright 2024-2026 Kaku Li (https://github.com/likaku)
# Licensed under the Apache License, Version 2.0 — see LICENSE and NOTICE.
# Part of "Mck-ppt-design-skill" (McKinsey PPT Design Framework).
# NOTICE: This file must be retained in all copies or substantial portions.
#
"""
PPT QA Engine — Automated visual quality analysis for mck_ppt outputs.

Analyzes generated PPTX files for common layout defects:
  1. Body overflow   — shapes extending beyond slide boundaries
  2. Text overflow   — text that doesn't fit in its textbox
  3. Dead whitespace — large unused regions in content area
  4. Shape overlap   — overlapping shapes that shouldn't overlap
  5. Font violations — wrong font family or size
  6. Guard rail violations — connector usage, missing cleanup, etc.

Usage:
    from mck_ppt.qa import PptQA
    report = PptQA("output.pptx").run()
    report.print_summary()
    report.to_json("qa_report.json")
"""

from __future__ import annotations

import json
import math
import os
import re
from dataclasses import dataclass, field, asdict
from typing import List, Optional, Tuple
from collections import defaultdict

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

# ── Slide constants (must match constants.py) ──────────────────────────
SW = Inches(13.333)   # slide width
SH = Inches(7.5)      # slide height
LM = Inches(0.8)      # left margin
RM = Inches(0.8)      # right margin
CW = Inches(11.733)   # content width

TITLE_TOP      = Inches(0.15)
TITLE_H        = Inches(0.9)
TITLE_LINE_Y   = Inches(1.05)
CONTENT_TOP    = Inches(1.3)
SOURCE_Y       = Inches(7.05)
BOTTOM_BAR_Y   = Inches(6.2)
BOTTOM_BAR_H   = Inches(0.65)

# Safe zone: all content shapes should stay within these bounds
SAFE_LEFT   = 0                 # shapes can touch left edge (cover etc.)
SAFE_RIGHT  = SW                # right boundary
SAFE_TOP    = 0                 # shapes can touch top (cover)
SAFE_BOTTOM = SH                # bottom boundary

# Content area (between title bar and source line)
CONTENT_AREA_TOP    = CONTENT_TOP
CONTENT_AREA_BOTTOM = SOURCE_Y
CONTENT_AREA_LEFT   = LM
CONTENT_AREA_RIGHT  = LM + CW

# Thresholds
OVERFLOW_TOLERANCE       = Emu(18288)   # 0.02" — ignore sub-pixel overflow
WHITESPACE_THRESHOLD     = 0.55         # warn if > 55% of content area is empty
TEXT_OVERFLOW_LINE_RATIO  = 1.15        # text height > 1.15x box height = overflow
MIN_FONT_SIZE            = Pt(8)        # anything smaller is suspicious
MAX_FONT_SIZE            = Pt(48)       # anything larger (except cover) is suspicious
OVERLAP_TOLERANCE        = Emu(18288)   # 0.02" — ignore trivial overlap


# ── Severity ───────────────────────────────────────────────────────────
class Severity:
    ERROR   = "ERROR"     # will cause visual defect or file corruption
    WARNING = "WARNING"   # likely visual issue
    INFO    = "INFO"      # minor observation


# ── Issue data class ───────────────────────────────────────────────────
@dataclass
class QAIssue:
    slide_num: int
    severity: str
    category: str
    message: str
    shape_name: str = ""
    details: dict = field(default_factory=dict)

    def to_dict(self):
        d = asdict(self)
        # Convert EMU values to inches for readability
        for k, v in d.get("details", {}).items():
            if isinstance(v, (int, float)) and abs(v) > 10000:
                d["details"][k] = f"{v / 914400:.3f} in"
        return d


# ── QA Report ──────────────────────────────────────────────────────────
@dataclass
class QAReport:
    filepath: str
    total_slides: int = 0
    issues: List[QAIssue] = field(default_factory=list)
    slide_scores: dict = field(default_factory=dict)  # slide_num → score (0-100)

    @property
    def errors(self) -> List[QAIssue]:
        return [i for i in self.issues if i.severity == Severity.ERROR]

    @property
    def warnings(self) -> List[QAIssue]:
        return [i for i in self.issues if i.severity == Severity.WARNING]

    @property
    def infos(self) -> List[QAIssue]:
        return [i for i in self.issues if i.severity == Severity.INFO]

    @property
    def overall_score(self) -> int:
        """0-100 score. 100 = perfect."""
        if not self.slide_scores:
            return 100
        return round(sum(self.slide_scores.values()) / len(self.slide_scores))

    @property
    def passed(self) -> bool:
        return len(self.errors) == 0

    def print_summary(self):
        """Print human-readable summary to stdout."""
        print(f"\n{'='*70}")
        print(f"  PPT QA Report: {os.path.basename(self.filepath)}")
        print(f"  Slides: {self.total_slides}  |  Score: {self.overall_score}/100")
        print(f"  Errors: {len(self.errors)}  |  Warnings: {len(self.warnings)}  |  Info: {len(self.infos)}")
        print(f"{'='*70}")

        if not self.issues:
            print("  ✅ All checks passed — no issues found.")
            print()
            return

        # Group by slide
        by_slide = defaultdict(list)
        for issue in self.issues:
            by_slide[issue.slide_num].append(issue)

        for slide_num in sorted(by_slide.keys()):
            score = self.slide_scores.get(slide_num, 100)
            print(f"\n  Slide {slide_num}  (score: {score}/100)")
            print(f"  {'─'*50}")
            for issue in by_slide[slide_num]:
                icon = {"ERROR": "❌", "WARNING": "⚠️ ", "INFO": "ℹ️ "}[issue.severity]
                print(f"    {icon} [{issue.category}] {issue.message}")
                if issue.shape_name:
                    print(f"       Shape: {issue.shape_name}")
                if issue.details:
                    for k, v in issue.details.items():
                        if isinstance(v, (int, float)) and abs(v) > 10000:
                            print(f"       {k}: {v/914400:.3f}\"")
                        else:
                            print(f"       {k}: {v}")

        print(f"\n{'='*70}\n")

    def to_json(self, outpath: str = None) -> str:
        """Serialize report to JSON. If outpath given, also write to file."""
        data = {
            "filepath": self.filepath,
            "total_slides": self.total_slides,
            "overall_score": self.overall_score,
            "passed": self.passed,
            "summary": {
                "errors": len(self.errors),
                "warnings": len(self.warnings),
                "info": len(self.infos),
            },
            "slide_scores": {str(k): v for k, v in self.slide_scores.items()},
            "issues": [i.to_dict() for i in self.issues],
        }
        text = json.dumps(data, indent=2, ensure_ascii=False)
        if outpath:
            os.makedirs(os.path.dirname(outpath) or ".", exist_ok=True)
            with open(outpath, "w", encoding="utf-8") as f:
                f.write(text)
        return text


# ── Shape bounding box helper ─────────────────────────────────────────
def _bbox(shape) -> Tuple[int, int, int, int]:
    """Return (left, top, right, bottom) in EMU."""
    return (shape.left, shape.top, shape.left + shape.width, shape.top + shape.height)


def _boxes_overlap(a, b, tol=OVERLAP_TOLERANCE) -> bool:
    """Check if two bounding boxes overlap beyond tolerance."""
    al, at, ar, ab = a
    bl, bt, br, bb = b
    overlap_x = max(0, min(ar, br) - max(al, bl) - tol)
    overlap_y = max(0, min(ab, bb) - max(at, bt) - tol)
    return overlap_x > 0 and overlap_y > 0


def _overlap_area(a, b) -> int:
    """Overlap area in EMU²."""
    al, at, ar, ab = a
    bl, bt, br, bb = b
    dx = max(0, min(ar, br) - max(al, bl))
    dy = max(0, min(ab, bb) - max(at, bt))
    return dx * dy


# ── Text overflow estimation ──────────────────────────────────────────
def _estimate_text_height(tf, box_width_emu: int) -> int:
    """
    Estimate required height for text content in a textbox.
    This is approximate — exact rendering depends on the OS font engine.
    """
    total_height = 0
    for para in tf.paragraphs:
        # Get font size for this paragraph — check run level first,
        # then paragraph level (p.font.size), then default 14pt
        font_size_pt = 14  # default fallback
        found = False
        for run in para.runs:
            if run.font.size:
                font_size_pt = run.font.size.pt
                found = True
                break
        if not found and para.font.size:
            font_size_pt = para.font.size.pt

        # Estimate line height (font size * line spacing factor)
        line_height_emu = int(font_size_pt * 1.4 * 12700)  # pt → EMU with ~1.4x spacing

        # Estimate text width and wrap count
        text = para.text or ""
        if not text:
            total_height += line_height_emu  # empty paragraph still takes space
            continue

        # Rough char width: ~0.6 * font_size for proportional fonts
        avg_char_width_emu = int(font_size_pt * 0.55 * 12700)
        if avg_char_width_emu <= 0:
            avg_char_width_emu = 1

        # Account for textbox internal padding (~0.1")
        usable_width = max(box_width_emu - Emu(91440), Emu(914400))

        # CJK chars are wider
        cjk_count = sum(1 for c in text if '\u4e00' <= c <= '\u9fff' or '\u3000' <= c <= '\u303f')
        latin_count = len(text) - cjk_count
        estimated_text_width = int((latin_count * 0.55 + cjk_count * 1.0) * font_size_pt * 12700)

        num_lines = max(1, math.ceil(estimated_text_width / usable_width))
        total_height += num_lines * line_height_emu

    return total_height


# ── PptQA Main Class ──────────────────────────────────────────────────
class PptQA:
    """
    Automated PPT quality analysis engine.

    Usage:
        report = PptQA("my_deck.pptx").run()
        report.print_summary()
    """

    def __init__(self, filepath: str):
        self.filepath = filepath
        self.prs = Presentation(filepath)
        self.issues: List[QAIssue] = []
        self.slide_scores: dict = {}

    def run(self) -> QAReport:
        """Run all QA checks and return report."""
        total = len(self.prs.slides)

        for idx, slide in enumerate(self.prs.slides, 1):
            slide_issues_before = len(self.issues)
            self._check_slide(idx, slide)
            slide_issues_after = len(self.issues)

            # Calculate slide score
            new_issues = self.issues[slide_issues_before:slide_issues_after]
            self.slide_scores[idx] = self._calc_slide_score(new_issues)

        # Global checks
        self._check_global()

        report = QAReport(
            filepath=self.filepath,
            total_slides=total,
            issues=self.issues,
            slide_scores=self.slide_scores,
        )
        return report

    def _check_slide(self, num: int, slide):
        """Run all per-slide checks."""
        shapes = list(slide.shapes)
        self._check_body_overflow(num, shapes)
        self._check_text_overflow(num, shapes)
        self._check_text_line_collision(num, shapes)
        self._check_whitespace(num, shapes)
        self._check_shape_overlap(num, shapes)
        self._check_fonts(num, shapes)
        self._check_peer_font_consistency(num, shapes)
        self._check_chart_legend_overflow(num, shapes)
        self._check_connectors(num, slide)

    # ── Check 1: Body Overflow ────────────────────────────────────────
    def _check_body_overflow(self, num: int, shapes):
        """Check if any shape extends beyond slide boundaries."""
        for shape in shapes:
            if not hasattr(shape, "left") or shape.left is None:
                continue
            left, top, right, bottom = _bbox(shape)
            name = getattr(shape, "name", "unknown")

            # Right overflow
            if right > SAFE_RIGHT + OVERFLOW_TOLERANCE:
                overflow = right - SAFE_RIGHT
                self.issues.append(QAIssue(
                    slide_num=num,
                    severity=Severity.ERROR,
                    category="body_overflow",
                    message=f"Shape overflows RIGHT edge by {overflow/914400:.3f}\"",
                    shape_name=name,
                    details={"right_edge": right, "slide_width": SW, "overflow_emu": overflow},
                ))

            # Bottom overflow
            if bottom > SAFE_BOTTOM + OVERFLOW_TOLERANCE:
                overflow = bottom - SAFE_BOTTOM
                self.issues.append(QAIssue(
                    slide_num=num,
                    severity=Severity.ERROR,
                    category="body_overflow",
                    message=f"Shape overflows BOTTOM edge by {overflow/914400:.3f}\"",
                    shape_name=name,
                    details={"bottom_edge": bottom, "slide_height": SH, "overflow_emu": overflow},
                ))

            # Left overflow (negative position)
            if left < SAFE_LEFT - OVERFLOW_TOLERANCE:
                overflow = SAFE_LEFT - left
                self.issues.append(QAIssue(
                    slide_num=num,
                    severity=Severity.WARNING,
                    category="body_overflow",
                    message=f"Shape extends LEFT beyond slide by {overflow/914400:.3f}\"",
                    shape_name=name,
                    details={"left_pos": left, "overflow_emu": overflow},
                ))

            # Top overflow (negative position)
            if top < SAFE_TOP - OVERFLOW_TOLERANCE:
                overflow = SAFE_TOP - top
                self.issues.append(QAIssue(
                    slide_num=num,
                    severity=Severity.WARNING,
                    category="body_overflow",
                    message=f"Shape extends TOP beyond slide by {overflow/914400:.3f}\"",
                    shape_name=name,
                    details={"top_pos": top, "overflow_emu": overflow},
                ))

    # ── Check 2: Text Overflow ────────────────────────────────────────
    def _check_text_overflow(self, num: int, shapes):
        """Check if text content likely overflows its textbox."""
        for shape in shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            text = tf.text.strip()
            if not text:
                continue

            name = getattr(shape, "name", "unknown")
            box_h = shape.height
            box_w = shape.width

            if box_h is None or box_h <= 0 or box_w is None or box_w <= 0:
                continue

            # Estimate required height
            est_height = _estimate_text_height(tf, box_w)

            if est_height > box_h * TEXT_OVERFLOW_LINE_RATIO:
                overflow_pct = (est_height - box_h) / box_h * 100
                sev = Severity.ERROR if overflow_pct > 30 else Severity.WARNING
                self.issues.append(QAIssue(
                    slide_num=num,
                    severity=sev,
                    category="text_overflow",
                    message=f"Text likely overflows box by ~{overflow_pct:.0f}%",
                    shape_name=name,
                    details={
                        "box_height": box_h,
                        "estimated_text_height": est_height,
                        "overflow_pct": round(overflow_pct, 1),
                        "text_preview": text[:80] + ("..." if len(text) > 80 else ""),
                    },
                ))

    # ── Check 3: Dead Whitespace ──────────────────────────────────────
    def _check_whitespace(self, num: int, shapes):
        """
        Check for excessive unused space in the content area.
        Uses a grid-based coverage approach.
        """
        # Filter to shapes in content area (skip title bar, source, page number)
        content_shapes = []
        for s in shapes:
            if not hasattr(s, "left") or s.left is None:
                continue
            l, t, r, b = _bbox(s)
            # Shape must overlap with content area
            if r > CONTENT_AREA_LEFT and l < CONTENT_AREA_RIGHT \
               and b > CONTENT_AREA_TOP and t < CONTENT_AREA_BOTTOM:
                content_shapes.append(s)

        if not content_shapes:
            # Slide has no content shapes — likely a cover or divider, skip
            return

        # Calculate coverage using bounding box union
        content_area_w = CONTENT_AREA_RIGHT - CONTENT_AREA_LEFT
        content_area_h = CONTENT_AREA_BOTTOM - CONTENT_AREA_TOP
        content_area_total = content_area_w * content_area_h

        if content_area_total <= 0:
            return

        # Grid-based coverage (20x20 grid for performance)
        GRID = 20
        cell_w = content_area_w // GRID
        cell_h = content_area_h // GRID
        covered = [[False] * GRID for _ in range(GRID)]

        for s in content_shapes:
            l, t, r, b = _bbox(s)
            # Clip to content area
            l = max(l, CONTENT_AREA_LEFT)
            t = max(t, CONTENT_AREA_TOP)
            r = min(r, CONTENT_AREA_RIGHT)
            b = min(b, CONTENT_AREA_BOTTOM)

            if l >= r or t >= b:
                continue

            col_start = max(0, int((l - CONTENT_AREA_LEFT) / cell_w))
            col_end = min(GRID, int(math.ceil((r - CONTENT_AREA_LEFT) / cell_w)))
            row_start = max(0, int((t - CONTENT_AREA_TOP) / cell_h))
            row_end = min(GRID, int(math.ceil((b - CONTENT_AREA_TOP) / cell_h)))

            for row in range(row_start, row_end):
                for col in range(col_start, col_end):
                    covered[row][col] = True

        covered_cells = sum(sum(row) for row in covered)
        total_cells = GRID * GRID
        coverage_ratio = covered_cells / total_cells
        empty_ratio = 1.0 - coverage_ratio

        if empty_ratio > WHITESPACE_THRESHOLD:
            # Identify where the dead space is
            dead_zones = self._identify_dead_zones(covered, GRID)
            self.issues.append(QAIssue(
                slide_num=num,
                severity=Severity.WARNING,
                category="dead_whitespace",
                message=f"Content area is {empty_ratio*100:.0f}% empty (threshold: {WHITESPACE_THRESHOLD*100:.0f}%)",
                details={
                    "coverage_pct": round(coverage_ratio * 100, 1),
                    "empty_pct": round(empty_ratio * 100, 1),
                    "dead_zones": dead_zones,
                },
            ))

    def _identify_dead_zones(self, covered, grid_size) -> List[str]:
        """Identify named regions of empty space."""
        zones = []
        # Check bottom third
        bottom_empty = sum(1 for col in range(grid_size) for row in range(grid_size*2//3, grid_size) if not covered[row][col])
        bottom_total = grid_size * (grid_size - grid_size*2//3)
        if bottom_total > 0 and bottom_empty / bottom_total > 0.8:
            zones.append("bottom_third")

        # Check right third
        right_empty = sum(1 for row in range(grid_size) for col in range(grid_size*2//3, grid_size) if not covered[row][col])
        right_total = grid_size * (grid_size - grid_size*2//3)
        if right_total > 0 and right_empty / right_total > 0.8:
            zones.append("right_third")

        # Check left third
        left_empty = sum(1 for row in range(grid_size) for col in range(0, grid_size//3) if not covered[row][col])
        left_total = grid_size * (grid_size // 3)
        if left_total > 0 and left_empty / left_total > 0.8:
            zones.append("left_third")

        # Check center
        mid = grid_size // 2
        r = grid_size // 4
        center_empty = sum(1 for row in range(mid-r, mid+r) for col in range(mid-r, mid+r) if not covered[row][col])
        center_total = (2*r) * (2*r)
        if center_total > 0 and center_empty / center_total > 0.8:
            zones.append("center")

        return zones if zones else ["scattered"]

    # ── Check 4: Shape Overlap ────────────────────────────────────────
    def _check_shape_overlap(self, num: int, shapes):
        """
        Check for suspicious shape overlaps.
        Some overlaps are intentional (text on top of rect background),
        so we only flag text-on-text overlaps.
        """
        text_shapes = []
        for s in shapes:
            if s.has_text_frame and s.text_frame.text.strip():
                if hasattr(s, "left") and s.left is not None:
                    text_shapes.append(s)

        for i in range(len(text_shapes)):
            for j in range(i + 1, len(text_shapes)):
                a = text_shapes[i]
                b = text_shapes[j]
                ba = _bbox(a)
                bb = _bbox(b)
                if _boxes_overlap(ba, bb):
                    area = _overlap_area(ba, bb)
                    # Only flag if overlap area is significant (> 0.5 sq inch)
                    min_area = min(
                        (ba[2]-ba[0]) * (ba[3]-ba[1]),
                        (bb[2]-bb[0]) * (bb[3]-bb[1])
                    )
                    if min_area > 0 and area / min_area > 0.15:
                        overlap_pct = area / min_area * 100
                        self.issues.append(QAIssue(
                            slide_num=num,
                            severity=Severity.WARNING,
                            category="shape_overlap",
                            message=f"Two text shapes overlap by {overlap_pct:.0f}%",
                            shape_name=f"{a.name} ↔ {b.name}",
                            details={
                                "shape_a": a.name,
                                "shape_b": b.name,
                                "overlap_pct": round(overlap_pct, 1),
                                "text_a_preview": a.text_frame.text[:40],
                                "text_b_preview": b.text_frame.text[:40],
                            },
                        ))

    # ── Check 5: Font Violations ──────────────────────────────────────
    def _check_fonts(self, num: int, shapes):
        """Check for font size anomalies."""
        for shape in shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.font.size:
                        continue
                    size = run.font.size
                    text = run.text.strip()
                    if not text:
                        continue

                    name = getattr(shape, "name", "unknown")

                    if size < MIN_FONT_SIZE:
                        self.issues.append(QAIssue(
                            slide_num=num,
                            severity=Severity.WARNING,
                            category="font_issue",
                            message=f"Font size {size.pt}pt is below minimum ({MIN_FONT_SIZE.pt}pt)",
                            shape_name=name,
                            details={"font_size_pt": size.pt, "text_preview": text[:40]},
                        ))

    # ── Check 6a: Text-Line Collision ────────────────────────────────
    TEXT_LINE_GAP_MIN = Emu(27432)  # 0.03" — minimum gap between text bottom and hline

    def _check_text_line_collision(self, num: int, shapes):
        """Check if text content visually collides with separator lines.

        Separator lines are thin rectangles (height ≤ 3pt).  If a text
        shape's estimated text bottom is within TEXT_LINE_GAP_MIN of
        a line's top edge, the text appears to "touch" or "overlap" the
        line — a common defect in dense layouts.
        """
        text_shapes = []
        line_shapes = []
        for s in shapes:
            if not hasattr(s, 'left') or s.left is None:
                continue
            if s.has_text_frame and s.text_frame.text.strip():
                text_shapes.append(s)
            # Identify thin rectangles (separator lines): height ≤ 3pt (38100 EMU)
            elif hasattr(s, 'height') and s.height is not None and s.height <= Emu(38100):
                if s.width is not None and s.width > Emu(914400):  # wider than 1"
                    line_shapes.append(s)

        for ts in text_shapes:
            if ts.height is None or ts.width is None:
                continue
            # Skip decorative elements (circle badges, single-char labels)
            if len(ts.text_frame.text.strip()) <= 2:
                continue
            # Skip very small shapes (badges, icons)
            if ts.height < Emu(228600) and ts.width < Emu(228600):  # < 0.25"
                continue
            # Skip vertically-centered shapes — text won't overflow downward
            try:
                bodyPr = ts.text_frame._txBody.find(qn('a:bodyPr'))
                if bodyPr is not None and bodyPr.get('anchor') in ('ctr', 'b'):
                    continue
            except Exception:
                pass
            est_text_h = _estimate_text_height(ts.text_frame, ts.width)
            text_bottom = ts.top + min(est_text_h, ts.height + Emu(45720))  # allow slight internal overflow

            for ls in line_shapes:
                line_top = ls.top
                if line_top is None or ls.left is None:
                    continue
                # Only check lines BELOW the text box (within reasonable range)
                if line_top < ts.top:
                    continue
                # Must horizontally overlap — skip lines in different columns
                ts_left = ts.left
                ts_right = ts.left + ts.width
                ls_left = ls.left
                ls_right = ls.left + ls.width
                h_overlap = min(ts_right, ls_right) - max(ts_left, ls_left)
                if h_overlap < Emu(457200):  # less than 0.5" horizontal overlap = different region
                    continue
                gap = line_top - text_bottom
                if gap < self.TEXT_LINE_GAP_MIN and gap > -Emu(182880):  # collision zone: -0.2" to +0.03"
                    self.issues.append(QAIssue(
                        slide_num=num,
                        severity=Severity.ERROR if gap < 0 else Severity.WARNING,
                        category="text_line_collision",
                        message=(f"Text {'overlaps' if gap < 0 else 'nearly touches'} "
                                 f"separator line (gap: {gap/914400:.3f}\")"),
                        shape_name=getattr(ts, 'name', ''),
                        details={
                            "text_bottom": text_bottom,
                            "line_top": line_top,
                            "gap_emu": gap,
                            "text_preview": ts.text_frame.text[:50],
                        },
                    ))

    # ── Check 6b: Peer Font Consistency ──────────────────────────────
    PEER_Y_TOLERANCE = Emu(18288)  # 0.02" — shapes within this Y range are peers

    def _check_peer_font_consistency(self, num: int, shapes):
        """Check that shapes at the same Y position (peer group) share
        consistent font size, family, and bold style.

        Peer groups are text shapes whose top positions match within a
        small tolerance.  Typical cases: card titles in meet_the_team,
        column headers in table layouts, numbered list items at the
        same row, etc.

        Any inconsistency within a peer group is an ERROR because it
        indicates the autofix (or generation) broke visual uniformity.
        """
        # Collect text shapes with resolved font info
        entries = []  # (top_emu, shape_name, font_size_pt, font_name, bold, text)
        for s in shapes:
            if not s.has_text_frame or not s.text_frame.text.strip():
                continue
            if s.top is None:
                continue
            # Resolve effective font from first non-empty paragraph
            for para in s.text_frame.paragraphs:
                if not para.text.strip():
                    continue
                size_pt = None
                fname = None
                bold = None
                # Run-level first, then paragraph-level
                for run in para.runs:
                    if run.font.size:
                        size_pt = run.font.size.pt
                    if run.font.name:
                        fname = run.font.name
                    if run.font.bold is not None:
                        bold = run.font.bold
                    if size_pt:
                        break
                if size_pt is None and para.font.size:
                    size_pt = para.font.size.pt
                if fname is None and para.font.name:
                    fname = para.font.name
                if bold is None and para.font.bold is not None:
                    bold = para.font.bold
                if size_pt is not None:
                    entries.append((s.top, getattr(s, 'name', ''),
                                    size_pt, fname, bold,
                                    s.text_frame.text.strip()[:30]))
                break  # only first paragraph

        if len(entries) < 2:
            return

        # Group by Y position (within tolerance)
        entries.sort(key=lambda e: e[0])
        groups: list[list] = []
        for entry in entries:
            placed = False
            for g in groups:
                if abs(entry[0] - g[0][0]) <= self.PEER_Y_TOLERANCE:
                    g.append(entry)
                    placed = True
                    break
            if not placed:
                groups.append([entry])

        # Check consistency within each group of 3+ peers
        for g in groups:
            if len(g) < 3:
                continue  # only flag groups with ≥3 items (card rows etc.)

            sizes = set(e[2] for e in g)
            fonts = set(e[3] for e in g if e[3])

            if len(sizes) > 1:
                names = [e[1] for e in g]
                size_list = [(e[1], e[2]) for e in g]
                self.issues.append(QAIssue(
                    slide_num=num,
                    severity=Severity.ERROR,
                    category="peer_font_inconsistency",
                    message=(f"Peer group at y≈{g[0][0]/914400:.2f}\" has "
                             f"inconsistent font sizes: "
                             f"{', '.join(f'{n}={s}pt' for n, s in size_list)}"),
                    shape_name=" / ".join(names),
                    details={
                        "sizes": {e[1]: e[2] for e in g},
                        "texts": {e[1]: e[5] for e in g},
                    },
                ))

            if len(fonts) > 1:
                names = [e[1] for e in g]
                font_list = [(e[1], e[3]) for e in g]
                self.issues.append(QAIssue(
                    slide_num=num,
                    severity=Severity.ERROR,
                    category="peer_font_inconsistency",
                    message=(f"Peer group at y≈{g[0][0]/914400:.2f}\" has "
                             f"inconsistent font families: "
                             f"{', '.join(f'{n}={f}' for n, f in font_list)}"),
                    shape_name=" / ".join(names),
                    details={
                        "fonts": {e[1]: e[3] for e in g},
                        "texts": {e[1]: e[5] for e in g},
                    },
                ))

    # ── Check: Chart Legend / Small Shape Cluster Overflow ───────────
    def _check_chart_legend_overflow(self, num: int, shapes):
        """Check if chart legends or small label clusters overflow the content area.

        Detects groups of small, horizontally-aligned shapes (typical of chart
        legends) where the rightmost element exceeds the content right boundary.
        Also catches any small text label whose right edge exceeds the slide width.
        """
        content_right = CONTENT_AREA_RIGHT

        # Collect small text shapes that look like legend items
        # (height <= 0.5", width <= 2.5", vertically aligned within 0.3")
        # Exclude page number shapes (pattern: "N/N" at bottom-right)
        import re
        page_num_re = re.compile(r'^\d+/\d+$')
        small_texts = []
        for shape in shapes:
            if not hasattr(shape, "left") or shape.left is None:
                continue
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            # Skip page numbers
            if page_num_re.match(text):
                continue
            w = (shape.width or 0) / 914400
            h = (shape.height or 0) / 914400
            top = (shape.top or 0) / 914400
            # Skip shapes in the footer area (source line, page number region)
            if top > 6.8:
                continue
            if h <= 0.5 and w <= 2.5 and len(text) <= 20:
                left, top_emu, right, bottom = _bbox(shape)
                small_texts.append((shape, left, top_emu, right, text))

        # Check each small text — if right edge exceeds content area, flag it
        for shape, left, top, right, text in small_texts:
            if right > content_right + OVERFLOW_TOLERANCE:
                overflow = (right - content_right) / 914400
                self.issues.append(QAIssue(
                    slide_num=num,
                    severity=Severity.ERROR,
                    category="chart_legend_overflow",
                    message=f"Legend/label '{text}' overflows content area RIGHT by {overflow:.2f}\"",
                    shape_name=getattr(shape, "name", ""),
                    details={"right_edge_in": right / 914400, "content_right_in": content_right / 914400},
                ))

    # ── Check 7: Connector Usage (Guard Rail #1) ─────────────────────
    def _check_connectors(self, num: int, slide):
        """Check for connector shapes that can cause file corruption."""
        xml = slide._element.xml
        if "<p:cxnSp" in xml:
            # Count connectors
            count = xml.count("<p:cxnSp")
            self.issues.append(QAIssue(
                slide_num=num,
                severity=Severity.ERROR,
                category="guard_rail",
                message=f"Slide contains {count} connector(s) — risk of file corruption (Guard Rail #1)",
                details={"connector_count": count},
            ))

    # ── Check 7: p:style remnants ─────────────────────────────────────
    def _check_global(self):
        """Global checks across the entire file."""
        # Check for p:style remnants in XML
        for idx, slide in enumerate(self.prs.slides, 1):
            xml = slide._element.xml
            pstyle_count = xml.count("<p:style")
            if pstyle_count > 0:
                self.issues.append(QAIssue(
                    slide_num=idx,
                    severity=Severity.WARNING,
                    category="guard_rail",
                    message=f"Found {pstyle_count} <p:style> remnant(s) — run full_cleanup()",
                    details={"pstyle_count": pstyle_count},
                ))

    # ── Slide score calculation ───────────────────────────────────────
    def _calc_slide_score(self, issues: List[QAIssue]) -> int:
        """Calculate quality score for a slide (0-100)."""
        score = 100
        for issue in issues:
            if issue.severity == Severity.ERROR:
                if issue.category == "body_overflow":
                    score -= 25
                elif issue.category == "text_overflow":
                    score -= 20
                elif issue.category == "guard_rail":
                    score -= 30
                else:
                    score -= 15
            elif issue.severity == Severity.WARNING:
                if issue.category == "dead_whitespace":
                    score -= 10
                elif issue.category == "text_overflow":
                    score -= 8
                elif issue.category == "shape_overlap":
                    score -= 10
                else:
                    score -= 5
            else:
                score -= 1
        return max(0, score)


# ── CLI entry point ───────────────────────────────────────────────────
def analyze(filepath: str, json_out: str = None, verbose: bool = True) -> QAReport:
    """
    Convenience function: analyze a PPTX and optionally save JSON report.

    Args:
        filepath: Path to .pptx file
        json_out: Optional path for JSON report output
        verbose: Print summary to stdout

    Returns:
        QAReport
    """
    qa = PptQA(filepath)
    report = qa.run()
    if verbose:
        report.print_summary()
    if json_out:
        report.to_json(json_out)
    return report


if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("Usage: python -m mck_ppt.qa <path.pptx> [--json report.json]")
        sys.exit(1)

    pptx_path = sys.argv[1]
    json_path = None
    if "--json" in sys.argv:
        idx = sys.argv.index("--json")
        if idx + 1 < len(sys.argv):
            json_path = sys.argv[idx + 1]

    report = analyze(pptx_path, json_out=json_path)
    sys.exit(0 if report.passed else 1)