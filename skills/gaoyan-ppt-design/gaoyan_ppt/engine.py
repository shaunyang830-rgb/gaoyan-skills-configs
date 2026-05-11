# GaoyanEngine — 高岩科技PPT设计引擎
# 基于高岩报告template--All-2026.pptx (255页) + PPT美化培训讲义 (42页)

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from .constants import *
from .core import (
    add_rect, add_flat_rect, add_oval, add_hline, add_vline,
    add_text, add_source, add_brand_footer, add_page_number,
    add_tag_label, add_sub_conclusion, add_action_title,
    add_image_placeholder, add_block_arc, full_cleanup, set_ea_font,
)


class GaoyanEngine:
    """高岩科技PPT生成引擎 — 高岩品牌视觉体系。

    特征：圆角矩形、微软雅黑、Century Gothic、高岩绿标签条、子结论条。
    """

    def __init__(self, total_slides=10):
        self.prs = Presentation()
        self.prs.slide_width = SW
        self.prs.slide_height = SH
        self.total_slides = total_slides
        self._page = 0

    def _new_slide(self):
        """Create a new blank slide and increment page counter."""
        layout = self.prs.slide_layouts[6]  # blank layout
        slide = self.prs.slides.add_slide(layout)
        self._page += 1
        return slide

    def _add_detail_chrome(self, slide, title, tag='', sub_conclusion='', source=''):
        """Add standard detail page chrome: tag + title + sub-conclusion + source + footer."""
        if tag:
            add_tag_label(slide, tag)
        add_action_title(slide, title)
        if sub_conclusion:
            add_sub_conclusion(slide, sub_conclusion)
        if source:
            add_source(slide, source)
        add_brand_footer(slide)
        add_page_number(slide, self._page, self.total_slides)

    # ══════════════════════════════════════════════
    # A. 结构导航
    # ══════════════════════════════════════════════

    def cover(self, title, subtitle='', author='', date='', project=''):
        """封面页 — 左侧信息 + 右侧图片区 + 左侧绿色渐变条。

        基于模板 Slide 10: 高岩科技公司介绍封面。
        """
        s = self._new_slide()

        # 左侧绿色渐变条
        add_flat_rect(s, GRAD_BAR_LEFT, GRAD_BAR_TOP, GRAD_BAR_W, GRAD_BAR_H, GAOYAN_GREEN)

        # 右侧图片占位区（约占右半）
        add_image_placeholder(s, Inches(5.24), Inches(0), Inches(8.10), SH, '封面图片')

        # Logo 占位
        add_rect(s, Inches(0.25), Inches(0.12), Inches(1.11), Inches(0.74),
                 BG_GRAY, corner_radius=Inches(0.05))
        add_text(s, Inches(0.25), Inches(0.12), Inches(1.11), Inches(0.74),
                 'LOGO', font_size=Pt(12), font_color=MED_GRAY,
                 alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 主标题
        n_lines = title.count('\n') + 1
        title_h = Inches(0.8 + 0.65 * max(n_lines - 1, 0))
        add_text(s, Inches(0.25), Inches(2.57), Inches(4.85), title_h,
                 title, font_size=Pt(40), font_color=GAOYAN_GRAY,
                 bold=True, font_name=FONT_CN)

        # 英文副标题
        sub_y = Inches(2.57) + title_h + Inches(0.1)
        if subtitle:
            add_text(s, Inches(0.25), sub_y, Inches(4.73), Inches(0.50),
                     subtitle, font_size=Pt(24), font_color=GAOYAN_GRAY,
                     font_name=FONT_EN)
            sub_y += Inches(0.60)

        # 日期/版本
        if date or author:
            info = f'{date}  {author}'.strip() if author else date
            add_text(s, Inches(0.38), Inches(5.04), Inches(2.97), Inches(0.34),
                     info, font_size=BODY_SIZE, font_color=MED_GRAY,
                     bold=True, font_name=FONT_EN)

        return s

    def cover_fullscreen(self, title, subtitle='', date=''):
        """全屏背景封面页 — 大图背景 + 右下角圆角信息区。

        基于模板 Slide 1-9: 中国中式餐饮白皮书封面。
        """
        s = self._new_slide()

        # 全屏图片占位
        add_image_placeholder(s, Inches(0), Inches(0), SW, SH, '封面背景图')

        # 右下角圆角信息区
        info_left = Inches(10.29)
        info_top = Inches(5.82)
        add_rect(s, info_left, info_top, Inches(3.05), Inches(1.68), GAOYAN_GREEN)

        # 标题
        add_text(s, Inches(6.24), Inches(3.35), Inches(6.71), Inches(0.84),
                 title, font_size=Pt(44), font_color=WHITE,
                 bold=True, font_name=FONT_CN)

        # 副标题
        if subtitle:
            add_text(s, Inches(6.24), Inches(4.22), Inches(6.36), Inches(1.18),
                     subtitle, font_size=Pt(32), font_color=WHITE,
                     font_name=FONT_EN)

        # 年份
        if date:
            add_text(s, Inches(11.15), Inches(6.34), Inches(2.19), Inches(0.84),
                     date, font_size=Pt(44), font_color=WHITE,
                     bold=True, font_name=FONT_EN)

        return s

    def toc(self, title='目录', items=None, source=''):
        """目录页 — Contents大字 + 编号列表 + 右侧绿色编号圆点。

        基于模板 Slide 11。
        items: list of (num, section_title, description)
        """
        s = self._new_slide()

        # "Contents" 大字
        add_text(s, Inches(0.43), Inches(0.23), Inches(5.95), Inches(1.58),
                 'Contents', font_size=Pt(88), font_color=GAOYAN_GRAY,
                 bold=True, font_name=FONT_EN)

        if not items:
            return s

        # 目录项
        y_start = Inches(0.50)
        row_h = min(Inches(1.13), Inches(6.40) / max(len(items), 1))

        for i, item in enumerate(items):
            num, sec_title = item[0], item[1]
            desc = item[2] if len(item) > 2 else ''
            y = y_start + row_h * i

            # 分隔线
            add_hline(s, Inches(6.81), Inches(1.18) + row_h * i,
                      Inches(6.03), LINE_GRAY, Pt(0.5))

            # 右侧绿色编号圆点
            add_oval(s, Inches(12.17), y + Inches(0.54), Inches(0.55), Inches(0.55),
                     GAOYAN_GREEN, text=str(num), font_size=Pt(16),
                     font_color=WHITE, bold=True)

            # 标题
            add_text(s, Inches(6.67), y + Inches(0.50), Inches(3.57), Inches(0.44),
                     sec_title, font_size=Pt(20), font_color=RGBColor(0x44, 0x54, 0x6A),
                     font_name=FONT_CN)

            # 英文描述
            if desc:
                add_text(s, Inches(6.69), y + Inches(0.85), Inches(3.04), Inches(0.34),
                         desc, font_size=BODY_SIZE, font_color=MED_GRAY,
                         font_name=FONT_EN)

        if source:
            add_source(s, source)
        add_brand_footer(s)

        return s

    def section_divider(self, section_label='', title='', subtitle=''):
        """章节分隔页 — 深蓝背景 + 大号标题。

        基于模板 Layout 9-11 (章节页01/02/03)。
        """
        s = self._new_slide()

        # 深蓝全屏背景
        add_flat_rect(s, Inches(0), Inches(0), SW, SH, GAOYAN_BLUE)

        # 绿色装饰线
        add_flat_rect(s, Inches(0.32), Inches(4.80), Inches(3.0), Inches(0.04), GAOYAN_GREEN)

        # 章节标签
        if section_label:
            add_text(s, Inches(0.32), Inches(4.20), Inches(7.47), Inches(0.50),
                     section_label, font_size=Pt(18), font_color=GAOYAN_GREEN,
                     font_name=FONT_EN)

        # 标题
        add_text(s, Inches(0.32), Inches(5.00), Inches(7.47), Inches(1.45),
                 title, font_size=Pt(40), font_color=WHITE,
                 bold=True, font_name=FONT_CN)

        # 副标题
        if subtitle:
            add_text(s, Inches(0.32), Inches(6.20), Inches(7.47), Inches(0.50),
                     subtitle, font_size=Pt(18), font_color=RGBColor(0xA0, 0xF3, 0xFF),
                     font_name=FONT_CN)

        return s

    def closing(self, title='THANKS!', message='', source_text=''):
        """结尾页 — 大号感谢文字 + 品牌信息。

        基于模板 Layout 15 (1_结尾)。
        """
        s = self._new_slide()

        # 绿色顶部装饰条
        add_flat_rect(s, Inches(0), Inches(0), SW, Inches(0.05), GAOYAN_GREEN)

        # 大号标题
        add_text(s, Inches(1.41), Inches(3.18), Inches(5.01), Inches(0.96),
                 title, font_size=Pt(80), font_color=GAOYAN_GRAY,
                 bold=False, font_name=FONT_EN)

        # 副标题/消息
        if message:
            add_text(s, Inches(1.41), Inches(4.30), Inches(8.0), Inches(0.60),
                     message, font_size=Pt(18), font_color=MED_GRAY,
                     font_name=FONT_CN)

        # 底部绿色装饰线
        add_flat_rect(s, Inches(0), Inches(7.45), SW, Inches(0.05), GAOYAN_GREEN)

        # 品牌信息
        add_text(s, Inches(1.41), Inches(6.50), Inches(4.0), Inches(0.40),
                 '数智餐饮 共创味来', font_size=BODY_SIZE, font_color=MED_GRAY,
                 font_name=FONT_CN)

        return s

    # ══════════════════════════════════════════════
    # B. 数据统计
    # ══════════════════════════════════════════════

    def big_number(self, title, number, unit='', description='', detail_items=None,
                   tag='', sub_conclusion='', source=''):
        """大数字展示页 — 中央大数字 + 单位 + 描述 + 可选明细列表。

        基于模板数据详情页：突出单一核心指标。
        detail_items: list of (label, value) 补充数据对。
        """
        s = self._new_slide()
        self._add_detail_chrome(s, title, tag, sub_conclusion, source)

        # 大数字
        add_text(s, Inches(1.0), Inches(2.40), Inches(10.0), Inches(2.0),
                 str(number), font_size=Pt(96), font_color=DATA_BLUE_4,
                 bold=True, font_name=FONT_NUM, alignment=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.BOTTOM)

        # 单位
        if unit:
            add_text(s, Inches(1.0), Inches(4.40), Inches(10.0), Inches(0.50),
                     unit, font_size=Pt(24), font_color=MED_GRAY,
                     font_name=FONT_CN, alignment=PP_ALIGN.CENTER)

        # 描述
        if description:
            add_text(s, Inches(2.0), Inches(5.10), Inches(8.0), Inches(0.60),
                     description, font_size=BODY_SIZE, font_color=GAOYAN_GRAY,
                     font_name=FONT_CN, alignment=PP_ALIGN.CENTER)

        # 明细列表
        if detail_items:
            n = len(detail_items)
            col_w = min(Inches(3.0), Inches(11.0) / max(n, 1))
            total_w = col_w * n
            start_x = (SW - total_w) // 2
            for i, (label, val) in enumerate(detail_items):
                x = start_x + col_w * i
                add_text(s, x, Inches(5.90), col_w, Inches(0.40),
                         str(val), font_size=Pt(20), font_color=DATA_BLUE_4,
                         bold=True, font_name=FONT_NUM, alignment=PP_ALIGN.CENTER)
                add_text(s, x, Inches(6.30), col_w, Inches(0.35),
                         label, font_size=BODY_SMALL_SIZE, font_color=MED_GRAY,
                         font_name=FONT_CN, alignment=PP_ALIGN.CENTER)

        return s

    def three_stat(self, title, stats, tag='', sub_conclusion='', source=''):
        """三指标仪表盘 — 三列大数字+标签+可选趋势箭头。

        stats: list of dict {number, unit, label, trend(optional: 'up'/'down'/'')}
        """
        s = self._new_slide()
        self._add_detail_chrome(s, title, tag, sub_conclusion, source)

        n = len(stats)
        col_w = Inches(11.0) / max(n, 1)
        start_x = Inches(1.17)

        for i, st in enumerate(stats):
            x = start_x + col_w * i
            cx = x + col_w // 2 - Inches(1.5)

            # 卡片背景
            card_w = col_w - Inches(0.30)
            add_rect(s, x + Inches(0.15), Inches(2.30), card_w, Inches(3.80),
                     BG_GRAY, corner_radius=CARD_CORNER_RADIUS)

            # 数字
            num_text = str(st.get('number', ''))
            add_text(s, cx, Inches(2.80), Inches(3.0), Inches(1.20),
                     num_text, font_size=Pt(54), font_color=DATA_BLUE_4,
                     bold=True, font_name=FONT_NUM, alignment=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.BOTTOM)

            # 单位
            unit = st.get('unit', '')
            if unit:
                add_text(s, cx, Inches(4.00), Inches(3.0), Inches(0.40),
                         unit, font_size=BODY_SIZE, font_color=MED_GRAY,
                         font_name=FONT_CN, alignment=PP_ALIGN.CENTER)

            # 趋势
            trend = st.get('trend', '')
            if trend:
                arrow = '▲' if trend == 'up' else '▼' if trend == 'down' else ''
                color = POSITIVE_GREEN if trend == 'up' else ACCENT_RED
                add_text(s, cx, Inches(4.40), Inches(3.0), Inches(0.35),
                         arrow, font_size=Pt(16), font_color=color,
                         alignment=PP_ALIGN.CENTER)

            # 标签
            label = st.get('label', '')
            add_text(s, cx, Inches(4.90), Inches(3.0), Inches(0.80),
                     label, font_size=BODY_SIZE, font_color=GAOYAN_GRAY,
                     font_name=FONT_CN, alignment=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.TOP)

        return s

    def data_table(self, title, headers, rows, col_widths=None, highlight_cols=None,
                   tag='', sub_conclusion='', source=''):
        """数据表格页 — 深蓝表头 + 斑马纹行 + 可选列高亮。

        headers: list of str
        rows: list of list[str]
        col_widths: list of Inches (auto if None)
        highlight_cols: list of int (0-indexed) — 用DATA_BLUE_4高亮
        """
        s = self._new_slide()
        self._add_detail_chrome(s, title, tag, sub_conclusion, source)

        n_cols = len(headers)
        n_rows = len(rows)
        table_w = Inches(11.80)
        table_left = Inches(0.77)
        table_top = CONTENT_TOP + Inches(0.20)
        row_h = min(Inches(0.55), Inches(4.50) / max(n_rows + 1, 2))

        if col_widths is None:
            cw = table_w / n_cols
            col_widths = [cw] * n_cols

        highlight_cols = highlight_cols or []

        # Header row background
        x = table_left
        for j, hdr in enumerate(headers):
            w = col_widths[j] if j < len(col_widths) else col_widths[-1]
            add_flat_rect(s, x, table_top, w, row_h, DATA_BLUE_1)
            add_text(s, x, table_top, w, row_h, hdr,
                     font_size=BODY_SMALL_SIZE, font_color=WHITE, bold=True,
                     font_name=FONT_CN, alignment=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)
            x += w

        # Data rows
        for i, row in enumerate(rows):
            y = table_top + row_h * (i + 1)
            bg = WHITE if i % 2 == 0 else BG_GRAY
            x = table_left
            for j, cell in enumerate(row):
                w = col_widths[j] if j < len(col_widths) else col_widths[-1]
                add_flat_rect(s, x, y, w, row_h, bg)
                fc = DATA_BLUE_4 if j in highlight_cols else GAOYAN_GRAY
                add_text(s, x, y, w, row_h, str(cell),
                         font_size=BODY_SMALL_SIZE, font_color=fc,
                         bold=(j in highlight_cols), font_name=FONT_CN,
                         alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                x += w

        return s

    # ══════════════════════════════════════════════
    # C. 洞察 / Key Findings
    # ══════════════════════════════════════════════

    def key_insight(self, title, left_content, right_takeaways, sub_conclusion='',
                    tag='', source=''):
        """左内容 + 右侧竖线 + 洞察面板。

        高岩核心版式：左侧放文字/数据，右侧放 Key Takeaways。
        left_content: str or list[str] — 左侧正文
        right_takeaways: list of str — 右侧要点列表
        """
        s = self._new_slide()
        self._add_detail_chrome(s, title, tag, sub_conclusion, source)

        # 左侧内容区
        left_w = Inches(7.80)
        if isinstance(left_content, list):
            left_content = '\n'.join(left_content)
        add_text(s, LM, CONTENT_TOP, left_w, Inches(4.80),
                 left_content, font_size=BODY_SIZE, font_color=GAOYAN_GRAY,
                 font_name=FONT_CN)

        # 右侧竖线
        add_vline(s, INSIGHT_PANEL_VLINE_X, CONTENT_TOP, Inches(4.80),
                  GAOYAN_GREEN, Pt(2))

        # 右侧洞察面板标题
        add_text(s, INSIGHT_PANEL_LEFT, CONTENT_TOP, INSIGHT_PANEL_W, Inches(0.45),
                 'Key Takeaways', font_size=BODY_SIZE, font_color=GAOYAN_BLUE,
                 bold=True, font_name=FONT_EN)

        # 要点列表
        y = CONTENT_TOP + Inches(0.55)
        for tk in right_takeaways:
            add_oval(s, INSIGHT_PANEL_LEFT, y + Inches(0.05),
                     Inches(0.12), Inches(0.12), GAOYAN_GREEN)
            add_text(s, INSIGHT_PANEL_LEFT + Inches(0.22), y,
                     INSIGHT_PANEL_W - Inches(0.22), Inches(0.60),
                     tk, font_size=BODY_SIZE, font_color=GAOYAN_GRAY,
                     font_name=FONT_CN)
            y += Inches(0.70)

        return s

    def chart_insight(self, title, chart_placeholder_label, takeaways,
                      sub_conclusion='', tag='', source=''):
        """图表 + 右侧洞察面板。

        左侧为图表占位区，右侧为 Key Takeaways。
        """
        s = self._new_slide()
        self._add_detail_chrome(s, title, tag, sub_conclusion, source)

        add_image_placeholder(s, LM, CONTENT_TOP, Inches(7.80), Inches(4.80),
                              chart_placeholder_label)

        add_vline(s, INSIGHT_PANEL_VLINE_X, CONTENT_TOP, Inches(4.80),
                  GAOYAN_GREEN, Pt(2))

        add_text(s, INSIGHT_PANEL_LEFT, CONTENT_TOP, INSIGHT_PANEL_W, Inches(0.45),
                 'Key Takeaways', font_size=BODY_SIZE, font_color=GAOYAN_BLUE,
                 bold=True, font_name=FONT_EN)

        y = CONTENT_TOP + Inches(0.55)
        for tk in takeaways:
            add_oval(s, INSIGHT_PANEL_LEFT, y + Inches(0.05),
                     Inches(0.12), Inches(0.12), GAOYAN_GREEN)
            add_text(s, INSIGHT_PANEL_LEFT + Inches(0.22), y,
                     INSIGHT_PANEL_W - Inches(0.22), Inches(0.60),
                     tk, font_size=BODY_SIZE, font_color=GAOYAN_GRAY,
                     font_name=FONT_CN)
            y += Inches(0.70)

        return s

    # ══════════════════════════════════════════════
    # D. 对比评估
    # ══════════════════════════════════════════════

    def side_by_side(self, title, options, tag='', sub_conclusion='', source=''):
        """左右对比页 — N列并排对比。

        options: list of dict {heading, points: list[str], color(optional)}
        """
        s = self._new_slide()
        self._add_detail_chrome(s, title, tag, sub_conclusion, source)

        n = len(options)
        gap = Inches(0.25)
        total_gap = gap * (n - 1)
        col_w = (Inches(11.80) - total_gap) / n
        start_x = Inches(0.77)

        for i, opt in enumerate(options):
            x = start_x + (col_w + gap) * i
            color = opt.get('color', CHART_PALETTE[i % len(CHART_PALETTE)])

            add_flat_rect(s, x, CONTENT_TOP, col_w, Inches(0.06), color)
            add_rect(s, x, CONTENT_TOP + Inches(0.06), col_w, Inches(4.60),
                     BG_GRAY, corner_radius=CARD_CORNER_RADIUS)
            add_text(s, x + Inches(0.20), CONTENT_TOP + Inches(0.20),
                     col_w - Inches(0.40), Inches(0.45),
                     opt.get('heading', ''), font_size=SUB_HEADER_SIZE,
                     font_color=GAOYAN_GRAY, bold=True, font_name=FONT_CN)

            y = CONTENT_TOP + Inches(0.80)
            for pt_text in opt.get('points', []):
                add_text(s, x + Inches(0.30), y, col_w - Inches(0.50), Inches(0.50),
                         f'• {pt_text}', font_size=BODY_SIZE, font_color=GAOYAN_GRAY,
                         font_name=FONT_CN)
                y += Inches(0.50)

        return s

    # ENGINE_SECTION_D_CONTINUE

    def before_after(self, title, before_title, before_points, after_title, after_points,
                     tag='', sub_conclusion='', source=''):
        """前后对比页 — 左Before + 右After + 中间箭头。"""
        s = self._new_slide()
        self._add_detail_chrome(s, title, tag, sub_conclusion, source)

        half_w = Inches(5.50)
        left_x = Inches(0.77)
        right_x = Inches(7.10)

        # Before
        add_rect(s, left_x, CONTENT_TOP, half_w, Inches(4.60),
                 BG_GRAY, corner_radius=CARD_CORNER_RADIUS)
        add_flat_rect(s, left_x, CONTENT_TOP, half_w, Inches(0.06), ACCENT_RED)
        add_text(s, left_x + Inches(0.25), CONTENT_TOP + Inches(0.20),
                 half_w - Inches(0.50), Inches(0.45),
                 before_title, font_size=SUB_HEADER_SIZE, font_color=ACCENT_RED,
                 bold=True, font_name=FONT_CN)
        y = CONTENT_TOP + Inches(0.80)
        for pt_text in before_points:
            add_text(s, left_x + Inches(0.35), y, half_w - Inches(0.60), Inches(0.50),
                     f'• {pt_text}', font_size=BODY_SIZE, font_color=GAOYAN_GRAY,
                     font_name=FONT_CN)
            y += Inches(0.50)

        # Arrow
        add_text(s, Inches(6.10), Inches(3.80), Inches(1.0), Inches(0.60),
                 '→', font_size=Pt(40), font_color=GAOYAN_GREEN,
                 bold=True, font_name=FONT_EN, alignment=PP_ALIGN.CENTER)

        # After
        add_rect(s, right_x, CONTENT_TOP, half_w, Inches(4.60),
                 BG_GRAY, corner_radius=CARD_CORNER_RADIUS)
        add_flat_rect(s, right_x, CONTENT_TOP, half_w, Inches(0.06), GAOYAN_GREEN)
        add_text(s, right_x + Inches(0.25), CONTENT_TOP + Inches(0.20),
                 half_w - Inches(0.50), Inches(0.45),
                 after_title, font_size=SUB_HEADER_SIZE, font_color=ACCENT_GREEN,
                 bold=True, font_name=FONT_CN)
        y = CONTENT_TOP + Inches(0.80)
        for pt_text in after_points:
            add_text(s, right_x + Inches(0.35), y, half_w - Inches(0.60), Inches(0.50),
                     f'• {pt_text}', font_size=BODY_SIZE, font_color=GAOYAN_GRAY,
                     font_name=FONT_CN)
            y += Inches(0.50)

        return s

    def pros_cons(self, title, pros, cons, conclusion='',
                  tag='', sub_conclusion='', source=''):
        """优劣分析页 — 左绿Pros + 右红Cons + 底部结论条。"""
        s = self._new_slide()
        self._add_detail_chrome(s, title, tag, sub_conclusion, source)

        half_w = Inches(5.70)
        left_x = Inches(0.77)
        right_x = Inches(6.80)
        card_h = Inches(3.80) if conclusion else Inches(4.60)

        # Pros
        add_rect(s, left_x, CONTENT_TOP, half_w, card_h,
                 BG_GRAY, corner_radius=CARD_CORNER_RADIUS)
        add_flat_rect(s, left_x, CONTENT_TOP, half_w, Inches(0.06), ACCENT_GREEN)
        add_text(s, left_x + Inches(0.25), CONTENT_TOP + Inches(0.15),
                 half_w - Inches(0.50), Inches(0.40),
                 '✓ 优势 / Pros', font_size=SUB_HEADER_SIZE, font_color=ACCENT_GREEN,
                 bold=True, font_name=FONT_CN)
        y = CONTENT_TOP + Inches(0.65)
        for p in pros:
            add_text(s, left_x + Inches(0.35), y, half_w - Inches(0.60), Inches(0.45),
                     f'• {p}', font_size=BODY_SIZE, font_color=GAOYAN_GRAY,
                     font_name=FONT_CN)
            y += Inches(0.45)

        # Cons
        add_rect(s, right_x, CONTENT_TOP, half_w, card_h,
                 BG_GRAY, corner_radius=CARD_CORNER_RADIUS)
        add_flat_rect(s, right_x, CONTENT_TOP, half_w, Inches(0.06), ACCENT_RED)
        add_text(s, right_x + Inches(0.25), CONTENT_TOP + Inches(0.15),
                 half_w - Inches(0.50), Inches(0.40),
                 '✗ 劣势 / Cons', font_size=SUB_HEADER_SIZE, font_color=ACCENT_RED,
                 bold=True, font_name=FONT_CN)
        y = CONTENT_TOP + Inches(0.65)
        for c in cons:
            add_text(s, right_x + Inches(0.35), y, half_w - Inches(0.60), Inches(0.45),
                     f'• {c}', font_size=BODY_SIZE, font_color=GAOYAN_GRAY,
                     font_name=FONT_CN)
            y += Inches(0.45)

        # 底部结论
        if conclusion:
            conc_y = CONTENT_TOP + card_h + Inches(0.20)
            add_rect(s, left_x, conc_y, Inches(11.80), Inches(0.60),
                     GAOYAN_GREEN, corner_radius=CARD_CORNER_RADIUS)
            add_text(s, left_x + Inches(0.25), conc_y, Inches(11.30), Inches(0.60),
                     conclusion, font_size=BODY_SIZE, font_color=GAOYAN_GRAY,
                     bold=True, font_name=FONT_CN, anchor=MSO_ANCHOR.MIDDLE)

        return s

    # ENGINE_SECTION_E_PLACEHOLDER

    # ══════════════════════════════════════════════
    # E. 框架矩阵
    # ══════════════════════════════════════════════

    def matrix_2x2(self, title, quadrants, axis_labels=None,
                   tag='', sub_conclusion='', source=''):
        """四象限矩阵页 — 2×2 grid + 轴标签。

        quadrants: list of 4 dict {heading, items: list[str]} — 顺序: TL, TR, BL, BR
        axis_labels: dict {top, bottom, left, right} (optional)
        """
        s = self._new_slide()
        self._add_detail_chrome(s, title, tag, sub_conclusion, source)

        grid_left = Inches(1.50)
        grid_top = CONTENT_TOP + Inches(0.30)
        cell_w = Inches(4.80)
        cell_h = Inches(2.10)
        gap = Inches(0.15)

        positions = [
            (grid_left, grid_top),                              # TL
            (grid_left + cell_w + gap, grid_top),               # TR
            (grid_left, grid_top + cell_h + gap),               # BL
            (grid_left + cell_w + gap, grid_top + cell_h + gap) # BR
        ]
        colors = [DATA_BLUE_4, DATA_BLUE_6, ACCENT_YELLOW, ACCENT_GREEN]

        for i, q in enumerate(quadrants[:4]):
            x, y = positions[i]
            add_rect(s, x, y, cell_w, cell_h, BG_GRAY, corner_radius=CARD_CORNER_RADIUS)
            add_flat_rect(s, x, y, cell_w, Inches(0.05), colors[i % 4])
            add_text(s, x + Inches(0.20), y + Inches(0.12), cell_w - Inches(0.40), Inches(0.35),
                     q.get('heading', ''), font_size=EMPHASIS_SIZE, font_color=GAOYAN_GRAY,
                     bold=True, font_name=FONT_CN)
            iy = y + Inches(0.55)
            for item in q.get('items', []):
                add_text(s, x + Inches(0.25), iy, cell_w - Inches(0.50), Inches(0.35),
                         f'• {item}', font_size=BODY_SMALL_SIZE, font_color=GAOYAN_GRAY,
                         font_name=FONT_CN)
                iy += Inches(0.35)

        # 轴标签
        if axis_labels:
            al = axis_labels
            mid_x = grid_left + cell_w + gap // 2
            mid_y = grid_top + cell_h + gap // 2
            if al.get('top'):
                add_text(s, mid_x - Inches(1.5), grid_top - Inches(0.35), Inches(3.0), Inches(0.30),
                         al['top'], font_size=TAG_SIZE, font_color=MED_GRAY,
                         font_name=FONT_CN, alignment=PP_ALIGN.CENTER)
            if al.get('bottom'):
                add_text(s, mid_x - Inches(1.5), grid_top + cell_h * 2 + gap + Inches(0.05),
                         Inches(3.0), Inches(0.30), al['bottom'],
                         font_size=TAG_SIZE, font_color=MED_GRAY,
                         font_name=FONT_CN, alignment=PP_ALIGN.CENTER)
            if al.get('left'):
                add_text(s, grid_left - Inches(1.20), mid_y - Inches(0.15),
                         Inches(1.10), Inches(0.30), al['left'],
                         font_size=TAG_SIZE, font_color=MED_GRAY,
                         font_name=FONT_CN, alignment=PP_ALIGN.RIGHT)
            if al.get('right'):
                add_text(s, grid_left + cell_w * 2 + gap + Inches(0.10), mid_y - Inches(0.15),
                         Inches(1.10), Inches(0.30), al['right'],
                         font_size=TAG_SIZE, font_color=MED_GRAY,
                         font_name=FONT_CN, alignment=PP_ALIGN.LEFT)

        return s

    def process_chevron(self, title, steps, tag='', sub_conclusion='', source=''):
        """流程箭头页 — Chevron 横向流程。

        steps: list of dict {number, heading, description}
        """
        s = self._new_slide()
        self._add_detail_chrome(s, title, tag, sub_conclusion, source)

        n = len(steps)
        total_w = Inches(11.80)
        step_w = total_w / max(n, 1)
        start_x = Inches(0.77)
        chev_top = CONTENT_TOP + Inches(0.50)
        chev_h = Inches(1.20)

        for i, step in enumerate(steps):
            x = start_x + step_w * i
            color = CHART_PALETTE[i % len(CHART_PALETTE)]

            # Chevron body (simplified as rounded rect with arrow overlay)
            add_rect(s, x + Inches(0.05), chev_top, step_w - Inches(0.10), chev_h,
                     color, corner_radius=Inches(0.08))

            # 编号
            num = step.get('number', i + 1)
            add_text(s, x + Inches(0.15), chev_top + Inches(0.10),
                     step_w - Inches(0.30), Inches(0.40),
                     str(num), font_size=Pt(28), font_color=WHITE,
                     bold=True, font_name=FONT_NUM, alignment=PP_ALIGN.CENTER)

            # 标题
            add_text(s, x + Inches(0.15), chev_top + Inches(0.55),
                     step_w - Inches(0.30), Inches(0.50),
                     step.get('heading', ''), font_size=BODY_SIZE, font_color=WHITE,
                     bold=True, font_name=FONT_CN, alignment=PP_ALIGN.CENTER)

            # 描述（chevron下方）
            desc = step.get('description', '')
            if desc:
                add_text(s, x + Inches(0.10), chev_top + chev_h + Inches(0.20),
                         step_w - Inches(0.20), Inches(1.80),
                         desc, font_size=BODY_SMALL_SIZE, font_color=GAOYAN_GRAY,
                         font_name=FONT_CN, alignment=PP_ALIGN.CENTER)

            # 箭头连接符
            if i < n - 1:
                arrow_x = x + step_w - Inches(0.15)
                add_text(s, arrow_x, chev_top + Inches(0.30), Inches(0.30), Inches(0.50),
                         '▶', font_size=Pt(18), font_color=MED_GRAY,
                         alignment=PP_ALIGN.CENTER)

        return s

    # ENGINE_SECTION_E2_PLACEHOLDER

    def vertical_steps(self, title, steps, tag='', sub_conclusion='', source=''):
        """垂直步骤页 — 左侧编号圆点 + 右侧内容卡片。

        steps: list of dict {number, heading, description}
        """
        s = self._new_slide()
        self._add_detail_chrome(s, title, tag, sub_conclusion, source)

        n = len(steps)
        step_h = min(Inches(1.10), Inches(4.60) / max(n, 1))
        start_y = CONTENT_TOP + Inches(0.20)
        circle_x = Inches(1.20)
        card_x = Inches(2.00)
        card_w = Inches(10.20)

        for i, step in enumerate(steps):
            y = start_y + step_h * i

            # 编号圆点
            num = step.get('number', i + 1)
            add_oval(s, circle_x, y + Inches(0.10), Inches(0.50), Inches(0.50),
                     CHART_PALETTE[i % len(CHART_PALETTE)],
                     text=str(num), font_size=Pt(18), font_color=WHITE, bold=True)

            # 连接线（非最后一步）
            if i < n - 1:
                add_vline(s, circle_x + Inches(0.24), y + Inches(0.60),
                          step_h - Inches(0.50), LINE_GRAY, Pt(1.5))

            # 标题
            add_text(s, card_x, y + Inches(0.05), card_w, Inches(0.40),
                     step.get('heading', ''), font_size=EMPHASIS_SIZE,
                     font_color=GAOYAN_GRAY, bold=True, font_name=FONT_CN)

            # 描述
            desc = step.get('description', '')
            if desc:
                add_text(s, card_x, y + Inches(0.45), card_w, Inches(0.55),
                         desc, font_size=BODY_SMALL_SIZE, font_color=MED_GRAY,
                         font_name=FONT_CN)

        return s

    def timeline(self, title, milestones, tag='', sub_conclusion='', source=''):
        """时间线页 — 水平时间轴 + 上下交替里程碑。

        milestones: list of dict {date, heading, description}
        """
        s = self._new_slide()
        self._add_detail_chrome(s, title, tag, sub_conclusion, source)

        n = len(milestones)
        line_y = Inches(4.20)
        start_x = Inches(1.20)
        end_x = Inches(12.10)
        total_w = end_x - start_x

        # 水平主线
        add_hline(s, start_x, line_y, total_w, GAOYAN_GREEN, Pt(3))

        step_w = total_w / max(n, 1)

        for i, ms in enumerate(milestones):
            cx = start_x + step_w * (i + 0.5)

            # 节点圆点
            add_oval(s, cx - Inches(0.15), line_y - Inches(0.15),
                     Inches(0.30), Inches(0.30), GAOYAN_GREEN)

            # 上下交替
            if i % 2 == 0:
                # 上方
                add_text(s, cx - step_w / 2, line_y - Inches(1.80),
                         step_w, Inches(0.35),
                         ms.get('date', ''), font_size=BODY_SMALL_SIZE,
                         font_color=DATA_BLUE_4, bold=True, font_name=FONT_EN,
                         alignment=PP_ALIGN.CENTER)
                add_text(s, cx - step_w / 2, line_y - Inches(1.45),
                         step_w, Inches(0.35),
                         ms.get('heading', ''), font_size=BODY_SIZE,
                         font_color=GAOYAN_GRAY, bold=True, font_name=FONT_CN,
                         alignment=PP_ALIGN.CENTER)
                desc = ms.get('description', '')
                if desc:
                    add_text(s, cx - step_w / 2, line_y - Inches(1.10),
                             step_w, Inches(0.70),
                             desc, font_size=BODY_SMALL_SIZE, font_color=MED_GRAY,
                             font_name=FONT_CN, alignment=PP_ALIGN.CENTER)
                # 竖线连接
                add_vline(s, cx, line_y - Inches(0.40), Inches(0.25), LINE_GRAY, Pt(1))
            else:
                # 下方
                add_vline(s, cx, line_y + Inches(0.15), Inches(0.25), LINE_GRAY, Pt(1))
                add_text(s, cx - step_w / 2, line_y + Inches(0.45),
                         step_w, Inches(0.35),
                         ms.get('date', ''), font_size=BODY_SMALL_SIZE,
                         font_color=DATA_BLUE_4, bold=True, font_name=FONT_EN,
                         alignment=PP_ALIGN.CENTER)
                add_text(s, cx - step_w / 2, line_y + Inches(0.80),
                         step_w, Inches(0.35),
                         ms.get('heading', ''), font_size=BODY_SIZE,
                         font_color=GAOYAN_GRAY, bold=True, font_name=FONT_CN,
                         alignment=PP_ALIGN.CENTER)
                desc = ms.get('description', '')
                if desc:
                    add_text(s, cx - step_w / 2, line_y + Inches(1.15),
                             step_w, Inches(0.70),
                             desc, font_size=BODY_SMALL_SIZE, font_color=MED_GRAY,
                             font_name=FONT_CN, alignment=PP_ALIGN.CENTER)

        return s

    # ══════════════════════════════════════════════
    # F. 团队 / 案例
    # ══════════════════════════════════════════════

    def meet_the_team(self, title, members, tag='', sub_conclusion='', source=''):
        """团队介绍页 — 头像占位 + 姓名 + 职位。

        members: list of dict {name, role, description(optional)}
        """
        s = self._new_slide()
        self._add_detail_chrome(s, title, tag, sub_conclusion, source)

        n = len(members)
        col_w = Inches(11.80) / max(n, 1)
        start_x = Inches(0.77)
        avatar_size = min(Inches(1.60), col_w - Inches(0.40))

        for i, m in enumerate(members):
            cx = start_x + col_w * i + col_w / 2

            # 头像占位圆
            add_oval(s, cx - avatar_size / 2, CONTENT_TOP + Inches(0.30),
                     avatar_size, avatar_size, BG_GRAY,
                     text='Photo', font_size=BODY_SMALL_SIZE, font_color=MED_GRAY)

            # 姓名
            name_y = CONTENT_TOP + Inches(0.50) + avatar_size
            add_text(s, cx - col_w / 2, name_y, col_w, Inches(0.40),
                     m.get('name', ''), font_size=EMPHASIS_SIZE, font_color=GAOYAN_GRAY,
                     bold=True, font_name=FONT_CN, alignment=PP_ALIGN.CENTER)

            # 职位
            add_text(s, cx - col_w / 2, name_y + Inches(0.40), col_w, Inches(0.35),
                     m.get('role', ''), font_size=BODY_SMALL_SIZE, font_color=MED_GRAY,
                     font_name=FONT_CN, alignment=PP_ALIGN.CENTER)

            # 描述
            desc = m.get('description', '')
            if desc:
                add_text(s, cx - col_w / 2 + Inches(0.10), name_y + Inches(0.80),
                         col_w - Inches(0.20), Inches(1.20),
                         desc, font_size=BODY_SMALL_SIZE, font_color=GAOYAN_GRAY,
                         font_name=FONT_CN, alignment=PP_ALIGN.CENTER)

        return s

    # ENGINE_SECTION_F2_PLACEHOLDER

    def case_study(self, title, sections, result_box=None,
                   tag='', sub_conclusion='', source=''):
        """案例研究页 — 左侧时间线节点 + 内容区 + 可选结果框。

        sections: list of dict {heading, content}
        result_box: dict {heading, content} (optional) — 底部结果框
        """
        s = self._new_slide()
        self._add_detail_chrome(s, title, tag, sub_conclusion, source)

        n = len(sections)
        avail_h = Inches(3.80) if result_box else Inches(4.60)
        sec_h = avail_h / max(n, 1)
        start_y = CONTENT_TOP + Inches(0.20)
        line_x = Inches(1.30)
        card_x = Inches(2.00)
        card_w = Inches(10.30)

        for i, sec in enumerate(sections):
            y = start_y + sec_h * i

            # 蓝色圆点
            add_oval(s, line_x - Inches(0.15), y + Inches(0.10),
                     Inches(0.30), Inches(0.30),
                     CHART_PALETTE[i % len(CHART_PALETTE)])

            # 连接线
            if i < n - 1:
                add_vline(s, line_x, y + Inches(0.40),
                          sec_h - Inches(0.30), LINE_GRAY, Pt(1.5))

            # 标题
            add_text(s, card_x, y + Inches(0.05), card_w, Inches(0.40),
                     sec.get('heading', ''), font_size=EMPHASIS_SIZE,
                     font_color=GAOYAN_GRAY, bold=True, font_name=FONT_CN)

            # 内容
            add_text(s, card_x, y + Inches(0.45), card_w, sec_h - Inches(0.55),
                     sec.get('content', ''), font_size=BODY_SIZE,
                     font_color=GAOYAN_GRAY, font_name=FONT_CN)

        # 结果框
        if result_box:
            rb_y = start_y + avail_h + Inches(0.15)
            add_rect(s, Inches(0.77), rb_y, Inches(11.80), Inches(0.70),
                     GAOYAN_GREEN, corner_radius=CARD_CORNER_RADIUS)
            heading = result_box.get('heading', '结果')
            content = result_box.get('content', '')
            add_text(s, Inches(1.00), rb_y, Inches(2.0), Inches(0.70),
                     heading, font_size=BODY_SIZE, font_color=GAOYAN_GRAY,
                     bold=True, font_name=FONT_CN, anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, Inches(3.00), rb_y, Inches(9.30), Inches(0.70),
                     content, font_size=BODY_SIZE, font_color=GAOYAN_GRAY,
                     font_name=FONT_CN, anchor=MSO_ANCHOR.MIDDLE)

        return s

    # ══════════════════════════════════════════════
    # G. 图表（高岩配色）
    # ══════════════════════════════════════════════

    def grouped_bar(self, title, categories, series, data,
                    tag='', sub_conclusion='', source=''):
        """分组柱状图 — 多系列柱状图。

        categories: list[str] — X轴分类
        series: list[str] — 系列名称
        data: list[list[float]] — data[series_idx][cat_idx]
        """
        s = self._new_slide()
        self._add_detail_chrome(s, title, tag, sub_conclusion, source)

        n_cat = len(categories)
        n_ser = len(series)
        chart_left = Inches(1.20)
        chart_bottom = Inches(6.50)
        chart_top = CONTENT_TOP + Inches(0.50)
        chart_w = Inches(10.80)
        chart_h = chart_bottom - chart_top

        # Find max value
        max_val = max(max(row) for row in data) if data else 1
        if max_val == 0:
            max_val = 1

        group_w = chart_w / max(n_cat, 1)
        bar_w = group_w * 0.70 / max(n_ser, 1)
        gap = group_w * 0.15

        # 基线
        add_hline(s, chart_left, chart_bottom, chart_w, MED_GRAY, Pt(1))

        for ci, cat in enumerate(categories):
            gx = chart_left + group_w * ci

            for si in range(n_ser):
                val = data[si][ci] if si < len(data) and ci < len(data[si]) else 0
                bh = (val / max_val) * chart_h * 0.85
                bx = gx + gap + bar_w * si
                by = chart_bottom - bh
                color = CHART_PALETTE[si % len(CHART_PALETTE)]

                add_rect(s, bx, by, bar_w, bh, color, corner_radius=Inches(0.04))

                # 数值标签
                add_text(s, bx, by - Inches(0.30), bar_w, Inches(0.25),
                         str(val), font_size=TAG_SIZE, font_color=GAOYAN_GRAY,
                         font_name=FONT_NUM, alignment=PP_ALIGN.CENTER)

            # 分类标签
            add_text(s, gx, chart_bottom + Inches(0.05), group_w, Inches(0.35),
                     cat, font_size=BODY_SMALL_SIZE, font_color=GAOYAN_GRAY,
                     font_name=FONT_CN, alignment=PP_ALIGN.CENTER)

        # 图例
        legend_y = CONTENT_TOP + Inches(0.10)
        lx = Inches(9.00)
        for si, ser_name in enumerate(series):
            color = CHART_PALETTE[si % len(CHART_PALETTE)]
            add_flat_rect(s, lx, legend_y, Inches(0.20), Inches(0.15), color)
            add_text(s, lx + Inches(0.25), legend_y - Inches(0.02), Inches(1.50), Inches(0.25),
                     ser_name, font_size=TAG_SIZE, font_color=GAOYAN_GRAY,
                     font_name=FONT_CN)
            lx += Inches(1.80)

        return s

    def stacked_bar(self, title, periods, series, data,
                    tag='', sub_conclusion='', source=''):
        """堆叠柱状图 — 同一柱子上堆叠多个系列。

        periods: list[str] — X轴（年份/时期）
        series: list[str] — 系列名称
        data: list[list[float]] — data[series_idx][period_idx]
        """
        s = self._new_slide()
        self._add_detail_chrome(s, title, tag, sub_conclusion, source)

        n_per = len(periods)
        n_ser = len(series)
        chart_left = Inches(1.20)
        chart_bottom = Inches(6.50)
        chart_top = CONTENT_TOP + Inches(0.50)
        chart_w = Inches(10.80)
        chart_h = chart_bottom - chart_top

        # Compute stacked totals
        totals = []
        for pi in range(n_per):
            t = sum(data[si][pi] for si in range(n_ser) if si < len(data) and pi < len(data[si]))
            totals.append(t)
        max_total = max(totals) if totals else 1
        if max_total == 0:
            max_total = 1

        bar_w = chart_w / max(n_per, 1) * 0.60
        gap = chart_w / max(n_per, 1) * 0.20

        add_hline(s, chart_left, chart_bottom, chart_w, MED_GRAY, Pt(1))

        for pi in range(n_per):
            bx = chart_left + (chart_w / n_per) * pi + gap
            cum_h = 0

            for si in range(n_ser):
                val = data[si][pi] if si < len(data) and pi < len(data[si]) else 0
                bh = (val / max_total) * chart_h * 0.85
                by = chart_bottom - cum_h - bh
                color = CHART_PALETTE[si % len(CHART_PALETTE)]

                radius = Inches(0.04) if si == n_ser - 1 else 0
                add_rect(s, bx, by, bar_w, bh, color, corner_radius=radius)

                # 值标签（大段才显示）
                if bh > Inches(0.30):
                    add_text(s, bx, by, bar_w, bh, str(val),
                             font_size=TAG_SIZE, font_color=WHITE,
                             font_name=FONT_NUM, alignment=PP_ALIGN.CENTER,
                             anchor=MSO_ANCHOR.MIDDLE)

                cum_h += bh

            add_text(s, bx - gap / 2, chart_bottom + Inches(0.05),
                     bar_w + gap, Inches(0.35), periods[pi],
                     font_size=BODY_SMALL_SIZE, font_color=GAOYAN_GRAY,
                     font_name=FONT_CN, alignment=PP_ALIGN.CENTER)

        # 图例
        legend_y = CONTENT_TOP + Inches(0.10)
        lx = Inches(9.00)
        for si, ser_name in enumerate(series):
            color = CHART_PALETTE[si % len(CHART_PALETTE)]
            add_flat_rect(s, lx, legend_y, Inches(0.20), Inches(0.15), color)
            add_text(s, lx + Inches(0.25), legend_y - Inches(0.02), Inches(1.50), Inches(0.25),
                     ser_name, font_size=TAG_SIZE, font_color=GAOYAN_GRAY,
                     font_name=FONT_CN)
            lx += Inches(1.80)

        return s

    # ENGINE_SECTION_G2_PLACEHOLDER

    def donut(self, title, segments, center_label='', summary='',
              tag='', sub_conclusion='', source=''):
        """环形图页 — BLOCK_ARC 环形 + 中心标签 + 右侧图例。

        segments: list of dict {label, value, color(optional)}
        center_label: str — 环心文字
        summary: str — 底部摘要文字
        """
        s = self._new_slide()
        self._add_detail_chrome(s, title, tag, sub_conclusion, source)

        total = sum(seg.get('value', 0) for seg in segments)
        if total == 0:
            total = 1

        donut_size = Inches(3.60)
        donut_left = Inches(2.00)
        donut_top = CONTENT_TOP + Inches(0.40)

        # Draw arcs
        cumulative = 0
        for i, seg in enumerate(segments):
            val = seg.get('value', 0)
            start_deg = (cumulative / total) * 360
            end_deg = ((cumulative + val) / total) * 360
            color = seg.get('color', CHART_PALETTE[i % len(CHART_PALETTE)])

            add_block_arc(s, donut_left, donut_top, donut_size, donut_size,
                          start_deg, end_deg, 25000, color)
            cumulative += val

        # 中心标签
        if center_label:
            add_text(s, donut_left, donut_top, donut_size, donut_size,
                     center_label, font_size=Pt(20), font_color=GAOYAN_GRAY,
                     bold=True, font_name=FONT_CN, alignment=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)

        # 右侧图例
        legend_x = Inches(7.00)
        legend_y = CONTENT_TOP + Inches(0.80)
        for i, seg in enumerate(segments):
            color = seg.get('color', CHART_PALETTE[i % len(CHART_PALETTE)])
            val = seg.get('value', 0)
            pct = f'{val / total * 100:.1f}%' if total > 0 else '0%'

            add_flat_rect(s, legend_x, legend_y, Inches(0.25), Inches(0.20), color)
            add_text(s, legend_x + Inches(0.35), legend_y - Inches(0.02),
                     Inches(3.50), Inches(0.30),
                     f'{seg.get("label", "")}  {pct}', font_size=BODY_SIZE,
                     font_color=GAOYAN_GRAY, font_name=FONT_CN)
            legend_y += Inches(0.45)

        # 摘要
        if summary:
            add_text(s, Inches(0.77), Inches(6.20), Inches(11.80), Inches(0.50),
                     summary, font_size=BODY_SIZE, font_color=GAOYAN_GRAY,
                     font_name=FONT_CN, alignment=PP_ALIGN.CENTER)

        return s

    def horizontal_bar(self, title, items, tag='', sub_conclusion='', source=''):
        """水平柱状图 — 横向条形图。

        items: list of dict {label, value, color(optional)}
        """
        s = self._new_slide()
        self._add_detail_chrome(s, title, tag, sub_conclusion, source)

        n = len(items)
        max_val = max(it.get('value', 0) for it in items) if items else 1
        if max_val == 0:
            max_val = 1

        bar_area_left = Inches(3.50)
        bar_area_w = Inches(8.50)
        label_left = Inches(0.77)
        label_w = Inches(2.60)
        start_y = CONTENT_TOP + Inches(0.30)
        bar_h_each = min(Inches(0.60), Inches(4.50) / max(n, 1))
        gap = Inches(0.10)

        for i, it in enumerate(items):
            y = start_y + (bar_h_each + gap) * i
            val = it.get('value', 0)
            bw = (val / max_val) * bar_area_w * 0.90
            color = it.get('color', CHART_PALETTE[i % len(CHART_PALETTE)])

            # 标签
            add_text(s, label_left, y, label_w, bar_h_each,
                     it.get('label', ''), font_size=BODY_SIZE, font_color=GAOYAN_GRAY,
                     font_name=FONT_CN, alignment=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

            # 条形
            add_rect(s, bar_area_left, y + Inches(0.05), bw, bar_h_each - Inches(0.10),
                     color, corner_radius=Inches(0.04))

            # 值标签
            add_text(s, bar_area_left + bw + Inches(0.10), y,
                     Inches(1.00), bar_h_each, str(val),
                     font_size=BODY_SIZE, font_color=GAOYAN_GRAY,
                     bold=True, font_name=FONT_NUM, anchor=MSO_ANCHOR.MIDDLE)

        return s

    # ══════════════════════════════════════════════
    # H. 高岩特色版式
    # ══════════════════════════════════════════════

    def egpm_flavor_card(self, title, chart_placeholder_label, flavor_cards,
                         tag='', sub_conclusion='', source=''):
        """EGPM风味卡片页 — 左侧散点图占位 + 右侧风味卡片。

        基于模板 Slide 155-159: EGPM趋势分析页。
        chart_placeholder_label: str
        flavor_cards: list of dict {title, description, image_label(optional)}
        """
        s = self._new_slide()
        self._add_detail_chrome(s, title, tag, sub_conclusion, source)

        # 左侧图表占位
        add_image_placeholder(s, LM, CONTENT_TOP, Inches(6.50), Inches(4.60),
                              chart_placeholder_label)

        # 右侧风味卡片
        card_left = Inches(7.70)
        card_w = Inches(4.80)
        n = len(flavor_cards)
        card_h = min(Inches(1.40), Inches(4.60) / max(n, 1))
        card_gap = Inches(0.15)

        for i, fc in enumerate(flavor_cards):
            y = CONTENT_TOP + (card_h + card_gap) * i

            # 浅蓝卡片背景 (#DAF1FC)
            add_rect(s, card_left, y, card_w, card_h,
                     LIGHT_BLUE_BG, corner_radius=CARD_CORNER_RADIUS)

            # 图片占位（左侧小方块）
            img_label = fc.get('image_label', '图')
            add_rect(s, card_left + Inches(0.10), y + Inches(0.10),
                     Inches(1.00), card_h - Inches(0.20),
                     BG_GRAY, corner_radius=Inches(0.06))
            add_text(s, card_left + Inches(0.10), y + Inches(0.10),
                     Inches(1.00), card_h - Inches(0.20),
                     img_label, font_size=TAG_SIZE, font_color=MED_GRAY,
                     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

            # 卡片标题
            add_text(s, card_left + Inches(1.20), y + Inches(0.10),
                     card_w - Inches(1.40), Inches(0.35),
                     fc.get('title', ''), font_size=BODY_SIZE, font_color=GAOYAN_GRAY,
                     bold=True, font_name=FONT_CN)

            # 描述
            desc = fc.get('description', '')
            if desc:
                add_text(s, card_left + Inches(1.20), y + Inches(0.45),
                         card_w - Inches(1.40), card_h - Inches(0.55),
                         desc, font_size=BODY_SMALL_SIZE, font_color=MED_GRAY,
                         font_name=FONT_CN)

        return s

    def three_column_insight(self, title, columns, takeaways=None,
                             tag='', sub_conclusion='', source=''):
        """三列 + 右侧洞察面板。

        columns: list of 3 dict {heading, content}
        takeaways: list of str (optional) — 右侧洞察
        """
        s = self._new_slide()
        self._add_detail_chrome(s, title, tag, sub_conclusion, source)

        has_takeaways = takeaways and len(takeaways) > 0
        cols_w = Inches(7.80) if has_takeaways else Inches(11.80)
        n = len(columns)
        col_w = cols_w / max(n, 1)
        start_x = LM

        for i, col in enumerate(columns):
            x = start_x + col_w * i
            color = CHART_PALETTE[i % len(CHART_PALETTE)]

            add_flat_rect(s, x + Inches(0.05), CONTENT_TOP,
                          col_w - Inches(0.10), Inches(0.05), color)

            add_text(s, x + Inches(0.15), CONTENT_TOP + Inches(0.15),
                     col_w - Inches(0.30), Inches(0.40),
                     col.get('heading', ''), font_size=EMPHASIS_SIZE,
                     font_color=GAOYAN_GRAY, bold=True, font_name=FONT_CN)

            add_text(s, x + Inches(0.15), CONTENT_TOP + Inches(0.60),
                     col_w - Inches(0.30), Inches(4.00),
                     col.get('content', ''), font_size=BODY_SIZE,
                     font_color=GAOYAN_GRAY, font_name=FONT_CN)

        # 右侧洞察面板
        if has_takeaways:
            add_vline(s, INSIGHT_PANEL_VLINE_X, CONTENT_TOP, Inches(4.80),
                      GAOYAN_GREEN, Pt(2))
            add_text(s, INSIGHT_PANEL_LEFT, CONTENT_TOP, INSIGHT_PANEL_W, Inches(0.45),
                     'Key Takeaways', font_size=BODY_SIZE, font_color=GAOYAN_BLUE,
                     bold=True, font_name=FONT_EN)
            y = CONTENT_TOP + Inches(0.55)
            for tk in takeaways:
                add_oval(s, INSIGHT_PANEL_LEFT, y + Inches(0.05),
                         Inches(0.12), Inches(0.12), GAOYAN_GREEN)
                add_text(s, INSIGHT_PANEL_LEFT + Inches(0.22), y,
                         INSIGHT_PANEL_W - Inches(0.22), Inches(0.60),
                         tk, font_size=BODY_SIZE, font_color=GAOYAN_GRAY,
                         font_name=FONT_CN)
                y += Inches(0.70)

        return s

    # ══════════════════════════════════════════════
    # I. 保存
    # ══════════════════════════════════════════════

    def save(self, outpath):
        """保存PPTX + 执行full_cleanup清理XML。"""
        self.prs.save(outpath)
        full_cleanup(outpath)
        return outpath
