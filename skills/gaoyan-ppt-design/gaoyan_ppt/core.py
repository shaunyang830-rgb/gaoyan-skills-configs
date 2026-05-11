# GaoyanEngine Core — 高岩科技PPT基础绘图原语
# Handles: shape creation, text rendering, CJK fonts, XML cleanup, rounded rects

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor

from .constants import *


def set_ea_font(run, typeface=FONT_CN):
    """Set East Asian font for Chinese/CJK text runs."""
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', typeface)


def _clean_shape(shape):
    """Remove p:style from shape to prevent file corruption from theme effects."""
    sp = shape._element
    style = sp.find(qn('p:style'))
    if style is not None:
        sp.remove(style)


def add_rect(slide, left, top, width, height, fill_color, corner_radius=None):
    """Add a rectangle (optionally rounded) with solid fill, no border, no p:style.

    高岩规范：默认使用圆角矩形（roundRect）。
    传 corner_radius=0 或 corner_radius=False 可强制直角。
    """
    if corner_radius is False or corner_radius == 0:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        # Set corner radius via XML adjustment
        if corner_radius is None:
            corner_radius = CORNER_RADIUS
        sp = shape._element.find(qn('p:spPr'))
        prstGeom = sp.find(qn('a:prstGeom'))
        if prstGeom is not None:
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = prstGeom.makeelement(qn('a:avLst'), {})
                prstGeom.append(avLst)
            for gd in avLst.findall(qn('a:gd')):
                avLst.remove(gd)
            # Compute ratio: corner_radius relative to min(width, height)
            min_dim = min(int(width), int(height))
            if min_dim > 0:
                ratio = int(int(corner_radius) / min_dim * 50000)
                ratio = max(0, min(ratio, 50000))
            else:
                ratio = 16667  # default ~1/3
            gd = avLst.makeelement(qn('a:gd'), {'name': 'adj', 'fmla': f'val {ratio}'})
            avLst.append(gd)

    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    _clean_shape(shape)
    return shape


def add_flat_rect(slide, left, top, width, height, fill_color):
    """Add a flat (non-rounded) rectangle. Shortcut for add_rect(..., corner_radius=0)."""
    return add_rect(slide, left, top, width, height, fill_color, corner_radius=0)


def add_oval(slide, left, top, width, height, fill_color, text='', font_size=Pt(14),
             font_color=WHITE, bold=True, font_name=FONT_NUM):
    """Add a circle/oval with centered text label."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    _clean_shape(shape)
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(text)
        run.font.size = font_size
        run.font.color.rgb = font_color
        run.font.bold = bold
        run.font.name = font_name
        set_ea_font(run, FONT_CN)
        bodyPr = tf._txBody.find(qn('a:bodyPr'))
        bodyPr.set('anchor', 'ctr')
        for attr in ['lIns', 'tIns', 'rIns', 'bIns']:
            bodyPr.set(attr, '0')
    return shape


def add_hline(slide, x, y, length, color=BLACK, thickness=Pt(0.5)):
    """Draw a horizontal line as a thin flat rectangle (no connector, no p:style).
    高岩规范：线条平直，无阴影。
    """
    h = max(int(thickness), Emu(6350))
    return add_flat_rect(slide, x, y, length, h, color)


def add_vline(slide, x, y, height, color=BLACK, thickness=Pt(0.5)):
    """Draw a vertical line as a thin flat rectangle."""
    w = max(int(thickness), Emu(6350))
    return add_flat_rect(slide, x, y, w, height, color)


def add_text(slide, left, top, width, height, text, font_size=BODY_SIZE,
             font_color=GAOYAN_GRAY, bold=False, font_name=None, alignment=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=None, space_before=None):
    """Add a text box with proper CJK font handling.

    font_name: If None, auto-selects based on content (FONT_CN for Chinese, FONT_EN for English).
    text: str or list[str] (joined with newlines).
    """
    if isinstance(text, list):
        text = '\n'.join(str(t) for t in text)
    text = str(text)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Body properties
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    anchor_map = {MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b', MSO_ANCHOR.TOP: 't'}
    bodyPr.set('anchor', anchor_map.get(anchor, 't'))
    for attr in ['lIns', 'tIns', 'rIns', 'bIns']:
        bodyPr.set(attr, '45720')  # ~0.05"

    # Determine font
    if font_name is None:
        # Check if text contains CJK characters
        has_cjk = any('\u4e00' <= ch <= '\u9fff' or '\u3000' <= ch <= '\u303f' for ch in text)
        font_name = FONT_CN if has_cjk else FONT_EN

    # Split by newlines and create paragraphs
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment

        if space_before is not None and i > 0:
            p.space_before = space_before

        # Line spacing
        if line_spacing is not None:
            p.line_spacing = line_spacing
        elif font_size and int(font_size) >= int(Pt(18)):
            p.line_spacing = 0.93  # Tighter for titles
        elif font_size:
            p.line_spacing = Pt(int(font_size) / 12700 * 1.35)  # Fixed for body

        run = p.add_run()
        run.text = line
        run.font.size = font_size
        run.font.color.rgb = font_color
        run.font.bold = bold
        run.font.name = font_name
        set_ea_font(run, FONT_CN)

    _clean_shape(txBox)
    return txBox


def add_source(slide, text='Source: 高岩餐观大数据'):
    """Add source attribution at standard position (bottom-left)."""
    return add_text(slide, SOURCE_LEFT, SOURCE_Y, SOURCE_W, SOURCE_H,
                    text, font_size=SOURCE_SIZE, font_color=RGBColor(0x06, 0x20, 0x32),
                    font_name=FONT_SOURCE)


def add_brand_footer(slide):
    """Add '高岩 GAOYAN' brand footer at bottom-right."""
    return add_text(slide, BRAND_LEFT, SOURCE_Y, BRAND_W, SOURCE_H,
                    '高岩 GAOYAN', font_size=SOURCE_SIZE, font_color=GAOYAN_GRAY,
                    font_name=FONT_CN, alignment=PP_ALIGN.RIGHT)


def add_page_number(slide, num, total):
    """Add page number at bottom-right corner."""
    add_text(slide, PAGE_NUM_LEFT, SOURCE_Y, PAGE_NUM_W, SOURCE_H,
             f'{num}', font_size=SOURCE_SIZE, font_color=MED_GRAY,
             alignment=PP_ALIGN.RIGHT)


def add_tag_label(slide, text, fill_color=GAOYAN_GREEN2, text_color=GAOYAN_GRAY):
    """Add the green tag label at top-left corner (高岩标签条).

    This is the signature element of Gaoyan slides — a small rounded-corner
    green label at the top-left that shows the section/chapter name.
    """
    # Green background bar
    tag = add_rect(slide, TAG_LEFT, TAG_TOP, TAG_W, TAG_H, fill_color,
                   corner_radius=Inches(0.05))
    # Text inside
    add_text(slide, TAG_LEFT, TAG_TOP, TAG_W, TAG_H, text,
             font_size=TAG_SIZE, font_color=text_color, bold=True,
             font_name=FONT_CN, alignment=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    return tag


def add_sub_conclusion(slide, text, has_icon=True):
    """Add the sub-conclusion bar below the title (高岩子结论条).

    White background bar with bold text, left vertical accent line, and green circle icon.
    """
    # White background
    add_flat_rect(slide, SUBCON_LEFT, SUBCON_TOP, SUBCON_W, SUBCON_H, WHITE)

    # Vertical accent line (left side)
    add_vline(slide, VLINE_X, VLINE_TOP, VLINE_H, GAOYAN_GREEN, Pt(2))

    # Green circle icon
    if has_icon:
        add_oval(slide, ICON_LEFT, ICON_TOP, ICON_SIZE, ICON_SIZE,
                 GAOYAN_GREEN, text='', font_size=Pt(1))

    # Text
    return add_text(slide, SUBCON_TEXT_LEFT, SUBCON_TOP, SUBCON_TEXT_W, SUBCON_TEXT_H,
                    text, font_size=BODY_SIZE, font_color=BLACK, bold=True,
                    font_name=FONT_CN, anchor=MSO_ANCHOR.MIDDLE)


def add_action_title(slide, text):
    """Add the action title at the standard title position.

    高岩规范：标题必须是一句完整结论，不是简单标签。
    字体：24pt 微软雅黑 Bold
    位置：(0.54, 0) 12.79×0.90
    """
    return add_text(slide, LM, TITLE_TOP, Inches(12.79), TITLE_H, text,
                    font_size=ACTION_TITLE_SIZE, font_color=GAOYAN_GRAY,
                    bold=True, font_name=FONT_CN, anchor=MSO_ANCHOR.BOTTOM)


def add_image_placeholder(slide, left, top, width, height, label='请插入图片'):
    """Add a rounded-corner image placeholder box (高岩风格：圆角).

    Triple-border style: outer rounded rect → inner white → light gray fill.
    """
    # Outer frame
    add_rect(slide, left, top, width, height, BG_GRAY)
    # Inner white border
    add_rect(slide, left + Inches(0.04), top + Inches(0.04),
             width - Inches(0.08), height - Inches(0.08), WHITE)
    # Inner light fill
    add_rect(slide, left + Inches(0.08), top + Inches(0.08),
             width - Inches(0.16), height - Inches(0.16),
             RGBColor(0xF8, 0xF8, 0xF8))
    # Label text
    add_text(slide, left, top + height // 2 - Inches(0.3), width, Inches(0.5),
             f'[ {label} ]', font_size=Pt(22), font_color=LINE_GRAY,
             bold=True, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, left, top + height // 2 + Inches(0.2), width, Inches(0.3),
             label, font_size=Pt(13), font_color=MED_GRAY,
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_block_arc(slide, left, top, width, height, start_deg, end_deg, inner_ratio, color):
    """Draw a BLOCK_ARC shape for donut/pie/gauge charts.

    PPT angle convention: clockwise from 12 o'clock, in 60000ths of a degree.
    """
    sh = slide.shapes.add_shape(MSO_SHAPE.BLOCK_ARC, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    _clean_shape(sh)

    sp = sh._element.find(qn('p:spPr'))
    prstGeom = sp.find(qn('a:prstGeom'))
    if prstGeom is not None:
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = prstGeom.makeelement(qn('a:avLst'), {})
            prstGeom.append(avLst)
        for gd in avLst.findall(qn('a:gd')):
            avLst.remove(gd)
        gd1 = avLst.makeelement(qn('a:gd'), {'name': 'adj1', 'fmla': f'val {int(start_deg * 60000)}'})
        gd2 = avLst.makeelement(qn('a:gd'), {'name': 'adj2', 'fmla': f'val {int(end_deg * 60000)}'})
        gd3 = avLst.makeelement(qn('a:gd'), {'name': 'adj3', 'fmla': f'val {inner_ratio}'})
        avLst.append(gd1)
        avLst.append(gd2)
        avLst.append(gd3)
    return sh


def full_cleanup(outpath):
    """Post-save XML sanitization: remove ALL p:style + theme shadows/3D effects."""
    import zipfile, os
    from lxml import etree

    tmppath = outpath + '.tmp'
    with zipfile.ZipFile(outpath, 'r') as zin:
        with zipfile.ZipFile(tmppath, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename.endswith('.xml'):
                    root = etree.fromstring(data)
                    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                    for style in root.findall(f'.//{{{ns_p}}}style'):
                        style.getparent().remove(style)
                    if 'theme' in item.filename.lower():
                        for tag in ['outerShdw', 'innerShdw', 'scene3d', 'sp3d']:
                            for el in root.findall(f'.//{{{ns_a}}}{tag}'):
                                el.getparent().remove(el)
                    data = etree.tostring(root, xml_declaration=True,
                                          encoding='UTF-8', standalone=True)
                zout.writestr(item, data)
    os.replace(tmppath, outpath)
